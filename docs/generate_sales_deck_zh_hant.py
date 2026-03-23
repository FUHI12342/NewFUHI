#!/usr/bin/env python3
"""
Sales Deck Generator — Traditional Chinese (繁體中文) Version
Outputs PPT + PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ==============================
# Color Palette
# ==============================
BROWN = RGBColor(0x8C, 0x87, 0x6C)
BEIGE = RGBColor(0xF1, 0xF0, 0xEC)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)

FONT_NAME = 'Microsoft JhengHei'


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="\u2713"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = spacing


def add_feature_card(slide, left, top, width, height, title, items, accent_color=BROWN):
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=True)
    line = add_shape(slide, left + Inches(0.15), top + Inches(0.08),
                     width - Inches(0.3), Inches(0.04), accent_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.2),
             width - Inches(0.4), Inches(0.4),
             title, font_size=15, bold=True, color=accent_color)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6),
                    width - Inches(0.4), height - Inches(0.8),
                    items, font_size=11, icon="\u25CF")


def add_picture_fit(slide, img_path, left, top, max_width, max_height):
    with Image.open(img_path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = max_width / max_height
    if ratio > box_ratio:
        w = max_width
        h = max_width / ratio
    else:
        h = max_height
        w = max_height * ratio
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    return slide.shapes.add_picture(img_path, int(x), int(y), width=int(w), height=int(h))


def add_number_highlight(slide, left, top, number, label, color=BROWN):
    add_text(slide, left, top, Inches(2), Inches(0.6),
             number, font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2), Inches(0.3),
             label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ==============================
# Slide Creation
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # Slide 1: Cover
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "全方位店舖DX平台",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "預約 \u2022 點餐 \u2022 支付 \u2022 員工管理 \u2022 薪資\n一個系統全部搞定。",
             font_size=22, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5), Inches(4.0), Inches(3.33), Inches(0.03), WHITE)

    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "TimeBaiBai  服務介紹",
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "2026  /  公司名稱",
             font_size=14, color=BEIGE, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 2: Pain Points
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "這些問題，是否似曾相識？", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    problems = [
        ("預約管道分散",
         "電話、LINE、電子郵件⋯⋯各種預約管道\n常常擔心重複預約。",
         "\U0001F4DE"),
        ("爽約與臨時取消",
         "客人預約了卻沒來。\n浪費了準備工作和營收。",
         "\U0001F630"),
        ("點餐手忙腳亂",
         "尖峰時段忙不過來。\n點餐錯誤頻頻發生。",
         "\U0001F4DD"),
        ("庫存盲點",
         "商品莫名就缺貨了。\n不知道何時該補貨。",
         "\U0001F4E6"),
    ]

    for i, (title, desc, icon) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.6) + row * Inches(2.6)

        card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), LIGHT_GRAY, corner_radius=True)

        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8),
                 icon, font_size=36, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK)

        add_text(slide, x + Inches(1.2), y + Inches(0.75), Inches(4), Inches(1.2),
                 desc, font_size=14, color=GRAY)

    # ==========================================
    # Slide 3: Solution — What is TimeBaiBai
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "TimeBaiBai 全部幫您解決", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "一站式平台，滿足店舖所有需求。只需一支智慧型手機即可開始使用。",
             font_size=16, color=DARK)

    highlights = [
        ("6+", "核心模組"),
        ("7", "支援語言"),
        ("24h", "自動預約"),
        ("$0", "省下紙張成本"),
    ]
    for i, (num, label) in enumerate(highlights):
        x = Inches(1.2) + i * Inches(3)
        add_shape(slide, x - Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.4), WHITE, corner_radius=True)
        add_number_highlight(slide, x, Inches(2.5), num, label)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
             "預約、桌邊點餐、庫存、排班管理、薪資計算、官網管理⋯⋯\n"
             "從單一儀表板管理所有業務——不再需要使用多種工具。\n\n"
             "設定簡單，開通帳號當天即可使用。\n"
             "不需要複雜的安裝或硬體設備。只需從手機或電腦存取即可。",
             font_size=15, color=DARK)

    # ==========================================
    # Slide 4: Features (1/4) Reservations & Orders
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主要功能  \u2460  預約 / 桌邊點餐",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "線上預約系統", [
            "透過LINE預約 — 直接在聊天中完成預約",
            "亦支援電子郵件預約，方便客戶選擇",
            "預付確認預約 — 付款後自動確認（防止爽約）",
            "管理後台一鍵切換預付/後付模式",
            "員工行事曆即時顯示可預約時段",
            "兩種預約流程：依日期或依員工選擇",
            "QR碼報到（零櫃台作業）",
            "預約確認後自動LINE通知員工",
            "支援7國語言（日文、英文、中文、韓文、西班牙文、葡萄牙文）",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "QR碼桌邊點餐系統", [
            "在每張桌上放置QR碼即可開始使用",
            "客人掃描手機 \u2192 瀏覽菜單 \u2192 下單",
            "圖片菜單幫助客人了解菜色",
            "分類標籤快速找到想點的品項",
            "購物車功能方便追加點餐",
            "即時訂單狀態追蹤",
            "支援現金、信用卡、電子支付",
            "管理後台可批量產生QR碼並列印下載",
        ], ACCENT_GREEN)

    # ==========================================
    # Slide 5: Features (2/4) Inventory & Staff
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主要功能  \u2461  庫存 / 員工排班管理",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "庫存管理", [
            "所有商品即時庫存數量追蹤",
            "訂單連動自動扣除庫存",
            "庫存不足時自動發出警示",
            "QR碼進貨（無需條碼掃描器）",
            "完整進出貨紀錄（稽核備查）",
            "依商品分類整理，管理更輕鬆",
            "EC商店與實體店庫存統一管理",
            "盤點功能支援批量調整",
        ], ACCENT_ORANGE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "員工排班管理", [
            "員工從手機提交班表偏好",
            "三級選項：可上班 / 希望上班 / 無法上班",
            "主管從管理後台確定排班",
            "自動排班功能，高效分配人力",
            "確定班表後自動LINE通知員工",
            "員工個人檔案含照片與介紹",
            "員工類型：服務人員 / 門市員工",
            "角色權限：店主 / 主管 / 員工 / 開發者",
        ], ACCENT_BLUE)

    # ==========================================
    # Slide 6: Features (3/4) Payroll & Website
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主要功能  \u2462  薪資 / 官網管理",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "薪資與出勤管理", [
            "從排班紀錄自動產生出勤資料",
            "自動區分一般時數、加班、夜班、假日時數",
            "時薪制與月薪制員工薪資計算",
            "自動計算社會保險費",
            "自動套用所得稅與住民稅扣除",
            "一鍵產生薪資明細",
            "銀行轉帳CSV匯出（全銀格式）",
            "輕鬆設定各項津貼（交通、住宿、家庭）",
        ], ACCENT_GREEN)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "官網與行銷管理", [
            "從管理後台編輯店舖官網",
            "主視覺輪播圖展示店舖魅力",
            "輕鬆更新最新消息與媒體報導",
            "廣告橫幅刊登與管理",
            "社群媒體整合（X / Instagram嵌入）",
            "自訂HTML區塊打造靈活版面",
            "隱私政策與法律頁面可從後台編輯",
            "內建EC商店支援線上銷售",
        ], ACCENT_ORANGE)

    # ==========================================
    # Slide 7: Features (4/4) Analytics & Security
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主要功能  \u2463  數據分析 / 資安",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "數據分析與AI", [
            "銷售儀表板提供每日/每週/每月/每年趨勢圖",
            "菜單工程分析（Star/Dog分類）提升獲利",
            "ABC分析、RFM分析、世代分析、購物籃分析",
            "NPS（顧客滿意度）自動彙整，提高回客率",
            "AI預測各時段最佳人力配置",
            "需求與營收預測，輔助採購與人力規劃",
            "商業洞察自動偵測異常與機會",
            "KPI計分卡一目了然掌握營運狀況",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "資安與稽核", [
            "完整存取日誌（登入、操作、異常偵測）",
            "IP速率限制自動阻擋未授權存取",
            "客戶個資以AES加密安全儲存",
            "自動化資安稽核（每日12項檢查）",
            "角色權限控管（店主 / 主管 / 員工 / 開發者）",
            "完整稽核軌跡：操作紀錄與排班變更日誌",
            "內建SSL加密與CSRF防護",
            "整合符合PCI DSS標準的外部支付服務",
        ], ACCENT_ORANGE)

    # ==========================================
    # Screenshot Slides
    # ==========================================
    screenshots_dir = os.path.dirname(os.path.abspath(__file__))

    def _add_desktop_slide(img_name, title, subtitle, bullets, header_prefix=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        full_title = f"{header_prefix}{title}" if header_prefix else title
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 full_title, font_size=28, bold=True, color=BROWN)
        add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)
        img_path = os.path.join(screenshots_dir, img_name)
        add_shape(slide, Inches(0.4), Inches(1.2), Inches(8.2), Inches(5.8),
                  LIGHT_GRAY, corner_radius=True)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.6))
            except Exception:
                add_text(slide, Inches(2), Inches(3.5), Inches(4), Inches(1),
                         f'[{title}]', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(0.5),
                 subtitle, font_size=18, bold=True, color=DARK)
        add_bullet_list(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(5),
                        bullets, font_size=12, icon="\u2713")

    # ── Customer-Facing Screens ──
    frontend_screens = [
        ('screenshots/zh-hant/front_top_desktop.png', '首頁',
         '客人看到的第一個畫面',
         [
             '主視覺輪播圖展示店舖魅力',
             '清晰的預約流程，提供多種選項',
             '員工排名推廣人氣服務人員',
             '自動顯示最新消息與媒體報導',
             '支援7國語言，服務國際客戶',
         ]),
        ('screenshots/zh-hant/front_staff_list_desktop.png', '員工簡介與店舖資訊',
         '精美呈現您的員工陣容',
         [
             '員工照片與個人介紹',
             '顯示專長與資格證照',
             '直接連結預約行事曆',
             '店舖資訊（營業時間、地圖、交通方式）',
             '從管理後台輕鬆更新內容',
         ]),
        ('screenshots/zh-hant/front_news_desktop.png', '最新消息',
         '讓客戶隨時掌握最新資訊',
         [
             '發布新菜單、活動與優惠促銷',
             '管理後台一鍵發布',
             '依日期自動排序，確保內容最新',
             '社群媒體整合擴大觸及率',
             '透過即時更新促進回客率',
         ]),
        ('screenshots/zh-hant/front_shop_desktop.png', '線上商店',
         '拓展您的銷售管道',
         [
             '在線上販售您的商品',
             '商品搜尋與分類篩選',
             '購物車支援多品項購買',
             '支援信用卡與電子支付',
             '統一庫存管理避免庫存差異',
         ]),
        ('screenshots/zh-hant/table_order_vp.png', 'QR碼桌邊點餐',
         '用智慧型手機輕鬆點餐',
         [
             '只要掃描桌上的QR碼',
             '圖片菜單直觀呈現菜色',
             '分類標籤快速找到品項',
             '購物車功能方便追加點餐',
             '即時訂單狀態追蹤',
         ]),
        ('screenshots/zh-hant/booking_calendar_desktop.png', '預約行事曆',
         '24小時全天候線上預約',
         [
             '行事曆一目了然查看員工可預約時段',
             '客人自選偏好的日期與時段',
             '預付確認機制防止爽約',
             '支援LINE與電子郵件預約',
             '確認後自動通知員工',
         ]),
    ]

    for img_name, title, subtitle, bullets in frontend_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="客戶頁面 \u2014 ")

    # ── Mobile-Responsive ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "完美支援手機瀏覽", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "客戶網站與管理後台皆可在智慧型手機上完美運作。",
             font_size=16, color=DARK)

    mobile_items = [
        ('screenshots/zh-hant/front_top_mobile_vp.png', '客戶首頁'),
        ('screenshots/zh-hant/admin_dashboard_mobile_vp.png', '管理後台'),
    ]

    for i, (img_name, label) in enumerate(mobile_items):
        x = Inches(1.5) + i * Inches(5.5)
        y = Inches(1.8)

        add_shape(slide, x, y, Inches(4.2), Inches(5.0), WHITE, corner_radius=True)

        img_path = os.path.join(screenshots_dir, img_name)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                x + Inches(0.6), y + Inches(0.15),
                                Inches(3.0), Inches(4.4))
            except Exception:
                add_text(slide, x + Inches(0.5), y + Inches(2), Inches(3.2), Inches(1),
                         f'[{label}]', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(4.5), Inches(4.2), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ── Admin Screens ──
    admin_screens = [
        ('screenshots/zh-hant/dashboard_sales_vp.png', '銷售儀表板',
         '一目了然掌握營運狀況',
         [
             '每日 / 每週 / 每月 / 每年營收趨勢',
             '預約KPI（數量、取消率）即時顯示',
             '三通路拆分：總計、EC商店、店內菜單',
             'KPI計分卡一覽關鍵指標',
             '商業洞察自動建議改善方向',
         ]),
        ('screenshots/zh-hant/dashboard_menu_eng_vp.png', '銷售分析（菜單分析）',
         'AI驅動的獲利優化',
         [
             '菜單工程：Star / Plowhorse / Puzzle / Dog 四象限分類',
             'ABC分析：帕累托法則找出主要營收來源',
             '營收預測：Prophet AI模型預測未來14天',
             '每小時銷售熱力圖輔助最佳人力安排',
             'AOV（客單價）趨勢供定價策略參考',
         ]),
        ('screenshots/zh-hant/dashboard_rfm_vp.png', '銷售分析（客戶分析）',
         '數據驅動的客戶洞察',
         [
             'RFM分析：依最近消費、消費頻率、消費金額自動分群',
             '世代分析：視覺化每月新客留存率',
             '購物籃分析：發掘交叉銷售模式與推薦',
             'AI摘要：以自然語言彙整分析結果',
             '建議行動：AI自動建議具體改善步驟',
         ]),
        ('screenshots/zh-hant/shift_calendar_vp.png', '排班行事曆',
         '簡化排班管理流程',
         [
             '員工從智慧型手機提交班表偏好',
             '行事曆檢視所有員工排班一目了然',
             '一鍵自動排班',
             '自動偵測人力不足的日期並發出警示',
             '確定班表後自動LINE通知',
         ]),
        ('screenshots/zh-hant/pos_vp.png', 'POS收銀機',
         '順暢的結帳體驗',
         [
             '分類標籤快速選擇商品',
             '支援現金、信用卡、行動支付與交通卡',
             '自動產生收據並列印',
             '廚房顯示器整合訂單通知',
             '自動產生每日與每月銷售報表',
         ]),
        ('screenshots/zh-hant/inventory_vp.png', '庫存管理',
         '防止缺貨',
         [
             '所有商品即時庫存數量',
             '收到訂單自動扣除庫存',
             '庫存不足警示與通知',
             'QR碼進貨（無需條碼掃描器）',
             '完整進出貨紀錄供稽核備查',
         ]),
        ('screenshots/zh-hant/customer_feedback_vp.png', '顧客回饋 (NPS)',
         '量化顧客心聲',
         [
             'NPS（淨推薦分數）自動彙整與趨勢追蹤',
             '推薦者 / 被動者 / 批評者 色彩標示',
             '依評論回饋提供可執行的改善方向',
             '串聯訂單資料分析服務品質',
             '每月及員工別滿意度趨勢',
         ]),
        ('screenshots/zh-hant/attendance_board_vp.png', '出勤看板',
         '即時掌握員工出勤狀況',
         [
             '一目了然查看上班中、休息中、未打卡狀態',
             '每30秒自動刷新即時追蹤',
             '三種打卡方式：QR碼、PIN碼、智慧型手機',
             '自動區分加班、夜班、假日時數',
             '出勤資料無縫串接薪資計算',
         ]),
        ('screenshots/zh-hant/iot_sensors_vp.png', 'IoT感測器監控',
         '即時監控店舖環境',
         [
             '即時顯示溫度、濕度、氣壓、氣體濃度',
             'PIR動作感測器自動計算來客人數',
             '異常情況（瓦斯外洩、高溫等）發出警示通知',
             '時序圖表追蹤感測器資料趨勢',
             'ESP32裝置連線狀態監控',
         ]),
    ]

    for img_name, title, subtitle, bullets in admin_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="管理後台 \u2014 ")

    # ==========================================
    # Benefits Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "導入效益", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    merits = [
        ("零爽約損失",
         "預付確認預約機制杜絕臨時取消。\n一鍵切換預付/後付模式。",
         "\U0001F6E1"),
        ("LINE輕鬆預約",
         "客人直接在LINE聊天中完成預約。\n無需安裝新的應用程式。",
         "\U0001F4AC"),
        ("提升營收",
         "QR點餐提高追加點餐率。\n多語言支援吸引國際旅客。",
         "\U0001F4C8"),
        ("節省工時",
         "預約、點餐、排班、薪資全面自動化。\n讓員工專注於顧客服務。",
         "\u23F1"),
        ("降低成本",
         "預約簿、點餐單、打卡卡⋯⋯\n省去紙張成本與人工作業。",
         "\U0001F4B0"),
        ("數據驅動決策",
         "視覺化銷售、預約與員工績效。\n以真實數據支持每一個決策。",
         "\U0001F4CA"),
    ]

    for i, (title, desc, icon) in enumerate(merits):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), WHITE, corner_radius=True)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                 icon, font_size=30, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.35),
                 title, font_size=16, bold=True, color=BROWN)

        add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # Getting Started Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "導入流程", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    steps = [
        ("STEP 1", "聯繫我們", "透過電話\n或表單與我們聯繫"),
        ("STEP 2", "需求諮詢", "我們瞭解您目前的\n工作流程與痛點"),
        ("STEP 3", "演示與報價", "觀看產品實際操作。\n我們為您推薦最佳方案"),
        ("STEP 4", "帳號設定", "註冊店舖資訊與員工。\n我們協助初始設定"),
        ("STEP 5", "正式上線", "經過簡短教學後，\n即可正式開始使用！"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.6)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BROWN
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            add_text(slide, x + Inches(1.9), y + Inches(0.2), Inches(0.6), Inches(0.6),
                     "\u2192", font_size=28, color=BROWN, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.2), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(1.5), LIGHT_GRAY, corner_radius=True)
    add_text(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1.2),
             "最快當天即可開始使用。\n"
             "無需特殊硬體或系統安裝。透過智慧型手機、平板或電腦即可立即存取。\n"
             "上線後我們持續提供技術支援、教育訓練與功能更新。",
             font_size=13, color=DARK)

    # ==========================================
    # Pricing Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "方案與價格", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "* 以下價格僅供參考。我們會依據您的需求量身訂製方案。",
             font_size=12, color=GRAY)

    plans = [
        ("輕量方案", "從 $XX/月 起", ACCENT_GREEN,
         "適合小型店舖", [
             "預約管理（LINE與電子郵件）",
             "員工管理（最多5人）",
             "官網管理",
             "電子郵件與線上客服",
         ]),
        ("標準方案", "從 $XX/月 起", ACCENT_BLUE,
         "推薦！", [
             "包含輕量方案所有功能",
             "QR碼桌邊點餐",
             "庫存管理",
             "排班管理",
             "員工人數無上限",
             "電話客服",
         ]),
        ("進階方案", "從 $XX/月 起", BROWN,
         "多店舖 / 完整功能", [
             "包含標準方案所有功能",
             "薪資與出勤管理",
             "銷售儀表板與數據分析",
             "多語言支援（7國語言）",
             "多店舖管理",
             "優先客服",
             "客製化開發",
         ]),
    ]

    for i, (name, price, color, badge, features) in enumerate(plans):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.2)

        card = add_shape(slide, x, y, w, h, LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, w, Inches(0.08), color)

        if badge == "推薦！":
            badge_shape = add_shape(slide, x + Inches(0.8), y - Inches(0.15), Inches(2.2), Inches(0.35), color, corner_radius=True)
            tf = badge_shape.text_frame
            p = tf.paragraphs[0]
            p.text = "\u2605 推薦！"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER

        add_text(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.5),
                 price, font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

        add_shape(slide, x + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(0.02), GRAY)

        add_text(slide, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(0.3),
                 badge if badge != "推薦！" else "適合中型店舖",
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9),
                        w - Inches(0.6), h - Inches(2.1),
                        features, font_size=12, icon="\u2713")

    add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "* 初始設定費與客製化開發費用另行報價。  * 價格未含稅。",
             font_size=10, color=GRAY)

    # ==========================================
    # FAQ Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "常見問題", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    faqs = [
        ("Q. 需要什麼特殊設備嗎？",
         "A. 不需要特殊硬體。使用您現有的智慧型手機、平板或電腦即可。\n"
         "   桌邊點餐功能只需從管理後台列印QR碼即可。"),
        ("Q. 不擅長操作電腦的員工也能使用嗎？",
         "A. 可以。只要會用智慧型手機就能使用TimeBaiBai。介面直覺且操作簡單。\n"
         "   我們在導入時也會提供實機教育訓練。"),
        ("Q. 之後可以新增或變更功能嗎？",
         "A. 當然可以。您可以隨時升級或變更方案。\n"
         "   建議先從核心功能開始，再依需求逐步擴充。"),
        ("Q. 我的資料安全嗎？",
         "A. 客戶資料在儲存時進行加密。所有通訊皆使用SSL加密。\n"
         "   信用卡資訊由符合PCI DSS標準的支付服務處理——我們的伺服器不儲存任何卡片資料。"),
        ("Q. 可以用於多間店舖嗎？",
         "A. 可以。進階方案支援多店舖管理。\n"
         "   可依各分店管理員工與庫存，同時查看所有店舖的彙整報表。"),
    ]

    for i, (q, a) in enumerate(faqs):
        y = Inches(1.2) + i * Inches(1.2)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(1.05), WHITE, corner_radius=True)

        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(11.73), Inches(0.3),
                 q, font_size=13, bold=True, color=BROWN)
        add_text(slide, Inches(0.8), y + Inches(0.4), Inches(11.73), Inches(0.6),
                 a, font_size=11, color=DARK)

    # ==========================================
    # Contact Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "聯繫我們",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             "歡迎洽詢演示、報價或任何問題。",
             font_size=18, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.33), Inches(0.03), WHITE)

    contact_card = add_shape(slide, Inches(3), Inches(3.8), Inches(7.33), Inches(2.5), WHITE, corner_radius=True)

    contacts = [
        "TEL:     000-0000-0000",
        "Email:   info@example.com",
        "Web:    https://timebaibai.com",
        "",
        "公司名稱",
        "地址",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, Inches(4), Inches(4.0) + i * Inches(0.35),
                 Inches(5), Inches(0.35),
                 line, font_size=15 if i < 3 else 13,
                 bold=(i < 3), color=DARK if i < 3 else GRAY,
                 align=PP_ALIGN.LEFT)

    add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "\u00A9 2026 TimeBaiBai. All Rights Reserved.",
             font_size=10, color=BEIGE, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# Output
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    pptx_path = os.path.join(output_dir, 'TimeBaiBai_Service_Overview_ZH_Hant.pptx')
    prs.save(pptx_path)
    print(f"\u2705 PPT saved: {pptx_path}")

    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            print(f"\u2705 PDF saved: {pdf_path}")
        else:
            raise FileNotFoundError("LibreOffice not found")
    except Exception as e:
        print(f"\u26A0\uFE0F  PDF conversion skipped ({e})")
        print(f"   To create PDF from PPTX:")
        print(f"   - Open in PowerPoint / Keynote and export as PDF")
        print(f"   - Or: brew install libreoffice && soffice --headless --convert-to pdf '{pptx_path}'")
