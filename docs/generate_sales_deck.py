#!/usr/bin/env python3
"""
営業用資料 (Sales Deck) 生成スクリプト
PPT + PDF 形式で出力
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ==============================
# カラーパレット
# ==============================
BROWN = RGBColor(0x8C, 0x87, 0x6C)       # salon-brown メインカラー
BEIGE = RGBColor(0xF1, 0xF0, 0xEC)       # salon-beige 背景
DARK = RGBColor(0x33, 0x33, 0x33)         # テキスト
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Meiryo'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="✓"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Meiryo'
        p.space_after = spacing


def add_feature_card(slide, left, top, width, height, title, items, accent_color=BROWN):
    """機能カード（角丸背景 + タイトル + 箇条書き）"""
    # カード背景
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=True)

    # アクセントライン（上部）
    line = add_shape(slide, left + Inches(0.15), top + Inches(0.08),
                     width - Inches(0.3), Inches(0.04), accent_color)

    # タイトル
    add_text(slide, left + Inches(0.2), top + Inches(0.2),
             width - Inches(0.4), Inches(0.4),
             title, font_size=15, bold=True, color=accent_color)

    # 項目
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6),
                    width - Inches(0.4), height - Inches(0.8),
                    items, font_size=11, icon="●")


def add_picture_fit(slide, img_path, left, top, max_width, max_height):
    """画像をアスペクト比を維持してバウンディングボックスに収める"""
    with Image.open(img_path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = max_width / max_height
    if ratio > box_ratio:
        # 横長 → 幅に合わせる
        w = max_width
        h = max_width / ratio
    else:
        # 縦長 → 高さに合わせる
        h = max_height
        w = max_height * ratio
    # 中央揃え
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    return slide.shapes.add_picture(img_path, int(x), int(y), width=int(w), height=int(h))


def add_number_highlight(slide, left, top, number, label, color=BROWN):
    """大きな数字 + ラベルのハイライト"""
    add_text(slide, left, top, Inches(2), Inches(0.6),
             number, font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2), Inches(0.3),
             label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ==============================
# スライド作成
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # スライド 1: 表紙
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BROWN)

    # メインタイトル
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "店舗DXオールインワンプラットフォーム",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # サブタイトル
    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "予約・注文・決済・スタッフ管理・給与計算を\nひとつのシステムで完結",
             font_size=22, color=BEIGE, align=PP_ALIGN.CENTER)

    # 区切り線
    add_shape(slide, Inches(5), Inches(4.0), Inches(3.33), Inches(0.03), WHITE)

    # サービス名
    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "TimeBaiBai  サービスご紹介資料",
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # 日付（編集用プレースホルダー）
    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "2026年  /  株式会社○○○○",
             font_size=14, color=BEIGE, align=PP_ALIGN.CENTER)

    # ==========================================
    # スライド 2: こんなお悩みありませんか？
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "こんなお悩みありませんか？", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    problems = [
        ("予約管理がバラバラ", "電話・LINE・メール…受付方法がバラバラで\nダブルブッキングが怖い", "📞"),
        ("ドタキャン・無断キャンセル", "予約を入れたのに来ない…\n売上損失と準備のムダが痛い", "😰"),
        ("注文対応が大変", "忙しい時間帯に注文を取りに行けない\nオーダーミスが発生する", "📝"),
        ("在庫がわからない", "いつの間にか商品が切れていた\n発注タイミングがわからない", "📦"),
    ]

    for i, (title, desc, icon) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.6) + row * Inches(2.6)

        card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), LIGHT_GRAY, corner_radius=True)

        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8),
                 icon, font_size=36, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK)

        add_text(slide, x + Inches(1.2), y + Inches(0.75), Inches(4), Inches(1.2),
                 desc, font_size=14, color=GRAY)

    # ==========================================
    # スライド 3: 解決策 — TimeBaiBaiとは
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "TimeBaiBai がすべて解決します", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "店舗運営に必要な機能をワンストップで提供。スマホひとつで始められます。",
             font_size=16, color=DARK)

    # 数字ハイライト
    highlights = [
        ("6+", "主要モジュール"),
        ("7言語", "多言語対応"),
        ("24h", "自動受付"),
        ("0円", "紙コスト削減"),
    ]
    for i, (num, label) in enumerate(highlights):
        x = Inches(1.2) + i * Inches(3)
        add_shape(slide, x - Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.4), WHITE, corner_radius=True)
        add_number_highlight(slide, x, Inches(2.5), num, label)

    # 概要テキスト
    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
             "予約管理・テーブル注文・在庫管理・シフト管理・給与計算・ホームページ管理…\n"
             "これまで別々のツールで行っていた業務を、ひとつの管理画面から操作できます。\n\n"
             "導入はカンタン。アカウント発行後、最短即日でご利用いただけます。\n"
             "面倒な設定や工事は不要。いつものスマホ・PCからアクセスするだけです。",
             font_size=15, color=DARK)

    # ==========================================
    # スライド 4: 機能一覧（1/3）予約・注文
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主な機能  ①  予約管理 / テーブル注文",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    # 予約管理カード
    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "オンライン予約システム", [
            "LINEから簡単予約 — トーク画面からそのまま予約完了",
            "メール予約にも対応、お客様の使いやすい方法で受付",
            "前払い予約確定制 — 決済完了で自動確定（ドタキャン防止）",
            "前払い/後払いは管理画面からワンクリックで切替可能",
            "スタッフ別カレンダーで空き状況をリアルタイム表示",
            "日付から選ぶ → スタッフから選ぶ、2つの予約導線",
            "QRコードでチェックイン（受付の手間ゼロ）",
            "予約確定時にスタッフへLINE自動通知",
            "7言語対応（日本語・英語・中韓・スペイン語・ポルトガル語）",
        ], ACCENT_BLUE)

    # テーブル注文カード
    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "QRテーブル注文システム", [
            "各テーブルにQRコードを設置するだけでスタート",
            "お客様がスマホでスキャン → メニュー閲覧 → 注文",
            "写真付きメニューで料理のイメージが伝わる",
            "カテゴリ別タブで欲しい商品がすぐ見つかる",
            "カート機能で追加注文もラクラク",
            "注文履歴をリアルタイムでステータス表示",
            "現金・クレジットカード・電子マネー決済対応",
            "管理画面からQRコード一括生成＆印刷用DL",
        ], ACCENT_GREEN)

    # ==========================================
    # スライド 5: 機能一覧（2/3）在庫・スタッフ
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主な機能  ②  在庫管理 / スタッフ・シフト管理",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "在庫管理", [
            "商品の在庫数をリアルタイムで把握",
            "注文と連動して自動で在庫を差し引き",
            "在庫が少なくなったら自動アラート通知",
            "QRコードで入庫処理（バーコードリーダー不要）",
            "入出庫の履歴をすべて記録（監査対応）",
            "商品カテゴリ別の管理で整理しやすい",
            "ECショップと店頭在庫を一元管理",
            "棚卸し機能で在庫を一括調整",
        ], ACCENT_ORANGE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "スタッフ・シフト管理", [
            "スタッフがスマホからシフト希望を提出",
            "「出勤可能」「希望」「出勤不可」の3段階で申請",
            "マネージャーが管理画面でシフトを確定",
            "自動スケジューリング機能で割り当てを効率化",
            "確定シフトをスタッフにLINEで自動通知",
            "スタッフ別のプロフィール・写真・紹介文を管理",
            "スタッフ種別（キャスト / 店舗スタッフ）で区別",
            "店長・オーナー・開発者など権限ロールを設定",
        ], ACCENT_BLUE)

    # ==========================================
    # スライド 6: 機能一覧（3/3）給与・ホームページ
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主な機能  ③  給与計算 / ホームページ管理",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "給与計算・勤怠管理", [
            "シフトから勤怠データを自動生成",
            "通常勤務・残業・深夜・休日を自動分類",
            "時給制・月給制に対応した給与計算",
            "社会保険料（厚生年金・健康保険等）を自動計算",
            "所得税・住民税の控除を自動反映",
            "給与明細をワンクリックで作成",
            "全銀フォーマットCSVで銀行振込データ出力",
            "手当（通勤・住宅・家族）の設定も簡単",
        ], ACCENT_GREEN)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "ホームページ・集客管理", [
            "管理画面からお店のホームページを編集",
            "ヒーローバナー（スライダー）で写真をアピール",
            "お知らせ・メディア掲載情報を簡単更新",
            "バナー広告の設置・管理",
            "SNS連携（X / Instagram）埋め込み表示",
            "カスタムHTMLブロックで自由にレイアウト",
            "プライバシーポリシー・特商法表記を管理画面から編集",
            "ECショップ機能でオンライン販売も可能",
        ], ACCENT_ORANGE)

    # ==========================================
    # スライド 7: 機能一覧（4/5）分析・セキュリティ — NEW
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "主な機能  ④  データ分析 / セキュリティ",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "データ分析・AI活用", [
            "売上ダッシュボードで月次/週次/日次の推移を可視化",
            "メニューエンジニアリング（Star/Dog分類）で利益改善",
            "ABC分析・RFM分析・コホート分析・バスケット分析",
            "NPS（顧客満足度）自動集計でリピーター獲得に貢献",
            "AIが時間帯別のスタッフ必要人数を自動予測",
            "来客予測・売上予測で仕入れや人員配置を最適化",
            "ビジネスインサイトが異常値や改善機会を自動検知",
            "KPIスコアカードで経営状態をひと目で把握",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "セキュリティ・監査", [
            "全アクセスログを自動記録（ログイン/操作/異常検知）",
            "IP別レートリミットで不正アクセスを自動ブロック",
            "顧客個人情報はAES暗号化で安全に保管",
            "セキュリティ自動監査（12項目を毎日チェック）",
            "権限別アクセス制御（オーナー/店長/スタッフ/開発者）",
            "操作履歴・シフト変更ログで完全な監査証跡",
            "SSL暗号化通信・CSRF保護を標準装備",
            "PCI DSS準拠の外部決済サービスと連携",
        ], ACCENT_ORANGE)

    # ==========================================
    # スライド 8〜: 画面イメージ（1画面1スライド）
    # 構成: メインページ → スマホ → 管理画面
    # ==========================================
    screenshots_dir = os.path.dirname(os.path.abspath(__file__))

    def _add_desktop_slide(img_name, title, subtitle, bullets, header_prefix=""):
        """デスクトップ画面スライド: 左にスクショ、右に説明"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        full_title = f"{header_prefix}{title}" if header_prefix else title
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 full_title, font_size=28, bold=True, color=BROWN)
        add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)
        img_path = os.path.join(screenshots_dir, img_name)
        add_shape(slide, Inches(0.4), Inches(1.2), Inches(8.2), Inches(5.8),
                  LIGHT_GRAY, corner_radius=True)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.6))
            except Exception:
                add_text(slide, Inches(2), Inches(3.5), Inches(4), Inches(1),
                         f'[{title}]', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(0.5),
                 subtitle, font_size=18, bold=True, color=DARK)
        add_bullet_list(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(5),
                        bullets, font_size=12, icon="✓")

    # ── 1. お客様向け画面（メインページ）──
    frontend_screens = [
        ('screenshots/front_top_desktop.png', 'トップページ',
         'お客様が最初に目にする画面',
         [
             'ヒーローバナーで店舗の魅力をアピール',
             '予約方法を選べるわかりやすい導線',
             'スタッフランキングで人気キャストをPR',
             'お知らせ・メディア掲載を自動表示',
             '7言語対応でインバウンド集客にも対応',
         ]),
        ('screenshots/front_staff_list_desktop.png', 'スタッフ紹介・店舗情報',
         'お店のスタッフを魅力的に紹介',
         [
             'スタッフの写真・プロフィールを掲載',
             '得意メニュー・資格を表示',
             'そのまま予約カレンダーへ遷移',
             '店舗情報（営業時間・地図・アクセス）を表示',
             '管理画面から簡単に内容を更新',
         ]),
        ('screenshots/front_news_desktop.png', 'お知らせ',
         '最新情報をお客様に発信',
         [
             '新メニュー・キャンペーン情報を簡単投稿',
             '管理画面からワンクリックで公開',
             '日付順に自動ソートで常に最新情報を表示',
             'SNS連携で情報拡散をサポート',
             'お客様のリピート来店を促進',
         ]),
        ('screenshots/front_shop_desktop.png', 'オンラインショップ',
         'ECで売上チャネルを拡大',
         [
             '店舗の商品をオンラインで販売',
             '商品検索・カテゴリ絞り込み機能',
             'カート機能で複数商品をまとめて購入',
             'クレジットカード・電子マネーで決済',
             '店頭在庫と一元管理で在庫ミスを防止',
         ]),
        ('screenshots/table_order_vp.png', 'QRテーブル注文',
         'スマホで簡単オーダー',
         [
             'テーブルのQRコードをスキャンするだけ',
             '写真付きメニューで料理のイメージが伝わる',
             'カテゴリ別タブで欲しい商品がすぐ見つかる',
             'カート機能で追加注文もラクラク',
             '注文状況をリアルタイムで確認',
         ]),
        ('screenshots/booking_calendar_desktop.png', '予約カレンダー',
         'オンラインで24時間予約受付',
         [
             'キャスト別の空き状況をカレンダーで一覧表示',
             'お客様が好きな日時・時間帯を選んで予約',
             '前払い決済で予約確定（ドタキャン防止）',
             'LINE・メールでの予約にも対応',
             '予約確定時にスタッフへ自動通知',
         ]),
    ]

    for img_name, title, subtitle, bullets in frontend_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="お客様向け画面 — ")

    # ── 2. スマホ対応（トップページ + 管理画面）──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "スマートフォン完全対応", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "お客様向けサイトも管理画面も、スマホからそのまま操作できます。",
             font_size=16, color=DARK)

    mobile_items = [
        ('screenshots/front_top_mobile_vp.png', 'お客様向けトップページ'),
        ('screenshots/admin_dashboard_mobile_vp.png', '管理画面ダッシュボード'),
    ]

    for i, (img_name, label) in enumerate(mobile_items):
        x = Inches(1.5) + i * Inches(5.5)
        y = Inches(1.8)

        add_shape(slide, x, y, Inches(4.2), Inches(5.0), WHITE, corner_radius=True)

        img_path = os.path.join(screenshots_dir, img_name)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                x + Inches(0.6), y + Inches(0.15),
                                Inches(3.0), Inches(4.4))
            except Exception:
                add_text(slide, x + Inches(0.5), y + Inches(2), Inches(3.2), Inches(1),
                         f'[{label}]', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(4.5), Inches(4.2), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ── 3. 管理画面（デスクトップ）──
    admin_screens = [
        ('screenshots/dashboard_sales_vp.png', '売上ダッシュボード',
         '経営状況をひと目で把握',
         [
             '日別・週別・月別の売上推移をグラフで表示',
             '予約KPI（件数・キャンセル率）をリアルタイム集計',
             'メニューエンジニアリングで人気商品・利益率を分析',
             'KPIスコアカードで主要指標を一覧表示',
             'ビジネスインサイトが改善機会を自動で提案',
         ]),
        ('screenshots/shift_calendar_vp.png', 'シフトカレンダー',
         'シフト管理を効率化',
         [
             'スタッフのシフト希望をスマホから受付',
             'カレンダー形式で全スタッフの予定を俯瞰',
             '自動スケジューリングでワンクリック割り当て',
             '人員不足の日を自動検知しアラート表示',
             '確定シフトをスタッフにLINEで自動通知',
         ]),
        ('screenshots/pos_vp.png', 'POS レジ',
         '会計をスムーズに',
         [
             'カテゴリ別タブで商品をすばやく選択',
             '現金・クレジットカード・PayPay・交通系ICに対応',
             'レシート自動発行・印刷対応',
             'キッチンディスプレイと連動してオーダー通知',
             '日次・月次の売上レポートを自動生成',
         ]),
        ('screenshots/inventory_vp.png', '在庫管理',
         '在庫切れを防止',
         [
             '全商品の在庫数をリアルタイムで一覧表示',
             '注文と連動して在庫を自動で差し引き',
             '在庫が少なくなったらアラート通知',
             'QRコードで入庫処理（バーコードリーダー不要）',
             '入出庫の履歴をすべて記録（監査対応）',
         ]),
        ('screenshots/customer_feedback_vp.png', '顧客フィードバック（NPS）',
         'お客様の声を数値で把握',
         [
             'NPS（顧客推奨度）を自動集計・トレンド表示',
             'Promoter / Passive / Detractor を色分け表示',
             'コメント付きフィードバックで具体的な改善点を把握',
             '注文データと紐付けてサービス品質を分析',
             '月別・スタッフ別の満足度推移を可視化',
         ]),
        ('screenshots/attendance_board_vp.png', '出退勤ボード',
         'スタッフの勤務状況をリアルタイム表示',
         [
             '出勤中・休憩中・未出勤のステータスを一覧表示',
             '30秒ごとに自動更新でリアルタイム把握',
             'QRコード / PIN入力 / スマホ打刻の3方式に対応',
             '残業・深夜・休日出勤を自動で分類',
             '勤怠データから給与計算へシームレスに連携',
         ]),
    ]

    for img_name, title, subtitle, bullets in admin_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="管理画面 — ")

    # ==========================================
    # スライド 10: 導入メリット
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "導入メリット", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    merits = [
        ("ドタキャン損失ゼロへ",
         "前払い予約確定制で無断キャンセルを防止。\n前払い/後払いはワンクリックで切替可能です。",
         "🛡"),
        ("LINEで予約がラクラク",
         "お客様はLINEのトーク画面から予約完了。\n新しいアプリのインストールは一切不要です。",
         "💬"),
        ("売上アップに貢献",
         "QR注文で追加オーダー率UP。\n多言語対応でインバウンド客も取りこぼしません。",
         "📈"),
        ("業務時間を大幅削減",
         "予約受付・注文・シフト管理・給与計算を自動化。\nスタッフが接客に集中できる環境を作ります。",
         "⏱"),
        ("コスト削減",
         "予約台帳・注文伝票・タイムカード…\n紙のコストと手間をゼロにします。",
         "💰"),
        ("データで経営判断",
         "売上・予約・スタッフ実績をダッシュボードで可視化。\nデータに基づいた意思決定をサポートします。",
         "📊"),
    ]

    for i, (title, desc, icon) in enumerate(merits):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), WHITE, corner_radius=True)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                 icon, font_size=30, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.35),
                 title, font_size=16, bold=True, color=BROWN)

        add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # スライド 8: ご利用の流れ
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "ご利用開始までの流れ", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    steps = [
        ("STEP 1", "お問い合わせ", "お電話またはフォームから\nお気軽にご連絡ください"),
        ("STEP 2", "ヒアリング", "現在の業務フローや\n課題をお伺いします"),
        ("STEP 3", "デモ・お見積り", "実際の画面をお見せしながら\n最適なプランをご提案"),
        ("STEP 4", "アカウント発行", "店舗情報・スタッフを登録\n初期設定をサポート"),
        ("STEP 5", "運用開始", "操作方法のレクチャー後\nすぐにご利用開始！"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.6)

        # ステップ番号（丸）
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BROWN
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Meiryo'
        p.alignment = PP_ALIGN.CENTER

        # 矢印（最後以外）
        if i < len(steps) - 1:
            add_text(slide, x + Inches(1.9), y + Inches(0.2), Inches(0.6), Inches(0.6),
                     "→", font_size=28, color=BROWN, align=PP_ALIGN.CENTER)

        # タイトル
        add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.2), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # 説明
        add_text(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # 補足テキスト
    add_shape(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(1.5), LIGHT_GRAY, corner_radius=True)
    add_text(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1.2),
             "最短即日でご利用開始いただけます。\n"
             "専用機器の購入やシステム工事は一切不要です。お手持ちのスマホ・タブレット・PCからすぐにアクセスできます。\n"
             "導入後も操作方法のサポートや機能アップデートを継続的にご提供します。",
             font_size=13, color=DARK)

    # ==========================================
    # スライド 9: 料金プラン
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "料金プラン", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "※ 以下は参考価格です。ご要望に応じてカスタマイズいたします。",
             font_size=12, color=GRAY)

    plans = [
        ("ライトプラン", "月額 ¥○○,○○○〜", ACCENT_GREEN,
         "小規模店舗向け", [
             "予約管理（LINE・メール）",
             "スタッフ管理（5名まで）",
             "ホームページ管理",
             "メール・チャットサポート",
         ]),
        ("スタンダードプラン", "月額 ¥○○,○○○〜", ACCENT_BLUE,
         "おすすめ！", [
             "ライトプランの全機能",
             "QRテーブル注文",
             "在庫管理",
             "シフト管理",
             "スタッフ無制限",
             "電話サポート",
         ]),
        ("プレミアムプラン", "月額 ¥○○,○○○〜", BROWN,
         "複数店舗・フル機能", [
             "スタンダードの全機能",
             "給与計算・勤怠管理",
             "売上ダッシュボード",
             "多言語対応（7言語）",
             "複数店舗管理",
             "優先サポート",
             "カスタマイズ対応",
         ]),
    ]

    for i, (name, price, color, badge, features) in enumerate(plans):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.2)

        # カード背景
        card = add_shape(slide, x, y, w, h, LIGHT_GRAY, corner_radius=True)

        # 上部アクセント
        add_shape(slide, x, y, w, Inches(0.08), color)

        # おすすめバッジ
        if badge == "おすすめ！":
            badge_shape = add_shape(slide, x + Inches(0.8), y - Inches(0.15), Inches(2.2), Inches(0.35), color, corner_radius=True)
            tf = badge_shape.text_frame
            p = tf.paragraphs[0]
            p.text = "★ おすすめ！"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Meiryo'
            p.alignment = PP_ALIGN.CENTER

        # プラン名
        add_text(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # 価格
        add_text(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.5),
                 price, font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

        # 区切り線
        add_shape(slide, x + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(0.02), GRAY)

        # 対象
        add_text(slide, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(0.3),
                 badge if badge != "おすすめ！" else "中規模店舗向け",
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        # 機能リスト
        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9),
                        w - Inches(0.6), h - Inches(2.1),
                        features, font_size=12, icon="✓")

    # 注意書き
    add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "※ 初期費用・カスタマイズ費用は別途お見積りとなります。  ※ 表示価格は税別です。",
             font_size=10, color=GRAY)

    # ==========================================
    # スライド 10: よくあるご質問
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "よくあるご質問", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    faqs = [
        ("Q. 導入に必要な機器はありますか？",
         "A. いいえ、専用機器は不要です。お手持ちのスマートフォン、タブレット、PCからすぐにご利用いただけます。\n"
         "　 テーブル注文のQRコードは、管理画面から生成して印刷するだけです。"),
        ("Q. ITに詳しくないスタッフでも使えますか？",
         "A. はい。スマホが使える方なら直感的に操作できるシンプルな画面設計です。\n"
         "　 導入時には操作レクチャーも実施いたします。"),
        ("Q. 途中で機能を追加・変更できますか？",
         "A. はい。プランの変更はいつでも可能です。\n"
         "　 まずは基本機能からスタートして、必要に応じて機能を追加していくことをおすすめしています。"),
        ("Q. データのセキュリティは大丈夫ですか？",
         "A. お客様の個人情報は暗号化して保管しています。通信はすべてSSL暗号化で保護されており、\n"
         "　 クレジットカード情報は外部の決済サービス（PCI DSS準拠）で処理するため、店舗側にカード情報は残りません。"),
        ("Q. 複数店舗で使えますか？",
         "A. はい。プレミアムプランでは複数店舗を一元管理できます。\n"
         "　 店舗ごとにスタッフや在庫を分けて管理しながら、全体の売上をまとめて確認できます。"),
    ]

    for i, (q, a) in enumerate(faqs):
        y = Inches(1.2) + i * Inches(1.2)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(1.05), WHITE, corner_radius=True)

        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(11.73), Inches(0.3),
                 q, font_size=13, bold=True, color=BROWN)
        add_text(slide, Inches(0.8), y + Inches(0.4), Inches(11.73), Inches(0.6),
                 a, font_size=11, color=DARK)

    # ==========================================
    # スライド 11: お問い合わせ
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "まずはお気軽にお問い合わせください",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             "デモのご依頼・お見積り・ご質問など、なんでもお気軽にどうぞ。",
             font_size=18, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.33), Inches(0.03), WHITE)

    # 連絡先カード
    contact_card = add_shape(slide, Inches(3), Inches(3.8), Inches(7.33), Inches(2.5), WHITE, corner_radius=True)

    contacts = [
        "TEL:     000-0000-0000",
        "Email:   info@example.com",
        "Web:    https://timebaibai.com",
        "",
        "株式会社○○○○",
        "〒000-0000  東京都○○区○○ 0-0-0",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, Inches(4), Inches(4.0) + i * Inches(0.35),
                 Inches(5), Inches(0.35),
                 line, font_size=15 if i < 3 else 13,
                 bold=(i < 3), color=DARK if i < 3 else GRAY,
                 align=PP_ALIGN.LEFT)

    # フッター
    add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "© 2026 TimeBaiBai. All Rights Reserved.",
             font_size=10, color=BEIGE, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# 出力
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    # PPT出力
    pptx_path = os.path.join(output_dir, 'TimeBaiBai_サービス紹介資料.pptx')
    prs.save(pptx_path)
    print(f"✅ PPT saved: {pptx_path}")

    # PDF出力を試みる（LibreOffice利用）
    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        # macOSの場合
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            print(f"✅ PDF saved: {pdf_path}")
        else:
            # LibreOffice がない場合のフォールバック
            raise FileNotFoundError("LibreOffice not found")
    except Exception as e:
        print(f"⚠️  PDF conversion skipped ({e})")
        print(f"   PPTXファイルからPDFを作成するには:")
        print(f"   - PowerPoint / Keynote で開いて「PDF として書き出し」")
        print(f"   - または: brew install libreoffice && soffice --headless --convert-to pdf '{pptx_path}'")
