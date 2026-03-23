#!/usr/bin/env python3
"""
Sales Deck Generator — Simplified Chinese (简体中文) Version
Outputs PPT + PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ==============================
# Color Palette
# ==============================
BROWN = RGBColor(0x8C, 0x87, 0x6C)
BEIGE = RGBColor(0xF1, 0xF0, 0xEC)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)

FONT_NAME = 'Microsoft YaHei'


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="\u2713"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = spacing


def add_feature_card(slide, left, top, width, height, title, items, accent_color=BROWN):
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=True)
    line = add_shape(slide, left + Inches(0.15), top + Inches(0.08),
                     width - Inches(0.3), Inches(0.04), accent_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.2),
             width - Inches(0.4), Inches(0.4),
             title, font_size=15, bold=True, color=accent_color)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6),
                    width - Inches(0.4), height - Inches(0.8),
                    items, font_size=11, icon="\u25CF")


def add_picture_fit(slide, img_path, left, top, max_width, max_height):
    with Image.open(img_path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = max_width / max_height
    if ratio > box_ratio:
        w = max_width
        h = max_width / ratio
    else:
        h = max_height
        w = max_height * ratio
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    return slide.shapes.add_picture(img_path, int(x), int(y), width=int(w), height=int(h))


def add_number_highlight(slide, left, top, number, label, color=BROWN):
    add_text(slide, left, top, Inches(2), Inches(0.6),
             number, font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2), Inches(0.3),
             label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ==============================
# Slide Creation
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # Slide 1: Cover
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "全方位门店DX平台",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "预约 \u2022 点餐 \u2022 支付 \u2022 员工管理 \u2022 工资\n一个系统全部搞定。",
             font_size=22, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5), Inches(4.0), Inches(3.33), Inches(0.03), WHITE)

    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "TimeBaiBai  服务介绍",
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "2026  /  公司名称",
             font_size=14, color=BEIGE, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 2: Pain Points
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "这些问题，您是否深有同感？", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    problems = [
        ("预约渠道分散",
         "电话、LINE、邮件……各种预约渠道\n总担心重复预约。",
         "\U0001F4DE"),
        ("爽约与临时取消",
         "客人预约了却不来。\n浪费了准备工作和营收。",
         "\U0001F630"),
        ("点餐手忙脚乱",
         "高峰时段忙不过来。\n点餐错误频频发生。",
         "\U0001F4DD"),
        ("库存盲区",
         "商品莫名就缺货了。\n不知道什么时候该补货。",
         "\U0001F4E6"),
    ]

    for i, (title, desc, icon) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.6) + row * Inches(2.6)

        card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), LIGHT_GRAY, corner_radius=True)

        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8),
                 icon, font_size=36, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK)

        add_text(slide, x + Inches(1.2), y + Inches(0.75), Inches(4), Inches(1.2),
                 desc, font_size=14, color=GRAY)

    # ==========================================
    # Slide 3: Solution — What is TimeBaiBai
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "TimeBaiBai 帮您全部解决", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "一站式平台，满足门店所有需求。只需一部智能手机即可开始使用。",
             font_size=16, color=DARK)

    highlights = [
        ("6+", "核心模块"),
        ("7", "支持语言"),
        ("24h", "自动预约"),
        ("$0", "省去纸张成本"),
    ]
    for i, (num, label) in enumerate(highlights):
        x = Inches(1.2) + i * Inches(3)
        add_shape(slide, x - Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.4), WHITE, corner_radius=True)
        add_number_highlight(slide, x, Inches(2.5), num, label)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
             "预约、桌边点餐、库存、排班管理、工资计算、官网管理……\n"
             "通过单一仪表板管理所有业务——无需再使用多种工具。\n\n"
             "设置简单，开通账号当天即可使用。\n"
             "无需复杂安装或硬件设备。只需通过手机或电脑访问即可。",
             font_size=15, color=DARK)

    # ==========================================
    # Slide 4: Features (1/4) Reservations & Orders
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "核心功能  \u2460  预约 / 桌边点餐",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "在线预约系统", [
            "通过LINE预约——在聊天中即可完成预约",
            "同时支持邮件预约，方便客户选择",
            "预付确认预约——支付后自动确认（防止爽约）",
            "管理后台一键切换预付/后付模式",
            "员工日历实时显示可预约时段",
            "两种预约方式：按日期 或 按员工",
            "二维码签到（零前台工作量）",
            "预约确认后自动通过LINE通知员工",
            "支持7种语言（日语、英语、中文、韩语、西班牙语、葡萄牙语）",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "二维码桌边点餐系统", [
            "在每张桌子上放置二维码即可开始使用",
            "客人扫码 \u2192 浏览菜单 \u2192 下单",
            "图片菜单帮助客人直观了解菜品",
            "分类标签快速查找菜品",
            "购物车功能方便追加点餐",
            "实时订单状态追踪",
            "支持现金、信用卡和电子支付",
            "管理后台批量生成二维码并可打印下载",
        ], ACCENT_GREEN)

    # ==========================================
    # Slide 5: Features (2/4) Inventory & Staff
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "核心功能  \u2461  库存 / 员工排班管理",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "库存管理", [
            "所有商品库存实时追踪",
            "订单联动自动扣减库存",
            "库存不足自动提醒",
            "二维码收货（无需条码扫描器）",
            "完整的出入库记录（审计就绪）",
            "按商品类别整理，管理轻松",
            "EC网店与门店库存统一管理",
            "盘点功能支持批量调整",
        ], ACCENT_ORANGE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "员工排班管理", [
            "员工通过手机提交排班意愿",
            "三级意愿：可以 / 希望 / 不可以",
            "管理者通过管理后台确定排班",
            "自动排班实现高效人员配置",
            "确定的排班自动通过LINE通知员工",
            "员工资料含照片和简介",
            "员工类型：接待人员 / 门店员工",
            "权限管理：老板 / 经理 / 员工 / 开发者",
        ], ACCENT_BLUE)

    # ==========================================
    # Slide 6: Features (3/4) Payroll & Website
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "核心功能  \u2462  工资 / 官网管理",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "工资与考勤管理", [
            "根据排班记录自动生成考勤数据",
            "自动区分普通工时、加班、夜班和节假日工时",
            "支持时薪和月薪员工的工资计算",
            "自动计算社会保险费",
            "自动适用所得税和住民税扣除",
            "一键生成工资单",
            "银行转账CSV导出（全银格式）",
            "轻松设置各类津贴（通勤、住房、家庭）",
        ], ACCENT_GREEN)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "官网与营销管理", [
            "通过管理后台编辑门店官网",
            "首页横幅（轮播图）展示门店魅力",
            "轻松更新新闻和媒体报道",
            "广告横幅投放与管理",
            "社交媒体集成（X / Instagram嵌入）",
            "自定义HTML模块实现灵活布局",
            "管理后台可编辑隐私政策和法律页面",
            "内置EC网店支持在线销售",
        ], ACCENT_ORANGE)

    # ==========================================
    # Slide 7: Features (4/4) Analytics & Security
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "核心功能  \u2463  数据分析 / 信息安全",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "数据分析与AI", [
            "销售仪表板：日/周/月/年趋势图表",
            "菜单工程学（Star/Dog分类）助力利润优化",
            "ABC分析、RFM分析、同期群分析、购物篮分析",
            "NPS（客户满意度）自动汇总，助力客户留存",
            "AI预测各时段最优人员配置",
            "需求与营收预测，助力采购和排班决策",
            "商业洞察自动检测异常和机会",
            "KPI记分卡一目了然掌握经营状况",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "信息安全与审计", [
            "完整的访问日志（登录、操作、异常检测）",
            "基于IP的速率限制自动拦截未授权访问",
            "客户个人信息采用AES加密安全存储",
            "自动安全审计（每日12项检查）",
            "基于角色的访问控制（老板 / 经理 / 员工 / 开发者）",
            "完整的审计追踪：操作历史和排班变更记录",
            "内置SSL加密和CSRF保护",
            "集成PCI DSS合规的外部支付服务",
        ], ACCENT_ORANGE)

    # ==========================================
    # Screenshot Slides
    # ==========================================
    screenshots_dir = os.path.dirname(os.path.abspath(__file__))

    def _add_desktop_slide(img_name, title, subtitle, bullets, header_prefix=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        full_title = f"{header_prefix}{title}" if header_prefix else title
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 full_title, font_size=28, bold=True, color=BROWN)
        add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)
        img_path = os.path.join(screenshots_dir, img_name)
        add_shape(slide, Inches(0.4), Inches(1.2), Inches(8.2), Inches(5.8),
                  LIGHT_GRAY, corner_radius=True)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.6))
            except Exception:
                add_text(slide, Inches(2), Inches(3.5), Inches(4), Inches(1),
                         f'[{title}]', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(0.5),
                 subtitle, font_size=18, bold=True, color=DARK)
        add_bullet_list(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(5),
                        bullets, font_size=12, icon="\u2713")

    # ── Customer-Facing Screens ──
    frontend_screens = [
        ('screenshots/zh-hans/front_top_desktop.png', '首页',
         '客人看到的第一个页面',
         [
             '首页横幅展示门店魅力',
             '清晰的预约流程，多种选择',
             '员工排名推广人气员工',
             '自动展示新闻和媒体报道',
             '支持7种语言，服务国际客户',
         ]),
        ('screenshots/zh-hans/front_staff_list_desktop.png', '员工介绍与门店信息',
         '精彩展示您的员工',
         [
             '员工照片和个人简介',
             '展示专长和资质',
             '直接链接到预约日历',
             '门店信息（营业时间、地图、交通方式）',
             '管理后台轻松更新内容',
         ]),
        ('screenshots/zh-hans/front_news_desktop.png', '新闻与公告',
         '让客户随时了解最新动态',
         [
             '发布新菜品、活动和促销信息',
             '管理后台一键发布',
             '按日期自动排序，保持内容新鲜',
             '社交媒体联动扩大传播',
             '通过及时更新促进回头客',
         ]),
        ('screenshots/zh-hans/front_shop_desktop.png', '在线商城',
         '拓展销售渠道',
         [
             '在线销售您的商品',
             '商品搜索和分类筛选',
             '购物车支持多件商品购买',
             '支持信用卡和电子支付',
             '统一库存防止库存差异',
         ]),
        ('screenshots/zh-hans/table_order_vp.png', '二维码桌边点餐',
         '手机轻松下单',
         [
             '只需扫描桌上的二维码',
             '图片菜单直观展示菜品',
             '分类标签快速查找',
             '购物车方便追加点餐',
             '实时订单状态追踪',
         ]),
        ('screenshots/zh-hans/booking_calendar_desktop.png', '预约日历',
         '24小时在线预约',
         [
             '日历一目了然查看员工空闲时段',
             '客人选择心仪的日期和时间',
             '预付确认防止爽约',
             '支持LINE和邮件预约',
             '确认后自动通知员工',
         ]),
    ]

    for img_name, title, subtitle, bullets in frontend_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="客户页面 \u2014 ")

    # ── Mobile-Responsive ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "完美支持手机浏览", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "客户端和管理后台均可在智能手机上完美运行。",
             font_size=16, color=DARK)

    mobile_items = [
        ('screenshots/zh-hans/front_top_mobile_vp.png', '客户首页'),
        ('screenshots/zh-hans/admin_dashboard_mobile_vp.png', '管理后台'),
    ]

    for i, (img_name, label) in enumerate(mobile_items):
        x = Inches(1.5) + i * Inches(5.5)
        y = Inches(1.8)

        add_shape(slide, x, y, Inches(4.2), Inches(5.0), WHITE, corner_radius=True)

        img_path = os.path.join(screenshots_dir, img_name)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                x + Inches(0.6), y + Inches(0.15),
                                Inches(3.0), Inches(4.4))
            except Exception:
                add_text(slide, x + Inches(0.5), y + Inches(2), Inches(3.2), Inches(1),
                         f'[{label}]', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(4.5), Inches(4.2), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ── Admin Screens ──
    admin_screens = [
        ('screenshots/zh-hans/dashboard_sales_vp.png', '销售仪表板',
         '一目了然的经营概况',
         [
             '日/周/月/年营收趋势',
             '预约KPI（数量、取消率）实时监控',
             '三渠道分解：合计、EC网店、店内菜单',
             'KPI记分卡概览关键指标',
             '商业洞察自动建议改进措施',
         ]),
        ('screenshots/zh-hans/dashboard_menu_eng_vp.png', '销售分析（菜单分析）',
         'AI驱动的利润优化',
         [
             '菜单工程学：Star / Plowhorse / Puzzle / Dog 四象限分类',
             'ABC分析：帕累托法则识别核心营收驱动',
             '营收预测：Prophet AI模型预测未来14天',
             '每小时销售热力图助力最优排班',
             'AOV（客单价）趋势用于定价策略评估',
         ]),
        ('screenshots/zh-hans/dashboard_rfm_vp.png', '销售分析（客户分析）',
         '数据驱动的客户洞察',
         [
             'RFM分析：按最近消费、消费频率、消费金额自动客户分群',
             '同期群分析：可视化每月新客户留存率',
             '购物篮分析：发现交叉销售模式和推荐',
             'AI摘要：分析结果以自然语言总结',
             '推荐行动：AI自动建议具体改进步骤',
         ]),
        ('screenshots/zh-hans/shift_calendar_vp.png', '排班日历',
         '高效的排班管理',
         [
             '员工通过手机提交排班意愿',
             '日历视图一目了然查看所有员工排班',
             '自动排班一键分配',
             '自动检测人手不足并发出提醒',
             '确定的排班自动通过LINE通知',
         ]),
        ('screenshots/zh-hans/pos_vp.png', 'POS收银',
         '顺畅的结账体验',
         [
             '分类标签快速选择商品',
             '支持现金、信用卡、PayPay和交通IC卡',
             '自动生成和打印收据',
             '厨房显示屏联动订单通知',
             '自动生成日报和月报',
         ]),
        ('screenshots/zh-hans/inventory_vp.png', '库存管理',
         '防止缺货',
         [
             '所有商品库存实时监控',
             '订单接收自动扣减',
             '库存不足提醒通知',
             '二维码收货（无需条码扫描器）',
             '完整的出入库记录便于审计',
         ]),
        ('screenshots/zh-hans/customer_feedback_vp.png', '客户反馈（NPS）',
         '量化客户心声',
         [
             'NPS（净推荐值）自动汇总和趋势分析',
             '推荐者/被动者/贬损者彩色标注显示',
             '基于评论的反馈助力改进',
             '关联订单数据进行服务质量分析',
             '月度和员工级别满意度趋势',
         ]),
        ('screenshots/zh-hans/attendance_board_vp.png', '考勤看板',
         '实时员工出勤状况',
         [
             '一目了然查看在岗、休息中、未打卡状态',
             '每30秒自动刷新实时追踪',
             '三种打卡方式：二维码、PIN码或手机',
             '自动区分加班、夜班和节假日工时',
             '考勤数据无缝对接工资系统',
         ]),
        ('screenshots/zh-hans/iot_sensors_vp.png', 'IoT传感器监控',
         '实时门店环境监测',
         [
             '温度、湿度、气压、气体浓度实时监测',
             'PIR人体传感器自动来客计数',
             '异常情况（燃气泄漏、高温等）告警通知',
             '时间序列图表追踪传感器数据趋势',
             'ESP32设备连接状态监控',
         ]),
    ]

    for img_name, title, subtitle, bullets in admin_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="管理后台 \u2014 ")

    # ==========================================
    # Benefits Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "导入效益", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    merits = [
        ("零爽约损失",
         "预付确认预约消除临时取消。\n管理后台一键切换预付/后付模式。",
         "\U0001F6E1"),
        ("LINE轻松预约",
         "客人通过LINE聊天直接预约。\n无需安装新应用。",
         "\U0001F4AC"),
        ("提升营收",
         "二维码点餐提高追加点餐率。\n多语言支持吸引国际客户。",
         "\U0001F4C8"),
        ("节省工时",
         "预约、点餐、排班、工资全自动化。\n让员工专注于客户服务。",
         "\u23F1"),
        ("降低成本",
         "预约簿、点餐单、考勤卡……\n告别纸张成本和人工作业。",
         "\U0001F4B0"),
        ("数据驱动决策",
         "可视化销售、预约和员工绩效。\n基于真实数据做出明智决策。",
         "\U0001F4CA"),
    ]

    for i, (title, desc, icon) in enumerate(merits):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), WHITE, corner_radius=True)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                 icon, font_size=30, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.35),
                 title, font_size=16, bold=True, color=BROWN)

        add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # Getting Started Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "导入流程", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    steps = [
        ("STEP 1", "联系我们", "通过电话或\n咨询表单联系我们"),
        ("STEP 2", "需求咨询", "了解您当前的\n工作流程和痛点"),
        ("STEP 3", "演示与报价", "体验产品功能。\n我们为您推荐最合适的方案"),
        ("STEP 4", "账号设置", "注册门店和员工信息。\n我们协助完成初始设置"),
        ("STEP 5", "正式上线", "经过简单培训后，\n即可正式投入使用！"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.6)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BROWN
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            add_text(slide, x + Inches(1.9), y + Inches(0.2), Inches(0.6), Inches(0.6),
                     "\u2192", font_size=28, color=BROWN, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.2), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(1.5), LIGHT_GRAY, corner_radius=True)
    add_text(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1.2),
             "最快当天即可开始使用。\n"
             "无需专用硬件或系统安装。通过智能手机、平板电脑或电脑即可立即访问。\n"
             "部署后我们持续提供技术支持、培训和功能更新。",
             font_size=13, color=DARK)

    # ==========================================
    # Pricing Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "方案与价格", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "* 以下价格仅供参考。我们会根据您的需求定制方案。",
             font_size=12, color=GRAY)

    plans = [
        ("轻量方案", "从 $XX/月 起", ACCENT_GREEN,
         "适合小型门店", [
             "预约管理（LINE和邮件）",
             "员工管理（最多5人）",
             "官网管理",
             "邮件和在线客服支持",
         ]),
        ("标准方案", "从 $XX/月 起", ACCENT_BLUE,
         "推荐！", [
             "包含轻量方案全部功能",
             "二维码桌边点餐",
             "库存管理",
             "排班管理",
             "员工人数不限",
             "电话支持",
         ]),
        ("高级方案", "从 $XX/月 起", BROWN,
         "多店铺 / 全功能", [
             "包含标准方案全部功能",
             "工资与考勤管理",
             "销售仪表板与数据分析",
             "多语言支持（7种语言）",
             "多店铺管理",
             "优先支持",
             "定制开发",
         ]),
    ]

    for i, (name, price, color, badge, features) in enumerate(plans):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.2)

        card = add_shape(slide, x, y, w, h, LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, w, Inches(0.08), color)

        if badge == "推荐！":
            badge_shape = add_shape(slide, x + Inches(0.8), y - Inches(0.15), Inches(2.2), Inches(0.35), color, corner_radius=True)
            tf = badge_shape.text_frame
            p = tf.paragraphs[0]
            p.text = "\u2605 推荐！"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER

        add_text(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.5),
                 price, font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

        add_shape(slide, x + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(0.02), GRAY)

        add_text(slide, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(0.3),
                 badge if badge != "推荐！" else "适合中型门店",
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9),
                        w - Inches(0.6), h - Inches(2.1),
                        features, font_size=12, icon="\u2713")

    add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "* 初始设置费和定制开发费另行报价。 * 价格不含税。",
             font_size=10, color=GRAY)

    # ==========================================
    # FAQ Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "常见问题", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    faqs = [
        ("Q. 需要专用设备吗？",
         "A. 不需要任何专用硬件。使用您现有的智能手机、平板电脑或电脑即可。\n"
         "   桌边点餐只需从管理后台打印二维码即可。"),
        ("Q. 不懂技术的员工也能用吗？",
         "A. 当然可以。只要会用智能手机，就能使用TimeBaiBai。界面直观简洁。\n"
         "   导入时我们还会提供实操培训。"),
        ("Q. 之后可以增加或更改功能吗？",
         "A. 完全可以。您可以随时升级或更换方案。\n"
         "   我们建议先从核心功能开始，根据需要逐步扩展。"),
        ("Q. 数据安全有保障吗？",
         "A. 客户数据在存储时全部加密。所有通信使用SSL加密。\n"
         "   信用卡信息由PCI DSS合规的支付服务处理——我们的服务器不存储任何卡片数据。"),
        ("Q. 可以用于多个店铺吗？",
         "A. 可以。高级方案支持多店铺管理。\n"
         "   按店铺管理员工和库存，同时查看所有店铺的汇总报告。"),
    ]

    for i, (q, a) in enumerate(faqs):
        y = Inches(1.2) + i * Inches(1.2)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(1.05), WHITE, corner_radius=True)

        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(11.73), Inches(0.3),
                 q, font_size=13, bold=True, color=BROWN)
        add_text(slide, Inches(0.8), y + Inches(0.4), Inches(11.73), Inches(0.6),
                 a, font_size=11, color=DARK)

    # ==========================================
    # Contact Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "联系我们",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             "欢迎咨询演示、报价或任何问题。",
             font_size=18, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.33), Inches(0.03), WHITE)

    contact_card = add_shape(slide, Inches(3), Inches(3.8), Inches(7.33), Inches(2.5), WHITE, corner_radius=True)

    contacts = [
        "TEL:     000-0000-0000",
        "Email:   info@example.com",
        "Web:    https://timebaibai.com",
        "",
        "公司名称",
        "地址信息",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, Inches(4), Inches(4.0) + i * Inches(0.35),
                 Inches(5), Inches(0.35),
                 line, font_size=15 if i < 3 else 13,
                 bold=(i < 3), color=DARK if i < 3 else GRAY,
                 align=PP_ALIGN.LEFT)

    add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "\u00A9 2026 TimeBaiBai. All Rights Reserved.",
             font_size=10, color=BEIGE, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# Output
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    pptx_path = os.path.join(output_dir, 'TimeBaiBai_Service_Overview_ZH_Hans.pptx')
    prs.save(pptx_path)
    print(f"\u2705 PPT saved: {pptx_path}")

    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            print(f"\u2705 PDF saved: {pdf_path}")
        else:
            raise FileNotFoundError("LibreOffice not found")
    except Exception as e:
        print(f"\u26A0\uFE0F  PDF conversion skipped ({e})")
        print(f"   To create PDF from PPTX:")
        print(f"   - Open in PowerPoint / Keynote and export as PDF")
        print(f"   - Or: brew install libreoffice && soffice --headless --convert-to pdf '{pptx_path}'")
