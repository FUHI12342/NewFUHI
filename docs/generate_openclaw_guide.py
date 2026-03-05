#!/usr/bin/env python3
"""
OpenClaw & ClawHub 解説資料 生成スクリプト
図を用いた初心者向けPDF資料
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==============================
# カラーパレット (OpenClaw = ロブスター/レッド系)
# ==============================
CLAW_RED = RGBColor(0xE0, 0x3E, 0x3E)       # OpenClawブランド赤
CLAW_DARK = RGBColor(0x1A, 0x1A, 0x2E)      # ダーク背景
NAVY = RGBColor(0x2D, 0x3A, 0x4A)           # ダークネイビー
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF8)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
GRAY = RGBColor(0x88, 0x88, 0x88)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
ACCENT_TEAL = RGBColor(0x1A, 0xBC, 0x9C)
SOFT_RED_BG = RGBColor(0xFD, 0xF0, 0xF0)
SOFT_BLUE_BG = RGBColor(0xF0, 0xF4, 0xFD)
SOFT_GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
SOFT_ORANGE_BG = RGBColor(0xFD, 0xF7, 0xF0)
SOFT_PURPLE_BG = RGBColor(0xF5, 0xF0, 0xFD)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bordered_shape(slide, left, top, width, height, fill_color, border_color, border_width=Pt(2), corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Meiryo'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_multiline_text(slide, left, top, width, height, lines, font_size=14,
                       color=DARK, spacing=Pt(4), font_name='Meiryo', align=PP_ALIGN.LEFT):
    """行ごとにスタイルを変えられるテキストボックス"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, col) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = col
        p.font.name = font_name
        p.alignment = align
        p.space_after = spacing
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="\u2022"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Meiryo'
        p.space_after = spacing


def add_arrow_right(slide, left, top, width, height, color):
    """右向き矢印"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_diagram_box(slide, left, top, width, height, text, fill_color, text_color=WHITE,
                    font_size=12, border_color=None, corner_radius=True):
    """図用のラベル付きボックス"""
    if border_color:
        shape = add_bordered_shape(slide, left, top, width, height, fill_color, border_color, corner_radius=corner_radius)
    else:
        shape = add_shape(slide, left, top, width, height, fill_color, corner_radius=corner_radius)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Meiryo'

    # 垂直方向の中央揃え
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return shape


# ==============================
# スライド作成
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # スライド 1: 表紙
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLAW_DARK)

    # ロブスターアイコン
    add_text(slide, Inches(5.5), Inches(0.8), Inches(2.33), Inches(1.2),
             "\U0001F99E", font_size=72, align=PP_ALIGN.CENTER, color=WHITE)

    add_text(slide, Inches(1), Inches(2.0), Inches(11.33), Inches(1.0),
             "OpenClaw & ClawHub",
             font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(3.2), Inches(11.33), Inches(0.6),
             "AI時代の「自分だけのアシスタント」を手に入れよう",
             font_size=24, color=CLAW_RED, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5.5), Inches(4.2), Inches(2.33), Inches(0.03), CLAW_RED)

    add_text(slide, Inches(1), Inches(4.6), Inches(11.33), Inches(0.8),
             "機能紹介・システム構成・使い方・革新性\n初心者向けガイド",
             font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(6.3), Inches(11.33), Inches(0.4),
             "2026年3月",
             font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ==========================================
    # スライド 2: OpenClawとは？
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "\U0001F99E  OpenClaw（オープンクロー）とは？",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    # 概要カード
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.33), Inches(1.6), SOFT_RED_BG, corner_radius=True)
    add_text(slide, Inches(1.0), Inches(1.5), Inches(11.33), Inches(1.2),
             "OpenClawは、あなた専用のAIアシスタントを自分のパソコンやサーバーで動かせる\n"
             "オープンソースのプロジェクトです。LINEやWhatsApp、Telegramなどの\n"
             "メッセージアプリを通じて、AIに様々な作業を頼めます。",
             font_size=16, color=DARK)

    # 基本情報カード（3列）
    infos = [
        ("開発者", "Peter Steinberger\n(PSPDFKit創業者)", ACCENT_BLUE),
        ("ライセンス", "MIT（完全無料）\n商用利用も自由", ACCENT_GREEN),
        ("人気度", "GitHub 220,000+\u2605\n600人以上が開発に参加", CLAW_RED),
    ]
    for i, (title, desc, accent) in enumerate(infos):
        x = Inches(0.5) + i * Inches(4.2)
        y = Inches(3.2)
        card = add_shape(slide, x, y, Inches(3.8), Inches(1.6), LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, Inches(3.8), Inches(0.06), accent)
        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.3),
                 title, font_size=14, bold=True, color=accent)
        add_text(slide, x + Inches(0.3), y + Inches(0.6), Inches(3.2), Inches(0.8),
                 desc, font_size=13, color=DARK)

    # 歴史タイムライン
    add_text(slide, Inches(0.8), Inches(5.1), Inches(11), Inches(0.4),
             "名称の変遷", font_size=16, bold=True, color=CLAW_DARK)

    timeline_items = [
        ("2025年11月", "Clawdbot\nとして公開"),
        ("2026年1月", "Moltbot\nに改名"),
        ("2026年1月末", "OpenClaw\nに再改名"),
        ("2026年2月〜", "22万\u2605突破\n爆発的成長"),
    ]
    for i, (date, label) in enumerate(timeline_items):
        x = Inches(0.8) + i * Inches(3.1)
        y = Inches(5.6)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y, Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = CLAW_RED if i == 3 else ACCENT_BLUE
        circle.line.fill.background()
        circle.shadow.inherit = False

        if i < 3:
            add_shape(slide, x + Inches(1.25), y + Inches(0.15), Inches(2.5), Inches(0.1), LIGHT_GRAY)

        add_text(slide, x, y + Inches(0.5), Inches(2.1), Inches(0.3),
                 date, font_size=11, bold=True, color=CLAW_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + Inches(0.85), Inches(2.1), Inches(0.5),
                 label, font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # ==========================================
    # スライド 3: OpenClawの主な機能
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "OpenClawでできること（主な機能）",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    features = [
        ("\U0001F4AC", "チャットで指示", "メッセージアプリ対応",
         "LINE・WhatsApp・Telegram・Discord・\niMessageなど、いつものアプリから\nAIに話しかけるだけで操作できます",
         ACCENT_BLUE, SOFT_BLUE_BG),
        ("\U0001F310", "ブラウザ操作", "Webサイトの自動操作",
         "フォームの入力、データの収集、\n予約の自動化など、ブラウザ上の\n作業をAIが代行します",
         ACCENT_GREEN, SOFT_GREEN_BG),
        ("\U0001F4C1", "ファイル管理", "読み書き・整理の自動化",
         "ドキュメントの作成・整理、\nCSVデータの処理、メールの下書き\nなどを自動でこなします",
         ACCENT_ORANGE, SOFT_ORANGE_BG),
        ("\U0001F3E0", "スマートホーム", "IoT機器の制御",
         "照明のON/OFF、エアコンの調整、\nセキュリティカメラの確認など\n家電をAIで声掛け操作",
         ACCENT_PURPLE, SOFT_PURPLE_BG),
        ("\U0001F9E0", "記憶（メモリ）", "会話を覚えてくれる",
         "過去に話した内容をMarkdownファイルに\n自動保存。あなたの好みや履歴を\n踏まえた応答をしてくれます",
         CLAW_RED, SOFT_RED_BG),
        ("\u23F0", "自動スケジュール", "自律的にタスク実行",
         "「毎朝ニュースをまとめて」など\n時間指定のタスクを自動実行。\n人間が指示しなくても動きます",
         ACCENT_TEAL, RGBColor(0xF0, 0xFD, 0xFA)),
    ]

    for i, (icon, title, subtitle, desc, accent, bg) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.6), bg, corner_radius=True)
        add_shape(slide, x + Inches(0.1), y + Inches(0.06), Inches(3.6), Inches(0.05), accent)

        add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=28, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.8), y + Inches(0.2), Inches(2.8), Inches(0.3),
                 title, font_size=16, bold=True, color=accent)
        add_text(slide, x + Inches(0.8), y + Inches(0.55), Inches(2.8), Inches(0.3),
                 subtitle, font_size=11, color=GRAY)
        add_text(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(1.4),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # スライド 4: システム構成図
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "システム構成図（アーキテクチャ）",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "OpenClawは「自分のPC/サーバーで動く」のが最大のポイント。データは外部に出ません。",
             font_size=14, color=MID_GRAY)

    # --- 図: 左側 = ユーザー入力チャネル ---
    add_text(slide, Inches(0.3), Inches(1.7), Inches(2.2), Inches(0.4),
             "あなた（ユーザー）", font_size=13, bold=True, color=CLAW_DARK, align=PP_ALIGN.CENTER)

    channels = [
        ("\U0001F4F1 WhatsApp", ACCENT_GREEN),
        ("\U0001F4AC Telegram", ACCENT_BLUE),
        ("\U0001F3AE Discord", ACCENT_PURPLE),
        ("\U0001F4E7 iMessage", GRAY),
        ("\U0001F310 Web UI", ACCENT_ORANGE),
    ]
    for i, (ch, col) in enumerate(channels):
        y = Inches(2.2) + i * Inches(0.7)
        add_diagram_box(slide, Inches(0.3), y, Inches(2.2), Inches(0.55), ch, col, WHITE, 11)

    # 矢印: チャネル → Gateway
    for i in range(5):
        y = Inches(2.35) + i * Inches(0.7)
        add_arrow_right(slide, Inches(2.6), y, Inches(0.7), Inches(0.25), CLAW_RED)

    # --- 中央: Gateway + Agent Runtime ---
    # Gateway ボックス
    gw_x = Inches(3.5)
    gw_y = Inches(1.7)
    add_bordered_shape(slide, gw_x, gw_y, Inches(3.0), Inches(5.0), WHITE, CLAW_RED, Pt(3), corner_radius=True)
    add_text(slide, gw_x + Inches(0.1), gw_y + Inches(0.1), Inches(2.8), Inches(0.35),
             "Gateway（中継サーバー）", font_size=13, bold=True, color=CLAW_RED, align=PP_ALIGN.CENTER)

    # セッション管理
    add_diagram_box(slide, gw_x + Inches(0.2), Inches(2.3), Inches(2.6), Inches(0.5),
                    "セッション管理", NAVY, WHITE, 11)

    # Agent Runtime
    add_bordered_shape(slide, gw_x + Inches(0.2), Inches(3.0), Inches(2.6), Inches(3.3),
                       SOFT_BLUE_BG, ACCENT_BLUE, Pt(2), corner_radius=True)
    add_text(slide, gw_x + Inches(0.3), Inches(3.1), Inches(2.4), Inches(0.3),
             "Agent Runtime", font_size=12, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    runtime_parts = [
        "コンテキスト構築",
        "LLM呼び出し",
        "ツール実行",
        "結果を保存",
    ]
    for i, part in enumerate(runtime_parts):
        y = Inches(3.5) + i * Inches(0.6)
        add_diagram_box(slide, gw_x + Inches(0.4), y, Inches(2.2), Inches(0.45),
                        part, ACCENT_BLUE, WHITE, 10)
        if i < 3:
            add_text(slide, gw_x + Inches(1.2), y + Inches(0.4), Inches(0.6), Inches(0.25),
                     "\u2193", font_size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # 矢印: Gateway → 外部LLM
    add_arrow_right(slide, Inches(6.7), Inches(3.7), Inches(0.7), Inches(0.25), ACCENT_ORANGE)
    add_text(slide, Inches(6.65), Inches(3.35), Inches(0.8), Inches(0.3),
             "API", font_size=10, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

    # --- 右側上: 外部LLM ---
    add_text(slide, Inches(7.6), Inches(1.7), Inches(2.5), Inches(0.35),
             "外部AIモデル（LLM）", font_size=13, bold=True, color=CLAW_DARK, align=PP_ALIGN.CENTER)

    llms = [
        ("Claude", ACCENT_ORANGE),
        ("GPT (OpenAI)", ACCENT_GREEN),
        ("DeepSeek", ACCENT_BLUE),
    ]
    for i, (name, col) in enumerate(llms):
        y = Inches(2.2) + i * Inches(0.6)
        add_diagram_box(slide, Inches(7.6), y, Inches(2.5), Inches(0.45), name, col, WHITE, 11)

    # --- 右側下: ツール/スキル ---
    add_text(slide, Inches(7.6), Inches(4.2), Inches(2.5), Inches(0.35),
             "ツール / スキル", font_size=13, bold=True, color=CLAW_DARK, align=PP_ALIGN.CENTER)

    # 矢印: Gateway → ツール
    add_arrow_right(slide, Inches(6.7), Inches(5.0), Inches(0.7), Inches(0.25), ACCENT_GREEN)

    tools = [
        ("\U0001F4C2 ファイル操作", ACCENT_GREEN),
        ("\U0001F310 ブラウザ", ACCENT_BLUE),
        ("\U0001F527 シェル/CLI", NAVY),
        ("\U0001F9E9 Skills (ClawHub)", CLAW_RED),
    ]
    for i, (name, col) in enumerate(tools):
        y = Inches(4.6) + i * Inches(0.55)
        add_diagram_box(slide, Inches(7.6), y, Inches(2.5), Inches(0.42), name, col, WHITE, 10)

    # --- 右端: メモリ ---
    add_text(slide, Inches(10.5), Inches(1.7), Inches(2.5), Inches(0.35),
             "メモリ（記憶）", font_size=13, bold=True, color=CLAW_DARK, align=PP_ALIGN.CENTER)

    add_bordered_shape(slide, Inches(10.5), Inches(2.2), Inches(2.5), Inches(2.2),
                       SOFT_GREEN_BG, ACCENT_GREEN, Pt(2), corner_radius=True)
    add_text(slide, Inches(10.7), Inches(2.35), Inches(2.1), Inches(0.3),
             "\U0001F4DD Markdownファイル", font_size=11, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.7), Inches(2.75), Inches(2.1), Inches(1.4),
             "会話履歴\n好みの設定\nタスク記録\n\n全てローカルに保存\n\u2192 プライバシー安全",
             font_size=10, color=DARK, align=PP_ALIGN.CENTER)

    # 矢印: ツール/LLM → メモリ
    add_text(slide, Inches(10.1), Inches(3.0), Inches(0.5), Inches(0.3),
             "\u2194", font_size=16, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    # 下部の説明
    add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.8), WHITE, corner_radius=True)
    add_text(slide, Inches(0.8), Inches(6.55), Inches(11.73), Inches(0.7),
             "\U0001F512 ポイント: Gateway（中心のエンジン）は自分のPCやサーバーで動作。"
             "メモリ（記憶）もローカルのMarkdownファイルに保存されるため、個人データが外部に漏れません。"
             "外部に送られるのはAIモデルへの問いかけだけです。",
             font_size=12, color=DARK)

    # ==========================================
    # スライド 5: ClawHubとは？
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "\U0001F9E9  ClawHub（クローハブ）とは？",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    # たとえ話
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.33), Inches(1.4), SOFT_BLUE_BG, corner_radius=True)
    add_text(slide, Inches(1.0), Inches(1.4), Inches(11.33), Inches(0.3),
             "\U0001F4A1 かんたんに言うと…", font_size=16, bold=True, color=ACCENT_BLUE)
    add_text(slide, Inches(1.0), Inches(1.8), Inches(11.33), Inches(0.7),
             "ClawHubは、OpenClawの「スキル（追加機能）」を探して導入できるマーケットプレイスです。\n"
             "スマホで言えば「App Store」や「Google Play」のような存在。\n"
             "AIアシスタントに新しい能力を追加したいとき、ここから1コマンドでインストールできます。",
             font_size=15, color=DARK)

    # ClawHub概要図
    # 左: 開発者 → 中央: ClawHub → 右: ユーザー
    add_text(slide, Inches(0.5), Inches(3.0), Inches(12), Inches(0.4),
             "ClawHubの仕組み", font_size=18, bold=True, color=CLAW_DARK)

    # 開発者ボックス
    add_diagram_box(slide, Inches(0.5), Inches(3.7), Inches(3.0), Inches(2.5),
                    "", SOFT_PURPLE_BG, DARK, 11, ACCENT_PURPLE)
    add_text(slide, Inches(0.7), Inches(3.8), Inches(2.6), Inches(0.3),
             "\U0001F468\u200D\U0001F4BB Skill開発者", font_size=14, bold=True, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)
    add_bullet_list(slide, Inches(0.7), Inches(4.2), Inches(2.6), Inches(1.5),
                    ["SKILL.mdを作成", "機能をコードで実装", "clawhub publishで公開"],
                    font_size=11, icon="\u2022", color=DARK)

    # 矢印
    add_arrow_right(slide, Inches(3.7), Inches(4.7), Inches(0.8), Inches(0.3), CLAW_RED)
    add_text(slide, Inches(3.7), Inches(4.4), Inches(0.8), Inches(0.3),
             "公開", font_size=10, bold=True, color=CLAW_RED, align=PP_ALIGN.CENTER)

    # ClawHub中央
    add_diagram_box(slide, Inches(4.7), Inches(3.5), Inches(3.8), Inches(3.0),
                    "", SOFT_RED_BG, DARK, 11, CLAW_RED)
    add_text(slide, Inches(4.9), Inches(3.6), Inches(3.4), Inches(0.35),
             "\U0001F9E9 ClawHub (clawhub.ai)", font_size=14, bold=True, color=CLAW_RED, align=PP_ALIGN.CENTER)

    hub_features = [
        "10,700+ スキル登録",
        "ベクトル検索で高速発見",
        "バージョン管理 (semver)",
        "\u2605 スター・コメント機能",
        "1コマンドでインストール",
    ]
    add_bullet_list(slide, Inches(5.0), Inches(4.1), Inches(3.3), Inches(2.0),
                    hub_features, font_size=11, icon="\u2713", color=DARK)

    # 矢印
    add_arrow_right(slide, Inches(8.7), Inches(4.7), Inches(0.8), Inches(0.3), ACCENT_GREEN)
    add_text(slide, Inches(8.7), Inches(4.4), Inches(0.8), Inches(0.3),
             "導入", font_size=10, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    # ユーザーボックス
    add_diagram_box(slide, Inches(9.7), Inches(3.7), Inches(3.0), Inches(2.5),
                    "", SOFT_GREEN_BG, DARK, 11, ACCENT_GREEN)
    add_text(slide, Inches(9.9), Inches(3.8), Inches(2.6), Inches(0.3),
             "\U0001F464 あなた（ユーザー）", font_size=14, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_bullet_list(slide, Inches(9.9), Inches(4.2), Inches(2.6), Inches(1.5),
                    ["clawhub install xxx", "すぐに新機能が使える", "不要ならuninstall"],
                    font_size=11, icon="\u2022", color=DARK)

    # 注意喚起
    add_shape(slide, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.7), RGBColor(0xFF, 0xF3, 0xCD), corner_radius=True)
    add_text(slide, Inches(0.8), Inches(6.7), Inches(11.73), Inches(0.5),
             "\u26A0\uFE0F 注意: 誰でもスキルを公開できるため、悪意のあるスキルが混在するリスクがあります。"
             "インストール前にSKILL.mdの内容を確認し、信頼できる開発者のものを選びましょう。",
             font_size=12, color=RGBColor(0x85, 0x6D, 0x0D))

    # ==========================================
    # スライド 6: Skillとは何か？
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Skill（スキル）の正体",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    # Skillの構造図
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5), SOFT_BLUE_BG, corner_radius=True)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(5.1), Inches(0.4),
             "\U0001F4C4 Skillの正体 = SKILL.md（Markdownファイル）", font_size=15, bold=True, color=ACCENT_BLUE)

    add_text(slide, Inches(0.7), Inches(1.9), Inches(5.1), Inches(0.5),
             "スキルはプログラムではなく、「AIへの指示書」です。\n"
             "YAMLヘッダー + Markdown本文で構成されます。",
             font_size=13, color=DARK)

    # SKILL.md のイメージ
    add_shape(slide, Inches(0.8), Inches(2.6), Inches(5.0), Inches(4.0), WHITE, corner_radius=True)
    code_lines = [
        ("---", 12, True, ACCENT_PURPLE),
        ("name: weather-report", 11, False, DARK),
        ("version: 1.0.0", 11, False, DARK),
        ("description: 天気予報を取得", 11, False, DARK),
        ("---", 12, True, ACCENT_PURPLE),
        ("", 8, False, DARK),
        ("# Weather Report Skill", 13, True, CLAW_DARK),
        ("", 8, False, DARK),
        ("## 概要", 12, True, MID_GRAY),
        ("指定した都市の天気予報を取得して", 11, False, DARK),
        ("わかりやすく要約します。", 11, False, DARK),
        ("", 8, False, DARK),
        ("## 使い方", 12, True, MID_GRAY),
        ("「東京の天気を教えて」と聞くだけ", 11, False, DARK),
    ]
    add_multiline_text(slide, Inches(1.0), Inches(2.7), Inches(4.6), Inches(3.8),
                       code_lines, font_size=11, font_name='Courier New')

    # 右側: スキルの種類
    add_text(slide, Inches(6.5), Inches(1.3), Inches(6), Inches(0.4),
             "スキルの例（人気カテゴリ）", font_size=18, bold=True, color=CLAW_DARK)

    skill_examples = [
        ("\U0001F4E7", "メール管理", "受信メールの要約、返信の下書き", ACCENT_BLUE),
        ("\U0001F4C5", "カレンダー連携", "予定の確認・登録・リマインド", ACCENT_GREEN),
        ("\U0001F4CA", "データ分析", "CSVファイルの集計・グラフ作成", ACCENT_ORANGE),
        ("\U0001F6D2", "ECサイト連携", "商品情報の取得・価格比較", ACCENT_PURPLE),
        ("\U0001F3B5", "音楽操作", "Spotify等の再生・プレイリスト管理", CLAW_RED),
        ("\U0001F4F0", "ニュース収集", "指定トピックの最新ニュース要約", ACCENT_TEAL),
        ("\U0001F4DD", "メモ・ノート", "議事録作成、読書メモの整理", NAVY),
        ("\U0001F3E0", "スマートホーム", "照明・エアコン・カメラの制御", RGBColor(0xE6, 0x7E, 0x22)),
    ]

    for i, (icon, title, desc, accent) in enumerate(skill_examples):
        col = i % 2
        row = i // 2
        x = Inches(6.5) + col * Inches(3.2)
        y = Inches(1.9) + row * Inches(1.3)

        card = add_shape(slide, x, y, Inches(2.9), Inches(1.1), LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, Inches(0.06), Inches(1.1), accent)
        add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.4), Inches(0.4),
                 icon, font_size=20, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.5), y + Inches(0.1), Inches(2.2), Inches(0.3),
                 title, font_size=13, bold=True, color=accent)
        add_text(slide, x + Inches(0.5), y + Inches(0.45), Inches(2.2), Inches(0.5),
                 desc, font_size=10, color=MID_GRAY)

    # ==========================================
    # スライド 7: 使い方（セットアップ）
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "使い方ガイド（セットアップ編）",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    steps = [
        ("STEP 1", "インストール",
         "npm install -g openclaw",
         "Node.jsがあればターミナルで\n1コマンドでインストール完了",
         ACCENT_BLUE),
        ("STEP 2", "初期設定",
         "openclaw init",
         "対話形式で設定ファイルを作成。\n使うAIモデルのAPIキーを入力",
         ACCENT_GREEN),
        ("STEP 3", "チャネル接続",
         "設定ファイルにアプリのトークンを記入",
         "WhatsApp / Telegram / Discord\nなど好きなアプリと接続",
         ACCENT_ORANGE),
        ("STEP 4", "起動",
         "openclaw start",
         "Gatewayが起動し、チャットで\nAIアシスタントとやり取り開始！",
         CLAW_RED),
    ]

    for i, (step, title, cmd, desc, accent) in enumerate(steps):
        x = Inches(0.4) + i * Inches(3.2)
        y = Inches(1.5)

        # ステップ番号
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.95), y, Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        circle.shadow.inherit = False
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Meiryo'
        p.alignment = PP_ALIGN.CENTER

        if i < 3:
            add_text(slide, x + Inches(2.1), y + Inches(0.25), Inches(1.0), Inches(0.5),
                     "\u27A1", font_size=24, color=accent, align=PP_ALIGN.CENTER)

        # タイトル
        add_text(slide, x + Inches(0.1), y + Inches(1.2), Inches(2.7), Inches(0.3),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # コマンド
        add_shape(slide, x + Inches(0.1), y + Inches(1.6), Inches(2.7), Inches(0.5),
                  CLAW_DARK, corner_radius=True)
        add_text(slide, x + Inches(0.2), y + Inches(1.65), Inches(2.5), Inches(0.4),
                 cmd, font_size=10, color=ACCENT_GREEN, align=PP_ALIGN.CENTER,
                 font_name='Courier New')

        # 説明
        add_text(slide, x + Inches(0.1), y + Inches(2.3), Inches(2.7), Inches(1.2),
                 desc, font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Skill導入セクション
    add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.33), Inches(2.2), WHITE, corner_radius=True)
    add_text(slide, Inches(0.8), Inches(5.1), Inches(11.73), Inches(0.4),
             "\U0001F9E9 Skillの導入方法（ClawHub）", font_size=18, bold=True, color=CLAW_DARK)

    skill_steps = [
        ("探す", "clawhub.ai でスキルを検索\nまたは clawhub search <キーワード>", ACCENT_BLUE),
        ("入れる", "clawhub install <スキル名>\n例: clawhub install weather-report", ACCENT_GREEN),
        ("使う", "チャットで「天気を教えて」と話すだけ\nAIが自動的にスキルを選んで実行", CLAW_RED),
        ("管理する", "clawhub list で一覧確認\nclawhub uninstall で削除", GRAY),
    ]
    for i, (label, desc, accent) in enumerate(skill_steps):
        x = Inches(0.8) + i * Inches(3.1)
        y = Inches(5.6)

        add_diagram_box(slide, x, y, Inches(1.0), Inches(0.5), label, accent, WHITE, 13)
        add_text(slide, x, y + Inches(0.6), Inches(2.8), Inches(0.9),
                 desc, font_size=10, color=DARK)

    # ==========================================
    # スライド 8: 何が画期的か？
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLAW_DARK)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "\u2728 OpenClaw & ClawHub の何が画期的なのか？",
             font_size=30, bold=True, color=WHITE)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    innovations = [
        ("\U0001F512", "完全ローカル＆プライバシー重視",
         "AIアシスタントなのに、データは全て自分のPCに保存。\n"
         "クラウドに個人情報を預ける不安がありません。\n"
         "テレメトリ（使用状況の送信）も一切なし。",
         CLAW_RED),
        ("\U0001F310", "どのAIモデルでも使える",
         "Claude・GPT・DeepSeekなど、好きなAIを選べます。\n"
         "特定の会社に縛られない「モデル非依存」設計。\n"
         "将来、新しいAIが出ても簡単に切り替え可能。",
         ACCENT_BLUE),
        ("\U0001F4AC", "いつものアプリからAIを使える",
         "新しいアプリのインストール不要。\n"
         "WhatsApp・LINE・Discord等、普段使いの\n"
         "チャットアプリがそのままAIの入口になります。",
         ACCENT_GREEN),
        ("\U0001F504", "AIが自分から動いてくれる",
         "指示を待つだけのAIではありません。\n"
         "「毎朝ニュースをまとめて送信」のように\n"
         "スケジュール設定で自律的にタスクを実行。",
         ACCENT_ORANGE),
        ("\U0001F9E9", "スキルで無限に拡張できる",
         "ClawHubで10,700以上のスキルが利用可能。\n"
         "必要な機能を1コマンドで追加。\n"
         "自作スキルの公開もカンタン。",
         ACCENT_PURPLE),
        ("\U0001F465", "世界中が育てるオープンソース",
         "600人以上の開発者が参加するコミュニティ。\n"
         "バグ修正や新機能が毎週のように追加。\n"
         "商用利用も含め、完全無料（MIT License）。",
         ACCENT_TEAL),
    ]

    for i, (icon, title, desc, accent) in enumerate(innovations):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.3) + row * Inches(2.0)

        card = add_shape(slide, x, y, Inches(5.9), Inches(1.75), NAVY, corner_radius=True)
        add_shape(slide, x, y, Inches(0.08), Inches(1.75), accent)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=24, align=PP_ALIGN.CENTER, color=WHITE)
        add_text(slide, x + Inches(0.7), y + Inches(0.15), Inches(5.0), Inches(0.3),
                 title, font_size=15, bold=True, color=accent)
        add_text(slide, x + Inches(0.7), y + Inches(0.5), Inches(5.0), Inches(1.1),
                 desc, font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC))

    # ==========================================
    # スライド 9: 従来のAIサービスとの比較
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "従来のAIサービスとの比較",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    # テーブルヘッダー
    cols = [Inches(0.5), Inches(3.5), Inches(6.5), Inches(9.7)]
    col_w = [Inches(2.8), Inches(2.8), Inches(3.0), Inches(3.0)]
    headers = ["比較項目", "ChatGPT等\nクラウドAI", "他のAIエージェント\n(AutoGPT等)", "OpenClaw\n+ ClawHub"]
    header_colors = [NAVY, GRAY, ACCENT_BLUE, CLAW_RED]

    for i, (header, col_x, w, hc) in enumerate(zip(headers, cols, col_w, header_colors)):
        add_diagram_box(slide, col_x, Inches(1.2), w, Inches(0.8), header, hc, WHITE, 12)

    # テーブル行
    rows_data = [
        ("データの保管場所", "クラウド\n（サービス提供者側）", "主にクラウド", "ローカル（自分のPC）\n完全に自分の管理下"),
        ("対応チャットアプリ", "専用Webサイト\nのみ", "専用UIが多い", "WhatsApp/Telegram/\nDiscord/LINE等"),
        ("拡張性", "プラグイン\n（限定的）", "ツール追加可能\n（手動設定が多い）", "ClawHubで10,700+\nスキルを1コマンド導入"),
        ("自律的な動作", "指示待ち\n（受動的）", "一部対応", "スケジュール実行\n（能動的に動く）"),
        ("コスト", "月額サブスク\n$20/月〜", "無料+API費用", "完全無料（OSS）\n+ API費用のみ"),
        ("AIモデルの選択", "固定\n（GPTのみ等）", "一部選択可", "Claude/GPT/DeepSeek\n等、自由に選べる"),
    ]

    for r, (label, v1, v2, v3) in enumerate(rows_data):
        y = Inches(2.1) + r * Inches(0.85)
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE

        add_shape(slide, cols[0], y, col_w[0], Inches(0.75), bg)
        add_text(slide, cols[0] + Inches(0.15), y + Inches(0.1), col_w[0] - Inches(0.3), Inches(0.55),
                 label, font_size=12, bold=True, color=DARK)

        add_shape(slide, cols[1], y, col_w[1], Inches(0.75), bg)
        add_text(slide, cols[1] + Inches(0.15), y + Inches(0.05), col_w[1] - Inches(0.3), Inches(0.65),
                 v1, font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

        add_shape(slide, cols[2], y, col_w[2], Inches(0.75), bg)
        add_text(slide, cols[2] + Inches(0.15), y + Inches(0.05), col_w[2] - Inches(0.3), Inches(0.65),
                 v2, font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

        # OpenClawの列を強調
        highlight_bg = SOFT_RED_BG if r % 2 == 0 else RGBColor(0xFF, 0xFA, 0xFA)
        add_shape(slide, cols[3], y, col_w[3], Inches(0.75), highlight_bg)
        add_text(slide, cols[3] + Inches(0.15), y + Inches(0.05), col_w[3] - Inches(0.3), Inches(0.65),
                 v3, font_size=10, bold=True, color=CLAW_RED, align=PP_ALIGN.CENTER)

    # ==========================================
    # スライド 10: セキュリティの注意点
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "\U0001F6E1\uFE0F  セキュリティに関する注意点",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    # 警告カード
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.33), Inches(1.4), RGBColor(0xFF, 0xF3, 0xCD), corner_radius=True)
    add_text(slide, Inches(1.0), Inches(1.4), Inches(11.33), Inches(0.3),
             "\u26A0\uFE0F 2026年2月に発覚した問題", font_size=16, bold=True, color=RGBColor(0x85, 0x6D, 0x0D))
    add_text(slide, Inches(1.0), Inches(1.8), Inches(11.33), Inches(0.7),
             "セキュリティ研究者の調査で、ClawHubに登録されていた2,857スキルのうち約341個（約12%）に\n"
             "悪意のあるコードが含まれていたことが報告されました。主にソーシャルエンジニアリングと\n"
             "情報窃取（インフォスティーラー）を目的としたものでした。",
             font_size=14, color=DARK)

    # 安全に使うためのポイント
    add_text(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.4),
             "安全に使うための5つのポイント", font_size=20, bold=True, color=CLAW_DARK)

    safety_tips = [
        ("\U0001F50D", "SKILL.mdを必ず確認",
         "インストール前にスキルの内容を読み、\n不審なコードや権限要求がないか確認",
         ACCENT_BLUE),
        ("\u2B50", "人気のスキルを選ぶ",
         "\u2605スターが多く、コメントで評判の良い\nスキルを優先的に選びましょう",
         ACCENT_GREEN),
        ("\U0001F464", "開発者の信頼性を確認",
         "GitHubプロフィールや過去の活動歴を\nチェック。新規アカウントは慎重に",
         ACCENT_ORANGE),
        ("\U0001F512", "権限を最小限に",
         "スキルに与える権限（ファイルアクセス等）は\n必要最小限に設定しましょう",
         ACCENT_PURPLE),
        ("\U0001F504", "定期的に更新",
         "OpenClaw本体とスキルを最新版に保つことで\nセキュリティパッチを適用",
         CLAW_RED),
    ]

    for i, (icon, title, desc, accent) in enumerate(safety_tips):
        x = Inches(0.3) + i * Inches(2.55)
        y = Inches(3.6)

        card = add_shape(slide, x, y, Inches(2.35), Inches(2.8), LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x + Inches(0.05), y + Inches(0.04), Inches(2.25), Inches(0.05), accent)

        add_text(slide, x + Inches(0.1), y + Inches(0.2), Inches(2.15), Inches(0.4),
                 icon, font_size=28, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), y + Inches(0.65), Inches(2.15), Inches(0.4),
                 title, font_size=13, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.15), Inches(1.5),
                 desc, font_size=10, color=DARK, align=PP_ALIGN.CENTER)

    # 下部メッセージ
    add_shape(slide, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.6), SOFT_GREEN_BG, corner_radius=True)
    add_text(slide, Inches(0.8), Inches(6.7), Inches(11.73), Inches(0.4),
             "\u2705 OpenClaw自体のアーキテクチャはローカルファースト設計のため、"
             "正しく運用すればデータの安全性は非常に高いです。",
             font_size=13, color=ACCENT_GREEN)

    # ==========================================
    # スライド 11: まとめ
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLAW_DARK)

    add_text(slide, Inches(1), Inches(0.5), Inches(11.33), Inches(0.6),
             "まとめ", font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5.5), Inches(1.2), Inches(2.33), Inches(0.03), CLAW_RED)

    summary_items = [
        ("OpenClawは", "自分のPCで動く、オープンソースのAIアシスタントプラットフォーム", CLAW_RED),
        ("ClawHubは", "AIの能力を拡張する10,700以上のスキルが集まるマーケットプレイス", ACCENT_BLUE),
        ("画期的な点は", "プライバシー重視・モデル非依存・自律動作・無限の拡張性", ACCENT_GREEN),
        ("注意点は", "スキルのセキュリティ確認を怠らないこと", ACCENT_ORANGE),
    ]

    for i, (prefix, desc, accent) in enumerate(summary_items):
        y = Inches(1.6) + i * Inches(1.15)
        add_shape(slide, Inches(1.5), y, Inches(10.33), Inches(0.95), NAVY, corner_radius=True)
        add_shape(slide, Inches(1.5), y, Inches(0.08), Inches(0.95), accent)

        add_text(slide, Inches(1.8), y + Inches(0.15), Inches(2.5), Inches(0.35),
                 prefix, font_size=16, bold=True, color=accent)
        add_text(slide, Inches(1.8), y + Inches(0.5), Inches(9.7), Inches(0.35),
                 desc, font_size=15, color=WHITE)

    # こんな人におすすめ
    add_text(slide, Inches(1), Inches(6.0), Inches(11.33), Inches(0.4),
             "\U0001F3AF こんな人におすすめ", font_size=18, bold=True, color=CLAW_RED, align=PP_ALIGN.CENTER)

    targets = [
        "AIを試してみたいけどプライバシーが心配な人",
        "日常の作業を自動化したいビジネスパーソン",
        "自分だけのAIアシスタントを作りたい開発者",
    ]
    for i, t in enumerate(targets):
        x = Inches(0.8) + i * Inches(4.0)
        add_text(slide, x, Inches(6.5), Inches(3.8), Inches(0.5),
                 f"\u2713  {t}", font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)

    # ==========================================
    # スライド 12: 参考リンク
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "\U0001F517  参考リンク・情報源",
             font_size=30, bold=True, color=CLAW_DARK)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04), CLAW_RED)

    links = [
        ("OpenClaw 公式サイト", "https://openclaw.ai/", ACCENT_BLUE),
        ("OpenClaw GitHub リポジトリ", "https://github.com/openclaw/openclaw", CLAW_DARK),
        ("OpenClaw 公式ドキュメント", "https://docs.openclaw.ai/", ACCENT_GREEN),
        ("ClawHub（スキルマーケットプレイス）", "https://clawhub.ai/", CLAW_RED),
        ("ClawHub GitHub リポジトリ", "https://github.com/openclaw/clawhub", NAVY),
        ("OpenClaw Wikipedia", "https://en.wikipedia.org/wiki/OpenClaw", ACCENT_ORANGE),
    ]

    for i, (title, url, accent) in enumerate(links):
        y = Inches(1.3) + i * Inches(0.85)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(0.7), LIGHT_GRAY, corner_radius=True)
        add_shape(slide, Inches(0.5), y, Inches(0.06), Inches(0.7), accent)
        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(5.0), Inches(0.3),
                 title, font_size=14, bold=True, color=accent)
        add_text(slide, Inches(0.8), y + Inches(0.38), Inches(11.73), Inches(0.25),
                 url, font_size=11, color=GRAY, font_name='Courier New')

    # フッター
    add_text(slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.4),
             "本資料は2026年3月時点の情報に基づいています。最新情報は公式サイトをご確認ください。",
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# 出力
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    # PPT出力
    pptx_path = os.path.join(output_dir, 'OpenClaw_ClawHub_解説資料.pptx')
    prs.save(pptx_path)
    print(f"\u2705 PPTX saved: {pptx_path}")

    # PDF出力（macOS Keynoteを利用）
    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        import time

        # Keynote で PPTX → PDF 変換
        applescript = f'''
        tell application "Keynote"
            activate
            delay 2
            set theDoc to open POSIX file "{pptx_path}"
            delay 3
            export theDoc to POSIX file "{pdf_path}" as PDF
            delay 1
            close theDoc saving no
            delay 1
            quit
        end tell
        '''
        result = subprocess.run(['osascript', '-e', applescript],
                                capture_output=True, text=True, timeout=60)
        if result.returncode == 0:
            print(f"\u2705 PDF saved: {pdf_path}")
        else:
            raise RuntimeError(f"Keynote error: {result.stderr}")
    except Exception as e:
        print(f"\u26A0\uFE0F  PDF conversion skipped ({e})")
        print(f"   PPTXファイルからPDFを作成するには:")
        print(f"   - PowerPoint / Keynote で開いて「PDFとして書き出し」")
