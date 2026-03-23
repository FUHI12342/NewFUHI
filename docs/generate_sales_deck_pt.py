#!/usr/bin/env python3
"""
Sales Deck Generator — Portuguese (Português Brasileiro) Version
Outputs PPT + PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ==============================
# Color Palette
# ==============================
BROWN = RGBColor(0x8C, 0x87, 0x6C)
BEIGE = RGBColor(0xF1, 0xF0, 0xEC)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="\u2713"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing


def add_feature_card(slide, left, top, width, height, title, items, accent_color=BROWN):
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=True)
    line = add_shape(slide, left + Inches(0.15), top + Inches(0.08),
                     width - Inches(0.3), Inches(0.04), accent_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.2),
             width - Inches(0.4), Inches(0.4),
             title, font_size=15, bold=True, color=accent_color)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6),
                    width - Inches(0.4), height - Inches(0.8),
                    items, font_size=11, icon="\u25CF")


def add_picture_fit(slide, img_path, left, top, max_width, max_height):
    with Image.open(img_path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = max_width / max_height
    if ratio > box_ratio:
        w = max_width
        h = max_width / ratio
    else:
        h = max_height
        w = max_height * ratio
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    return slide.shapes.add_picture(img_path, int(x), int(y), width=int(w), height=int(h))


def add_number_highlight(slide, left, top, number, label, color=BROWN):
    add_text(slide, left, top, Inches(2), Inches(0.6),
             number, font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2), Inches(0.3),
             label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ==============================
# Slide Creation
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # Slide 1: Cover
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Plataforma DX Tudo-em-Um para Lojas",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "Reservas \u2022 Pedidos \u2022 Pagamentos \u2022 Gest\u00e3o de Equipe \u2022 Folha de Pagamento\nTudo em um \u00fanico sistema.",
             font_size=22, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5), Inches(4.0), Inches(3.33), Inches(0.03), WHITE)

    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "TimeBaiBai  Apresenta\u00e7\u00e3o do Servi\u00e7o",
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "2026  /  Nome da Empresa",
             font_size=14, color=BEIGE, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 2: Pain Points
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Algum desses problemas \u00e9 familiar?", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    problems = [
        ("Reservas Dispersas",
         "Telefone, LINE, e-mail... diferentes canais\ncausam preocupa\u00e7\u00e3o com reservas duplicadas.",
         "\U0001F4DE"),
        ("Aus\u00eancias e Cancelamentos de \u00daltima Hora",
         "Clientes reservam mas n\u00e3o aparecem.\nPerda de receita e prepara\u00e7\u00e3o desperdi\u00e7ada.",
         "\U0001F630"),
        ("Recebimento de Pedidos Sobrecarregado",
         "N\u00e3o consegue dar conta nos hor\u00e1rios de pico.\nErros nos pedidos acontecem com frequ\u00eancia.",
         "\U0001F4DD"),
        ("Pontos Cegos no Estoque",
         "Produtos esgotam sem aviso.\nSem saber quando fazer novos pedidos.",
         "\U0001F4E6"),
    ]

    for i, (title, desc, icon) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.6) + row * Inches(2.6)

        card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), LIGHT_GRAY, corner_radius=True)

        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8),
                 icon, font_size=36, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK)

        add_text(slide, x + Inches(1.2), y + Inches(0.75), Inches(4), Inches(1.2),
                 desc, font_size=14, color=GRAY)

    # ==========================================
    # Slide 3: Solution — What is TimeBaiBai
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "TimeBaiBai resolve tudo", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Plataforma completa para tudo que sua loja precisa. Comece apenas com um smartphone.",
             font_size=16, color=DARK)

    highlights = [
        ("6+", "M\u00f3dulos Principais"),
        ("7", "Idiomas Suportados"),
        ("24h", "Reserva Autom\u00e1tica"),
        ("R$0", "Economia em Papel"),
    ]
    for i, (num, label) in enumerate(highlights):
        x = Inches(1.2) + i * Inches(3)
        add_shape(slide, x - Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.4), WHITE, corner_radius=True)
        add_number_highlight(slide, x, Inches(2.5), num, label)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
             "Reservas, pedidos na mesa, estoque, gest\u00e3o de turnos, folha de pagamento, gest\u00e3o do site...\n"
             "Gerencie tudo em um \u00fanico painel \u2014 sem precisar usar v\u00e1rias ferramentas.\n\n"
             "A configura\u00e7\u00e3o \u00e9 simples. Comece a usar no mesmo dia da cria\u00e7\u00e3o da conta.\n"
             "Sem instala\u00e7\u00e3o complicada ou hardware especial. Acesse pelo celular ou computador.",
             font_size=15, color=DARK)

    # ==========================================
    # Slide 4: Features (1/4) Reservations & Orders
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funcionalidades Principais  \u2460  Reservas / Pedidos na Mesa",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Sistema de Reservas Online", [
            "Reserve pelo LINE \u2014 complete a reserva direto no chat",
            "Reservas por e-mail tamb\u00e9m dispon\u00edveis",
            "Confirma\u00e7\u00e3o com pr\u00e9-pagamento \u2014 confirma\u00e7\u00e3o autom\u00e1tica ao pagar (previne aus\u00eancias)",
            "Alterne entre pr\u00e9-pagamento / p\u00f3s-pagamento com um clique",
            "Calend\u00e1rios da equipe mostram disponibilidade em tempo real",
            "Dois fluxos de reserva: por data ou por membro da equipe",
            "Check-in por QR code (zero esfor\u00e7o na recep\u00e7\u00e3o)",
            "Notifica\u00e7\u00e3o autom\u00e1tica por LINE \u00e0 equipe na confirma\u00e7\u00e3o",
            "Suporte em 7 idiomas (japon\u00eas, ingl\u00eas, chin\u00eas, coreano, espanhol, portugu\u00eas)",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Sistema de Pedidos QR na Mesa", [
            "Coloque QR codes em cada mesa para come\u00e7ar",
            "Clientes escaneiam com o celular \u2192 veem o menu \u2192 fazem pedido",
            "Menus com fotos ajudam a visualizar os pratos",
            "Abas por categoria para busca r\u00e1pida",
            "Fun\u00e7\u00e3o de carrinho para pedidos adicionais f\u00e1ceis",
            "Acompanhamento do status do pedido em tempo real",
            "Suporte a pagamento: dinheiro, cart\u00e3o de cr\u00e9dito e dinheiro eletr\u00f4nico",
            "Gera\u00e7\u00e3o em massa de QR codes e download para impress\u00e3o",
        ], ACCENT_GREEN)

    # ==========================================
    # Slide 5: Features (2/4) Inventory & Staff
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funcionalidades Principais  \u2461  Estoque / Gest\u00e3o de Equipe e Turnos",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gest\u00e3o de Estoque", [
            "Acompanhamento de n\u00edveis de estoque em tempo real",
            "Dedu\u00e7\u00e3o autom\u00e1tica de estoque vinculada a pedidos",
            "Alertas autom\u00e1ticos quando o estoque est\u00e1 baixo",
            "Recebimento por QR code (sem scanner de c\u00f3digo de barras)",
            "Hist\u00f3rico completo de entradas/sa\u00eddas de estoque (pronto para auditoria)",
            "Organizado por categoria de produto",
            "Estoque unificado para loja EC e loja f\u00edsica",
            "Fun\u00e7\u00e3o de invent\u00e1rio para ajustes em massa",
        ], ACCENT_ORANGE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gest\u00e3o de Equipe e Turnos", [
            "Equipe envia prefer\u00eancias de turno pelo celular",
            "Tr\u00eas n\u00edveis: Dispon\u00edvel / Preferido / Indispon\u00edvel",
            "Gerentes finalizam turnos pelo painel de administra\u00e7\u00e3o",
            "Agendamento autom\u00e1tico para aloca\u00e7\u00e3o eficiente",
            "Turnos confirmados notificados automaticamente por LINE",
            "Perfis da equipe com fotos e descri\u00e7\u00f5es",
            "Tipos: Cast (atendimento) / Equipe da Loja",
            "Permiss\u00f5es por fun\u00e7\u00e3o: Propriet\u00e1rio / Gerente / Equipe / Desenvolvedor",
        ], ACCENT_BLUE)

    # ==========================================
    # Slide 6: Features (3/4) Payroll & Website
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funcionalidades Principais  \u2462  Folha de Pagamento / Gest\u00e3o do Site",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gest\u00e3o de Folha e Frequ\u00eancia", [
            "Gera\u00e7\u00e3o autom\u00e1tica de dados de frequ\u00eancia a partir de registros de turno",
            "Classifica\u00e7\u00e3o autom\u00e1tica de horas regulares, extras, noturnas e feriados",
            "C\u00e1lculo de folha para funcion\u00e1rios horistas e assalariados",
            "C\u00e1lculo autom\u00e1tico de encargos sociais",
            "Aplica\u00e7\u00e3o autom\u00e1tica de dedu\u00e7\u00f5es de impostos",
            "Gera\u00e7\u00e3o de holerites com um clique",
            "Exporta\u00e7\u00e3o CSV para transfer\u00eancias banc\u00e1rias",
            "Configura\u00e7\u00e3o f\u00e1cil de benef\u00edcios (transporte, moradia, fam\u00edlia)",
        ], ACCENT_GREEN)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gest\u00e3o do Site e Marketing", [
            "Edite o site da loja pelo painel de administra\u00e7\u00e3o",
            "Banners hero (sliders) para destacar sua loja",
            "Atualiza\u00e7\u00e3o f\u00e1cil de not\u00edcias e apari\u00e7\u00f5es na m\u00eddia",
            "Posicionamento e gest\u00e3o de banners publicit\u00e1rios",
            "Integra\u00e7\u00e3o com redes sociais (X / Instagram embeds)",
            "Blocos HTML personalizados para layouts flex\u00edveis",
            "P\u00e1ginas legais edit\u00e1veis pelo administrador",
            "Loja EC integrada para vendas online",
        ], ACCENT_ORANGE)

    # ==========================================
    # Slide 7: Features (4/4) Analytics & Security
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funcionalidades Principais  \u2463  An\u00e1lise de Dados / Seguran\u00e7a",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "An\u00e1lise de Dados e IA", [
            "Painel de vendas com gr\u00e1ficos de tend\u00eancia di\u00e1rio/semanal/mensal/anual",
            "Engenharia de Menu (classifica\u00e7\u00e3o Star/Dog) para melhoria de lucro",
            "An\u00e1lise ABC, RFM, coorte e de cesta",
            "NPS (satisfa\u00e7\u00e3o do cliente) agrega\u00e7\u00e3o autom\u00e1tica",
            "IA prev\u00ea equipe ideal por faixa hor\u00e1ria",
            "Previs\u00e3o de demanda e receita para compras e pessoal",
            "Insights de neg\u00f3cio detectam anomalias e oportunidades",
            "Scorecards KPI para vis\u00e3o r\u00e1pida da sa\u00fade do neg\u00f3cio",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Seguran\u00e7a e Auditoria", [
            "Log completo de acessos (login, opera\u00e7\u00f5es, detec\u00e7\u00e3o de anomalias)",
            "Limita\u00e7\u00e3o de taxa por IP bloqueia acessos n\u00e3o autorizados",
            "PII de clientes criptografada com AES",
            "Auditoria de seguran\u00e7a automatizada (12 verifica\u00e7\u00f5es di\u00e1rias)",
            "Controle de acesso baseado em fun\u00e7\u00f5es",
            "Rastreabilidade completa: hist\u00f3rico de opera\u00e7\u00f5es e mudan\u00e7as de turno",
            "Criptografia SSL e prote\u00e7\u00e3o CSRF integrados",
            "Integra\u00e7\u00e3o com servi\u00e7o de pagamento compat\u00edvel com PCI DSS",
        ], ACCENT_ORANGE)

    # ==========================================
    # Screenshot Slides
    # ==========================================
    screenshots_dir = os.path.dirname(os.path.abspath(__file__))

    def _add_desktop_slide(img_name, title, subtitle, bullets, header_prefix=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        full_title = f"{header_prefix}{title}" if header_prefix else title
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 full_title, font_size=28, bold=True, color=BROWN)
        add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)
        img_path = os.path.join(screenshots_dir, img_name)
        add_shape(slide, Inches(0.4), Inches(1.2), Inches(8.2), Inches(5.8),
                  LIGHT_GRAY, corner_radius=True)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.6))
            except Exception:
                add_text(slide, Inches(2), Inches(3.5), Inches(4), Inches(1),
                         f'[{title}]', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(0.5),
                 subtitle, font_size=18, bold=True, color=DARK)
        add_bullet_list(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(5),
                        bullets, font_size=12, icon="\u2713")

    # ── Customer-Facing Screens ──
    frontend_screens = [
        ('screenshots/pt/front_top_desktop.png', 'P\u00e1gina Inicial',
         'A primeira tela que os clientes veem',
         [
             'Banners hero destacam o apelo da sua loja',
             'Fluxo de reserva claro com m\u00faltiplas op\u00e7\u00f5es',
             'Ranking de equipe para promover membros populares',
             'Exibi\u00e7\u00e3o autom\u00e1tica de not\u00edcias e apari\u00e7\u00f5es na m\u00eddia',
             'Suporte em 7 idiomas para clientes internacionais',
         ]),
        ('screenshots/pt/front_staff_list_desktop.png', 'Perfis da Equipe e Info da Loja',
         'Apresente sua equipe de forma atraente',
         [
             'Fotos e descri\u00e7\u00f5es de perfil da equipe',
             'Exiba especialidades e qualifica\u00e7\u00f5es',
             'Link direto para o calend\u00e1rio de reservas',
             'Info da loja (hor\u00e1rio, mapa, como chegar)',
             'Atualiza\u00e7\u00e3o f\u00e1cil do conte\u00fado pelo painel admin',
         ]),
        ('screenshots/pt/front_news_desktop.png', 'Not\u00edcias e An\u00fancios',
         'Mantenha os clientes informados',
         [
             'Publique novos itens do menu, campanhas e promo\u00e7\u00f5es',
             'Publica\u00e7\u00e3o com um clique pelo painel admin',
             'Ordena\u00e7\u00e3o autom\u00e1tica por data para conte\u00fado atualizado',
             'Integra\u00e7\u00e3o com redes sociais para maior alcance',
             'Incentive visitas recorrentes com atualiza\u00e7\u00f5es oportunas',
         ]),
        ('screenshots/pt/front_shop_desktop.png', 'Loja Online',
         'Expanda seus canais de venda com EC',
         [
             'Venda seus produtos online',
             'Busca de produtos e filtragem por categoria',
             'Carrinho de compras para pedidos m\u00faltiplos',
             'Pagamento com cart\u00e3o de cr\u00e9dito e dinheiro eletr\u00f4nico',
             'Estoque unificado evita discrep\u00e2ncias',
         ]),
        ('screenshots/pt/table_order_vp.png', 'Pedidos QR na Mesa',
         'Pedido f\u00e1cil pelo smartphone',
         [
             'Basta escanear o QR code na mesa',
             'Menus com fotos transmitem o prato visualmente',
             'Abas por categoria para busca r\u00e1pida',
             'Fun\u00e7\u00e3o de carrinho para pedidos adicionais f\u00e1ceis',
             'Acompanhamento do status do pedido em tempo real',
         ]),
        ('screenshots/pt/booking_calendar_desktop.png', 'Calend\u00e1rio de Reservas',
         'Reservas online 24/7',
         [
             'Veja a disponibilidade da equipe no calend\u00e1rio de relance',
             'Clientes escolhem data e hor\u00e1rio preferidos',
             'Confirma\u00e7\u00e3o com pr\u00e9-pagamento previne aus\u00eancias',
             'Suporte a reservas por LINE e e-mail',
             'Notifica\u00e7\u00e3o autom\u00e1tica \u00e0 equipe na confirma\u00e7\u00e3o',
         ]),
    ]

    for img_name, title, subtitle, bullets in frontend_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="Vis\u00e3o do Cliente \u2014 ")

    # ── Mobile-Responsive ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Totalmente Responsivo para Celular", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "Tanto o site do cliente quanto o painel de administra\u00e7\u00e3o funcionam perfeitamente em smartphones.",
             font_size=16, color=DARK)

    mobile_items = [
        ('screenshots/pt/front_top_mobile_vp.png', 'P\u00e1gina Inicial do Cliente'),
        ('screenshots/pt/admin_dashboard_mobile_vp.png', 'Painel de Administra\u00e7\u00e3o'),
    ]

    for i, (img_name, label) in enumerate(mobile_items):
        x = Inches(1.5) + i * Inches(5.5)
        y = Inches(1.8)

        add_shape(slide, x, y, Inches(4.2), Inches(5.0), WHITE, corner_radius=True)

        img_path = os.path.join(screenshots_dir, img_name)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                x + Inches(0.6), y + Inches(0.15),
                                Inches(3.0), Inches(4.4))
            except Exception:
                add_text(slide, x + Inches(0.5), y + Inches(2), Inches(3.2), Inches(1),
                         f'[{label}]', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(4.5), Inches(4.2), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ── Admin Screens ──
    admin_screens = [
        ('screenshots/pt/dashboard_sales_vp.png', 'Painel de Vendas',
         'Vis\u00e3o geral do neg\u00f3cio de relance',
         [
             'Tend\u00eancias de receita di\u00e1ria / semanal / mensal / anual',
             'KPIs de reservas (quantidade, taxa de cancelamento) em tempo real',
             'Divis\u00e3o em tr\u00eas canais: Total, Loja EC, Menu na Loja',
             'Scorecards KPI para vis\u00e3o geral das principais m\u00e9tricas',
             'Insights de Neg\u00f3cio sugerem melhorias automaticamente',
         ]),
        ('screenshots/pt/dashboard_menu_eng_vp.png', 'An\u00e1lise de Vendas (An\u00e1lise de Menu)',
         'Otimiza\u00e7\u00e3o de lucro com IA',
         [
             'Engenharia de Menu: classifica\u00e7\u00e3o em 4 quadrantes Star / Plowhorse / Puzzle / Dog',
             'An\u00e1lise ABC: princ\u00edpio de Pareto para identificar os maiores geradores de receita',
             'Previs\u00e3o de Receita: modelo de IA Prophet prev\u00ea 14 dias \u00e0 frente',
             'Mapa de calor de vendas por hora para decis\u00f5es \u00f3timas de pessoal',
             'Tend\u00eancias de AOV (Ticket M\u00e9dio) para revis\u00e3o de estrat\u00e9gia de pre\u00e7os',
         ]),
        ('screenshots/pt/dashboard_rfm_vp.png', 'An\u00e1lise de Vendas (An\u00e1lise de Clientes)',
         'Compreens\u00e3o de clientes baseada em dados',
         [
             'An\u00e1lise RFM: segmenta\u00e7\u00e3o autom\u00e1tica de clientes por recencia, frequ\u00eancia e valor',
             'An\u00e1lise de Coorte: visualize taxas de reten\u00e7\u00e3o mensal de novos clientes',
             'An\u00e1lise de Cesta: descubra padr\u00f5es de venda cruzada e recomenda\u00e7\u00f5es',
             'Resumo por IA: resultados da an\u00e1lise resumidos em linguagem natural',
             'A\u00e7\u00f5es Recomendadas: IA sugere passos de melhoria espec\u00edficos automaticamente',
         ]),
        ('screenshots/pt/shift_calendar_vp.png', 'Calend\u00e1rio de Turnos',
         'Gest\u00e3o de turnos simplificada',
         [
             'Equipe envia prefer\u00eancias de turno pelo smartphone',
             'Vis\u00e3o de calend\u00e1rio com todas as escalas de relance',
             'Agendamento autom\u00e1tico com aloca\u00e7\u00e3o em um clique',
             'Detec\u00e7\u00e3o autom\u00e1tica de dias com pessoal insuficiente com alertas',
             'Turnos confirmados notificados automaticamente via LINE',
         ]),
        ('screenshots/pt/pos_vp.png', 'Caixa POS',
         'Experi\u00eancia de checkout suave',
         [
             'Abas por categoria para sele\u00e7\u00e3o r\u00e1pida de produtos',
             'Suporte a dinheiro, cart\u00e3o de cr\u00e9dito, PayPay e cart\u00e3o de transporte',
             'Gera\u00e7\u00e3o e impress\u00e3o autom\u00e1tica de recibos',
             'Integra\u00e7\u00e3o com display da cozinha para notifica\u00e7\u00f5es de pedidos',
             'Relat\u00f3rios de vendas di\u00e1rios e mensais gerados automaticamente',
         ]),
        ('screenshots/pt/inventory_vp.png', 'Gest\u00e3o de Estoque',
         'Previna rupturas',
         [
             'N\u00edveis de estoque em tempo real para todos os produtos',
             'Dedu\u00e7\u00e3o autom\u00e1tica ao receber pedidos',
             'Alertas e notifica\u00e7\u00f5es de estoque baixo',
             'Recebimento por QR code (sem scanner de c\u00f3digo de barras)',
             'Hist\u00f3rico completo de entradas/sa\u00eddas para conformidade de auditoria',
         ]),
        ('screenshots/pt/customer_feedback_vp.png', 'Feedback do Cliente (NPS)',
         'Quantifique a voz dos seus clientes',
         [
             'NPS (Net Promoter Score) agrega\u00e7\u00e3o autom\u00e1tica e tend\u00eancias',
             'Exibi\u00e7\u00e3o com cores: Promotor / Passivo / Detrator',
             'Feedback baseado em coment\u00e1rios para melhorias acion\u00e1veis',
             'Vinculado a dados de pedidos para an\u00e1lise de qualidade do servi\u00e7o',
             'Tend\u00eancias mensais e por membro da equipe de satisfa\u00e7\u00e3o',
         ]),
        ('screenshots/pt/attendance_board_vp.png', 'Quadro de Frequ\u00eancia',
         'Status de frequ\u00eancia em tempo real',
         [
             'Veja status de servi\u00e7o, pausa e n\u00e3o registrado de relance',
             'Atualiza\u00e7\u00e3o autom\u00e1tica a cada 30 segundos para rastreamento ao vivo',
             'Tr\u00eas m\u00e9todos de registro: QR code, PIN ou smartphone',
             'Classifica\u00e7\u00e3o autom\u00e1tica de horas extras, noturnas e feriados',
             'Fluxo cont\u00ednuo de dados de frequ\u00eancia para folha de pagamento',
         ]),
        ('screenshots/pt/iot_sensors_vp.png', 'Monitor de Sensores IoT',
         'Monitoramento ambiental em tempo real',
         [
             'Temperatura, umidade, press\u00e3o e n\u00edveis de g\u00e1s em tempo real',
             'Sensor de movimento PIR para contagem autom\u00e1tica de visitantes',
             'Notifica\u00e7\u00f5es de alerta em anomalias (vazamentos de g\u00e1s, altas temperaturas, etc.)',
             'Gr\u00e1ficos de s\u00e9ries temporais para acompanhar tend\u00eancias dos sensores',
             'Monitoramento de conectividade de dispositivos ESP32',
         ]),
    ]

    for img_name, title, subtitle, bullets in admin_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="Painel Admin \u2014 ")

    # ==========================================
    # Benefits Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Benef\u00edcios", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    merits = [
        ("Zero Perdas com Aus\u00eancias",
         "A confirma\u00e7\u00e3o com pr\u00e9-pagamento elimina cancelamentos de \u00faltima hora.\nAlterne entre pr\u00e9-pagamento / p\u00f3s-pagamento com um clique.",
         "\U0001F6E1"),
        ("Reservas F\u00e1ceis pelo LINE",
         "Clientes reservam direto no chat do LINE.\nSem necessidade de instalar novo aplicativo.",
         "\U0001F4AC"),
        ("Aumente sua Receita",
         "Pedidos QR aumentam a taxa de pedidos adicionais.\nSuporte multilíngue captura visitantes internacionais.",
         "\U0001F4C8"),
        ("Economize Horas de Trabalho",
         "Automatize reservas, pedidos, turnos e folha de pagamento.\nLibere sua equipe para focar no atendimento ao cliente.",
         "\u23F1"),
        ("Reduza Custos",
         "Livros de reserva, notas de pedido, cart\u00f5es de ponto...\nElimine custos com papel e esfor\u00e7o manual.",
         "\U0001F4B0"),
        ("Decis\u00f5es Baseadas em Dados",
         "Visualize vendas, reservas e desempenho da equipe.\nTome decis\u00f5es informadas com dados reais.",
         "\U0001F4CA"),
    ]

    for i, (title, desc, icon) in enumerate(merits):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), WHITE, corner_radius=True)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                 icon, font_size=30, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.35),
                 title, font_size=16, bold=True, color=BROWN)

        add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # Getting Started Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Como Come\u00e7ar", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    steps = [
        ("STEP 1", "Contate-nos", "Ligue ou use nosso\nformul\u00e1rio de contato"),
        ("STEP 2", "Consultoria", "Conhecemos seu fluxo\nde trabalho e desafios"),
        ("STEP 3", "Demo e Or\u00e7amento", "Veja o produto em a\u00e7\u00e3o.\nPropomos o melhor plano"),
        ("STEP 4", "Configura\u00e7\u00e3o", "Registre info da loja e equipe.\nAuxiliamos na configura\u00e7\u00e3o"),
        ("STEP 5", "No Ar!", "Ap\u00f3s uma breve orienta\u00e7\u00e3o,\nvoc\u00ea est\u00e1 pronto!"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.6)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BROWN
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            add_text(slide, x + Inches(1.9), y + Inches(0.2), Inches(0.6), Inches(0.6),
                     "\u2192", font_size=28, color=BROWN, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.2), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(1.5), LIGHT_GRAY, corner_radius=True)
    add_text(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1.2),
             "Comece a usar no mesmo dia.\n"
             "Sem hardware especial ou instala\u00e7\u00e3o de sistema. Acesse instantaneamente do seu smartphone, tablet ou PC.\n"
             "Oferecemos suporte cont\u00ednuo, treinamento e atualiza\u00e7\u00f5es ap\u00f3s a implanta\u00e7\u00e3o.",
             font_size=13, color=DARK)

    # ==========================================
    # Pricing Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Planos e Pre\u00e7os", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "* Os pre\u00e7os s\u00e3o para refer\u00eancia. Personalizamos planos conforme suas necessidades.",
             font_size=12, color=GRAY)

    plans = [
        ("Plano B\u00e1sico", "A partir de R$XX/m\u00eas", ACCENT_GREEN,
         "Para lojas pequenas", [
             "Gest\u00e3o de reservas (LINE e e-mail)",
             "Gest\u00e3o de equipe (at\u00e9 5)",
             "Gest\u00e3o do site",
             "Suporte por e-mail e chat",
         ]),
        ("Plano Padr\u00e3o", "A partir de R$XX/m\u00eas", ACCENT_BLUE,
         "Recomendado!", [
             "Todas as funcionalidades do Plano B\u00e1sico",
             "Pedidos QR na mesa",
             "Gest\u00e3o de estoque",
             "Gest\u00e3o de turnos",
             "Equipe ilimitada",
             "Suporte telef\u00f4nico",
         ]),
        ("Plano Premium", "A partir de R$XX/m\u00eas", BROWN,
         "Multi-loja / Todas as funcionalidades", [
             "Todas as funcionalidades do Plano Padr\u00e3o",
             "Folha de pagamento e frequ\u00eancia",
             "Painel de vendas e an\u00e1lises",
             "Suporte multilíngue (7 idiomas)",
             "Gest\u00e3o multi-loja",
             "Suporte priorit\u00e1rio",
             "Desenvolvimento personalizado",
         ]),
    ]

    for i, (name, price, color, badge, features) in enumerate(plans):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.2)

        card = add_shape(slide, x, y, w, h, LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, w, Inches(0.08), color)

        if badge == "Recomendado!":
            badge_shape = add_shape(slide, x + Inches(0.8), y - Inches(0.15), Inches(2.2), Inches(0.35), color, corner_radius=True)
            tf = badge_shape.text_frame
            p = tf.paragraphs[0]
            p.text = "\u2605 Recomendado!"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER

        add_text(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.5),
                 price, font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

        add_shape(slide, x + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(0.02), GRAY)

        add_text(slide, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(0.3),
                 badge if badge != "Recomendado!" else "Para lojas de m\u00e9dio porte",
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9),
                        w - Inches(0.6), h - Inches(2.1),
                        features, font_size=12, icon="\u2713")

    add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "* Taxas de configura\u00e7\u00e3o e personaliza\u00e7\u00e3o cotadas separadamente.  * Pre\u00e7os n\u00e3o incluem impostos.",
             font_size=10, color=GRAY)

    # ==========================================
    # FAQ Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Perguntas Frequentes", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    faqs = [
        ("P. Preciso de algum equipamento especial?",
         "R. N\u00e3o \u00e9 necess\u00e1rio hardware especial. Use no seu smartphone, tablet ou PC existente.\n"
         "   Para pedidos na mesa, basta imprimir QR codes pelo painel de administra\u00e7\u00e3o."),
        ("P. Funcion\u00e1rios sem conhecimento t\u00e9cnico conseguem usar?",
         "R. Sim. Se voc\u00ea sabe usar um smartphone, sabe usar o TimeBaiBai. A interface \u00e9 intuitiva e simples.\n"
         "   Tamb\u00e9m oferecemos treinamento pr\u00e1tico durante a implanta\u00e7\u00e3o."),
        ("P. Posso adicionar ou mudar funcionalidades depois?",
         "R. Com certeza. Voc\u00ea pode fazer upgrade ou mudar seu plano a qualquer momento.\n"
         "   Recomendamos come\u00e7ar com as funcionalidades b\u00e1sicas e expandir conforme necess\u00e1rio."),
        ("P. Meus dados est\u00e3o seguros?",
         "R. Os dados dos clientes s\u00e3o criptografados em repouso. Toda comunica\u00e7\u00e3o usa criptografia SSL.\n"
         "   Informa\u00e7\u00f5es de cart\u00e3o s\u00e3o processadas por servi\u00e7os compat\u00edveis com PCI DSS \u2014 nenhum dado de cart\u00e3o \u00e9 armazenado em nossos servidores."),
        ("P. Posso usar para m\u00faltiplas unidades?",
         "R. Sim. O Plano Premium suporta gest\u00e3o multi-loja.\n"
         "   Gerencie equipe e estoque por unidade enquanto visualiza relat\u00f3rios consolidados de todas as lojas."),
    ]

    for i, (q, a) in enumerate(faqs):
        y = Inches(1.2) + i * Inches(1.2)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(1.05), WHITE, corner_radius=True)

        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(11.73), Inches(0.3),
                 q, font_size=13, bold=True, color=BROWN)
        add_text(slide, Inches(0.8), y + Inches(0.4), Inches(11.73), Inches(0.6),
                 a, font_size=11, color=DARK)

    # ==========================================
    # Contact Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "Entre em Contato",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             "Solicite uma demo, or\u00e7amento ou tire suas d\u00favidas.",
             font_size=18, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.33), Inches(0.03), WHITE)

    contact_card = add_shape(slide, Inches(3), Inches(3.8), Inches(7.33), Inches(2.5), WHITE, corner_radius=True)

    contacts = [
        "TEL:     000-0000-0000",
        "Email:   info@example.com",
        "Web:    https://timebaibai.com",
        "",
        "Nome da Empresa",
        "Endere\u00e7o",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, Inches(4), Inches(4.0) + i * Inches(0.35),
                 Inches(5), Inches(0.35),
                 line, font_size=15 if i < 3 else 13,
                 bold=(i < 3), color=DARK if i < 3 else GRAY,
                 align=PP_ALIGN.LEFT)

    add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "\u00A9 2026 TimeBaiBai. All Rights Reserved.",
             font_size=10, color=BEIGE, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# Output
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    pptx_path = os.path.join(output_dir, 'TimeBaiBai_Service_Overview_PT.pptx')
    prs.save(pptx_path)
    print(f"\u2705 PPT saved: {pptx_path}")

    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            print(f"\u2705 PDF saved: {pdf_path}")
        else:
            raise FileNotFoundError("LibreOffice not found")
    except Exception as e:
        print(f"\u26A0\uFE0F  PDF conversion skipped ({e})")
        print(f"   To create PDF from PPTX:")
        print(f"   - Open in PowerPoint / Keynote and export as PDF")
        print(f"   - Or: brew install libreoffice && soffice --headless --convert-to pdf '{pptx_path}'")
