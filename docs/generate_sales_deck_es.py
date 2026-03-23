#!/usr/bin/env python3
"""
Sales Deck Generator — Spanish (Español) Version
Outputs PPT + PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ==============================
# Color Palette
# ==============================
BROWN = RGBColor(0x8C, 0x87, 0x6C)
BEIGE = RGBColor(0xF1, 0xF0, 0xEC)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="\u2713"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing


def add_feature_card(slide, left, top, width, height, title, items, accent_color=BROWN):
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=True)
    line = add_shape(slide, left + Inches(0.15), top + Inches(0.08),
                     width - Inches(0.3), Inches(0.04), accent_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.2),
             width - Inches(0.4), Inches(0.4),
             title, font_size=15, bold=True, color=accent_color)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6),
                    width - Inches(0.4), height - Inches(0.8),
                    items, font_size=11, icon="\u25CF")


def add_picture_fit(slide, img_path, left, top, max_width, max_height):
    with Image.open(img_path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = max_width / max_height
    if ratio > box_ratio:
        w = max_width
        h = max_width / ratio
    else:
        h = max_height
        w = max_height * ratio
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    return slide.shapes.add_picture(img_path, int(x), int(y), width=int(w), height=int(h))


def add_number_highlight(slide, left, top, number, label, color=BROWN):
    add_text(slide, left, top, Inches(2), Inches(0.6),
             number, font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2), Inches(0.3),
             label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ==============================
# Slide Creation
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # Slide 1: Cover
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Plataforma DX Todo-en-Uno para Tiendas",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "Reservas \u2022 Pedidos \u2022 Pagos \u2022 Gesti\u00f3n de Personal \u2022 N\u00f3mina\nTodo en un solo sistema.",
             font_size=22, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5), Inches(4.0), Inches(3.33), Inches(0.03), WHITE)

    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "TimeBaiBai  Presentaci\u00f3n del Servicio",
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "2026  /  Nombre de la Empresa",
             font_size=14, color=BEIGE, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 2: Pain Points
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "\u00bfLe resultan familiares estos problemas?", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    problems = [
        ("Reservas Dispersas",
         "Tel\u00e9fono, LINE, email... diferentes canales\ncausan ansiedad por reservas duplicadas.",
         "\U0001F4DE"),
        ("Ausencias y Cancelaciones de \u00daltimo Momento",
         "Los clientes reservan pero no se presentan.\nP\u00e9rdida de ingresos y preparaci\u00f3n en vano.",
         "\U0001F630"),
        ("Toma de Pedidos Sobrecargada",
         "No se dan abasto en horas punta.\nLos errores en pedidos se repiten.",
         "\U0001F4DD"),
        ("Puntos Ciegos en Inventario",
         "Los productos se agotan sin aviso.\nNo se sabe cu\u00e1ndo hacer nuevos pedidos.",
         "\U0001F4E6"),
    ]

    for i, (title, desc, icon) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.6) + row * Inches(2.6)

        card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), LIGHT_GRAY, corner_radius=True)

        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8),
                 icon, font_size=36, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK)

        add_text(slide, x + Inches(1.2), y + Inches(0.75), Inches(4), Inches(1.2),
                 desc, font_size=14, color=GRAY)

    # ==========================================
    # Slide 3: Solution — What is TimeBaiBai
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "TimeBaiBai lo resuelve todo", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Plataforma integral para todo lo que su tienda necesita. Comience con solo un smartphone.",
             font_size=16, color=DARK)

    highlights = [
        ("6+", "M\u00f3dulos Principales"),
        ("7", "Idiomas Soportados"),
        ("24h", "Reserva Autom\u00e1tica"),
        ("$0", "Ahorro en Papel"),
    ]
    for i, (num, label) in enumerate(highlights):
        x = Inches(1.2) + i * Inches(3)
        add_shape(slide, x - Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.4), WHITE, corner_radius=True)
        add_number_highlight(slide, x, Inches(2.5), num, label)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
             "Reservas, pedidos en mesa, inventario, gesti\u00f3n de turnos, n\u00f3mina, gesti\u00f3n del sitio web...\n"
             "Administre todo desde un solo panel \u2014 sin necesidad de usar m\u00faltiples herramientas.\n\n"
             "La configuraci\u00f3n es sencilla. Comience a usarlo el mismo d\u00eda que se crea su cuenta.\n"
             "No requiere instalaci\u00f3n complicada ni hardware especial. Acceda desde su tel\u00e9fono o computadora.",
             font_size=15, color=DARK)

    # ==========================================
    # Slide 4: Features (1/4) Reservations & Orders
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funciones Clave  \u2460  Reservas / Pedidos en Mesa",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Sistema de Reservas en L\u00ednea", [
            "Reserve por LINE \u2014 complete la reserva directamente desde el chat",
            "Reservas por email tambi\u00e9n disponibles para conveniencia del cliente",
            "Confirmaci\u00f3n con prepago \u2014 confirmaci\u00f3n autom\u00e1tica al pagar (previene ausencias)",
            "Cambio entre prepago / postpago con un clic en el administrador",
            "Calendarios de personal muestran disponibilidad en tiempo real",
            "Dos flujos de reserva: por fecha o por miembro del personal",
            "Check-in por c\u00f3digo QR (sin esfuerzo de recepci\u00f3n)",
            "Notificaci\u00f3n autom\u00e1tica por LINE al personal al confirmar reserva",
            "Soporte en 7 idiomas (japon\u00e9s, ingl\u00e9s, chino, coreano, espa\u00f1ol, portugu\u00e9s)",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Sistema de Pedidos QR en Mesa", [
            "Coloque c\u00f3digos QR en cada mesa para comenzar",
            "Los clientes escanean con su tel\u00e9fono \u2192 ven el men\u00fa \u2192 hacen pedido",
            "Men\u00fas con fotos ayudan a visualizar los platos",
            "Pesta\u00f1as por categor\u00eda para b\u00fasqueda r\u00e1pida",
            "Funci\u00f3n de carrito para pedidos adicionales f\u00e1ciles",
            "Seguimiento del estado del pedido en tiempo real",
            "Soporte de pago: efectivo, tarjeta de cr\u00e9dito y dinero electr\u00f3nico",
            "Generaci\u00f3n masiva de c\u00f3digos QR e impresi\u00f3n desde el administrador",
        ], ACCENT_GREEN)

    # ==========================================
    # Slide 5: Features (2/4) Inventory & Staff
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funciones Clave  \u2461  Inventario / Gesti\u00f3n de Personal y Turnos",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gesti\u00f3n de Inventario", [
            "Seguimiento de niveles de stock en tiempo real",
            "Descuento autom\u00e1tico de inventario vinculado a pedidos",
            "Alertas autom\u00e1ticas cuando el stock es bajo",
            "Recepci\u00f3n por c\u00f3digo QR (sin esc\u00e1ner de c\u00f3digo de barras)",
            "Historial completo de entradas/salidas de stock (listo para auditor\u00eda)",
            "Organizado por categor\u00eda de producto",
            "Inventario unificado para tienda EC y tienda f\u00edsica",
            "Funci\u00f3n de inventario para ajustes masivos",
        ], ACCENT_ORANGE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gesti\u00f3n de Personal y Turnos", [
            "El personal env\u00eda preferencias de turno desde su tel\u00e9fono",
            "Tres niveles de solicitud: Disponible / Preferido / No disponible",
            "Los gerentes finalizan turnos desde el panel de administraci\u00f3n",
            "Programaci\u00f3n autom\u00e1tica para asignaci\u00f3n eficiente",
            "Turnos confirmados notificados autom\u00e1ticamente por LINE",
            "Perfiles de personal con fotos y descripciones",
            "Tipos de personal: Cast (servicio) / Personal de Tienda",
            "Permisos por rol: Propietario / Gerente / Personal / Desarrollador",
        ], ACCENT_BLUE)

    # ==========================================
    # Slide 6: Features (3/4) Payroll & Website
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funciones Clave  \u2462  N\u00f3mina / Gesti\u00f3n del Sitio Web",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gesti\u00f3n de N\u00f3mina y Asistencia", [
            "Generaci\u00f3n autom\u00e1tica de datos de asistencia desde registros de turnos",
            "Clasificaci\u00f3n autom\u00e1tica de horas regulares, extras, nocturnas y festivas",
            "C\u00e1lculo de n\u00f3mina para empleados por hora y asalariados",
            "C\u00e1lculo autom\u00e1tico de primas de seguro social",
            "Aplicaci\u00f3n autom\u00e1tica de deducciones de impuestos",
            "Generaci\u00f3n de recibos de pago con un clic",
            "Exportaci\u00f3n CSV para transferencias bancarias",
            "Configuraci\u00f3n f\u00e1cil de asignaciones (transporte, vivienda, familia)",
        ], ACCENT_GREEN)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Gesti\u00f3n del Sitio Web y Marketing", [
            "Edite el sitio web de su tienda desde el panel de administraci\u00f3n",
            "Banners hero (sliders) para mostrar su tienda",
            "Actualizaci\u00f3n f\u00e1cil de noticias y apariciones en medios",
            "Colocaci\u00f3n y gesti\u00f3n de banners publicitarios",
            "Integraci\u00f3n de redes sociales (X / Instagram embeds)",
            "Bloques HTML personalizados para dise\u00f1os flexibles",
            "P\u00e1ginas legales editables desde el administrador",
            "Tienda EC integrada para ventas en l\u00ednea",
        ], ACCENT_ORANGE)

    # ==========================================
    # Slide 7: Features (4/4) Analytics & Security
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Funciones Clave  \u2463  An\u00e1lisis de Datos / Seguridad",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "An\u00e1lisis de Datos e IA", [
            "Panel de ventas con gr\u00e1ficos de tendencia diario/semanal/mensual/anual",
            "Ingenier\u00eda de Men\u00fa (clasificaci\u00f3n Star/Dog) para mejorar rentabilidad",
            "An\u00e1lisis ABC, RFM, cohorte y de canasta",
            "NPS (satisfacci\u00f3n del cliente) agregaci\u00f3n autom\u00e1tica",
            "IA predice dotaci\u00f3n \u00f3ptima de personal por franja horaria",
            "Previsi\u00f3n de demanda e ingresos para compras y personal",
            "Insights de negocio detectan anomal\u00edas y oportunidades autom\u00e1ticamente",
            "Cuadros de mando KPI para visi\u00f3n r\u00e1pida del negocio",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Seguridad y Auditor\u00eda", [
            "Registro completo de accesos (login, operaciones, detecci\u00f3n de anomal\u00edas)",
            "Limitaci\u00f3n de tasa basada en IP bloquea accesos no autorizados",
            "PII de clientes cifrada con AES",
            "Auditor\u00eda de seguridad automatizada (12 verificaciones diarias)",
            "Control de acceso basado en roles",
            "Trazabilidad completa: historial de operaciones y cambios de turno",
            "Cifrado SSL y protecci\u00f3n CSRF integrados",
            "Integraci\u00f3n con servicio de pago compatible con PCI DSS",
        ], ACCENT_ORANGE)

    # ==========================================
    # Screenshot Slides
    # ==========================================
    screenshots_dir = os.path.dirname(os.path.abspath(__file__))

    def _add_desktop_slide(img_name, title, subtitle, bullets, header_prefix=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        full_title = f"{header_prefix}{title}" if header_prefix else title
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 full_title, font_size=28, bold=True, color=BROWN)
        add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)
        img_path = os.path.join(screenshots_dir, img_name)
        add_shape(slide, Inches(0.4), Inches(1.2), Inches(8.2), Inches(5.8),
                  LIGHT_GRAY, corner_radius=True)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.6))
            except Exception:
                add_text(slide, Inches(2), Inches(3.5), Inches(4), Inches(1),
                         f'[{title}]', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(0.5),
                 subtitle, font_size=18, bold=True, color=DARK)
        add_bullet_list(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(5),
                        bullets, font_size=12, icon="\u2713")

    # -- Customer-Facing Screens --
    frontend_screens = [
        ('screenshots/es/front_top_desktop.png', 'P\u00e1gina Principal',
         'La primera pantalla que ven los clientes',
         [
             'Banners hero muestran el atractivo de su tienda',
             'Flujo de reserva claro con m\u00faltiples opciones',
             'Ranking de personal para promover miembros populares',
             'Noticias y apariciones en medios se muestran autom\u00e1ticamente',
             'Soporte en 7 idiomas para clientes internacionales',
         ]),
        ('screenshots/es/front_staff_list_desktop.png', 'Perfiles del Personal e Info de Tienda',
         'Presente a su personal de forma atractiva',
         [
             'Fotos y descripciones del perfil del personal',
             'Muestre especialidades y calificaciones',
             'Enlace directo al calendario de reservas',
             'Info de tienda (horarios, mapa, indicaciones de acceso)',
             'Actualizaci\u00f3n f\u00e1cil de contenido desde el panel de admin',
         ]),
        ('screenshots/es/front_news_desktop.png', 'Noticias y Anuncios',
         'Mantenga informados a sus clientes',
         [
             'Publique nuevos platos, campa\u00f1as y promociones',
             'Publicaci\u00f3n con un clic desde el panel de admin',
             'Ordenado autom\u00e1ticamente por fecha para contenido fresco',
             'Integraci\u00f3n con redes sociales para mayor alcance',
             'Impulse visitas recurrentes con actualizaciones oportunas',
         ]),
        ('screenshots/es/front_shop_desktop.png', 'Tienda en L\u00ednea',
         'Ampl\u00ede sus canales de venta con EC',
         [
             'Venda sus productos en l\u00ednea',
             'B\u00fasqueda de productos y filtrado por categor\u00eda',
             'Carrito de compras para compras de m\u00faltiples art\u00edculos',
             'Pagos con tarjeta de cr\u00e9dito y dinero electr\u00f3nico',
             'Inventario unificado previene discrepancias de stock',
         ]),
        ('screenshots/es/table_order_vp.png', 'Pedidos QR en Mesa',
         'Pedidos f\u00e1ciles desde el smartphone',
         [
             'Solo escanee el c\u00f3digo QR en la mesa',
             'Men\u00fas con fotos transmiten el plato visualmente',
             'Pesta\u00f1as por categor\u00eda para descubrir art\u00edculos r\u00e1pidamente',
             'Funci\u00f3n de carrito para pedidos adicionales f\u00e1ciles',
             'Seguimiento del estado del pedido en tiempo real',
         ]),
        ('screenshots/es/booking_calendar_desktop.png', 'Calendario de Reservas',
         'Reservas en l\u00ednea 24/7',
         [
             'Vea disponibilidad del personal en el calendario de un vistazo',
             'Los clientes eligen su fecha y horario preferido',
             'Confirmaci\u00f3n con prepago previene ausencias',
             'Soporte de reservas por LINE y email',
             'Notificaci\u00f3n autom\u00e1tica al personal al confirmar',
         ]),
    ]

    for img_name, title, subtitle, bullets in frontend_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="Vista del Cliente \u2014 ")

    # -- Mobile-Responsive --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Totalmente Adaptable a M\u00f3viles", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "Tanto el sitio del cliente como el panel de administraci\u00f3n funcionan perfectamente en smartphones.",
             font_size=16, color=DARK)

    mobile_items = [
        ('screenshots/es/front_top_mobile_vp.png', 'P\u00e1gina Principal del Cliente'),
        ('screenshots/es/admin_dashboard_mobile_vp.png', 'Panel de Administraci\u00f3n'),
    ]

    for i, (img_name, label) in enumerate(mobile_items):
        x = Inches(1.5) + i * Inches(5.5)
        y = Inches(1.8)

        add_shape(slide, x, y, Inches(4.2), Inches(5.0), WHITE, corner_radius=True)

        img_path = os.path.join(screenshots_dir, img_name)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                x + Inches(0.6), y + Inches(0.15),
                                Inches(3.0), Inches(4.4))
            except Exception:
                add_text(slide, x + Inches(0.5), y + Inches(2), Inches(3.2), Inches(1),
                         f'[{label}]', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(4.5), Inches(4.2), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # -- Admin Screens --
    admin_screens = [
        ('screenshots/es/dashboard_sales_vp.png', 'Panel de Ventas',
         'Visi\u00f3n general del negocio de un vistazo',
         [
             'Tendencias de ingresos diario / semanal / mensual / anual',
             'KPIs de reservas (cantidad, tasa de cancelaci\u00f3n) en tiempo real',
             'Desglose de tres canales: Total, Tienda EC, Men\u00fa en Tienda',
             'Cuadros de mando KPI para visi\u00f3n general de m\u00e9tricas clave',
             'Insights de Negocio sugieren mejoras autom\u00e1ticamente',
         ]),
        ('screenshots/es/dashboard_menu_eng_vp.png', 'An\u00e1lisis de Ventas (An\u00e1lisis de Men\u00fa)',
         'Optimizaci\u00f3n de rentabilidad con IA',
         [
             'Ingenier\u00eda de Men\u00fa: clasificaci\u00f3n en 4 cuadrantes Star / Plowhorse / Puzzle / Dog',
             'An\u00e1lisis ABC: principio de Pareto para identificar principales generadores de ingresos',
             'Previsi\u00f3n de Ingresos: modelo IA Prophet predice 14 d\u00edas adelante',
             'Mapa de calor de ventas por hora para decisiones \u00f3ptimas de personal',
             'Tendencias de AOV (Valor Promedio de Pedido) para revisi\u00f3n de estrategia de precios',
         ]),
        ('screenshots/es/dashboard_rfm_vp.png', 'An\u00e1lisis de Ventas (An\u00e1lisis de Clientes)',
         'Comprensi\u00f3n del cliente basada en datos',
         [
             'An\u00e1lisis RFM: segmentaci\u00f3n autom\u00e1tica de clientes por recencia, frecuencia, monto',
             'An\u00e1lisis de Cohorte: visualice tasas de retenci\u00f3n mensual de nuevos clientes',
             'An\u00e1lisis de Canasta: descubra patrones de venta cruzada y recomendaciones',
             'Resumen IA: resultados del an\u00e1lisis resumidos en lenguaje natural',
             'Acciones Recomendadas: la IA sugiere autom\u00e1ticamente pasos espec\u00edficos de mejora',
         ]),
        ('screenshots/es/shift_calendar_vp.png', 'Calendario de Turnos',
         'Gesti\u00f3n de turnos simplificada',
         [
             'El personal env\u00eda preferencias de turno desde su smartphone',
             'Vista de calendario de todos los horarios del personal de un vistazo',
             'Programaci\u00f3n autom\u00e1tica con asignaci\u00f3n en un clic',
             'Detecci\u00f3n autom\u00e1tica de d\u00edas con falta de personal con alertas',
             'Turnos confirmados notificados autom\u00e1ticamente por LINE',
         ]),
        ('screenshots/es/pos_vp.png', 'Caja Registradora POS',
         'Experiencia de cobro fluida',
         [
             'Pesta\u00f1as por categor\u00eda para selecci\u00f3n r\u00e1pida de productos',
             'Soporte de efectivo, tarjeta de cr\u00e9dito, PayPay e IC de transporte',
             'Generaci\u00f3n e impresi\u00f3n autom\u00e1tica de recibos',
             'Integraci\u00f3n con pantalla de cocina para notificaciones de pedidos',
             'Reportes de ventas diarios y mensuales generados autom\u00e1ticamente',
         ]),
        ('screenshots/es/inventory_vp.png', 'Gesti\u00f3n de Inventario',
         'Prevenga agotamientos de stock',
         [
             'Niveles de stock en tiempo real para todos los productos',
             'Descuento autom\u00e1tico al recibir pedidos',
             'Alertas y notificaciones de stock bajo',
             'Recepci\u00f3n por c\u00f3digo QR (sin esc\u00e1ner de c\u00f3digo de barras)',
             'Historial completo de entradas/salidas para cumplimiento de auditor\u00eda',
         ]),
        ('screenshots/es/customer_feedback_vp.png', 'Retroalimentaci\u00f3n de Clientes (NPS)',
         'Cuantifique la voz de sus clientes',
         [
             'NPS (Net Promoter Score) agregaci\u00f3n autom\u00e1tica y tendencias',
             'Visualizaci\u00f3n con c\u00f3digos de color: Promotor / Pasivo / Detractor',
             'Retroalimentaci\u00f3n basada en comentarios para mejoras accionables',
             'Vinculado a datos de pedidos para an\u00e1lisis de calidad del servicio',
             'Tendencias de satisfacci\u00f3n mensual y por miembro del personal',
         ]),
        ('screenshots/es/attendance_board_vp.png', 'Tablero de Asistencia',
         'Estado de asistencia del personal en tiempo real',
         [
             'Vea de un vistazo qui\u00e9n est\u00e1 de turno, en descanso o sin fichar',
             'Actualizaci\u00f3n autom\u00e1tica cada 30 segundos para seguimiento en vivo',
             'Tres m\u00e9todos de fichaje: c\u00f3digo QR, PIN o smartphone',
             'Clasificaci\u00f3n autom\u00e1tica de horas extra, nocturnas y festivas',
             'Flujo de datos continuo desde asistencia hasta n\u00f3mina',
         ]),
        ('screenshots/es/iot_sensors_vp.png', 'Monitor de Sensores IoT',
         'Monitoreo del ambiente de la tienda en tiempo real',
         [
             'Temperatura, humedad, presi\u00f3n, niveles de gas en tiempo real',
             'Sensor de movimiento PIR para conteo autom\u00e1tico de visitantes',
             'Notificaciones de alerta ante anomal\u00edas (fugas de gas, altas temperaturas, etc.)',
             'Gr\u00e1ficos de series de tiempo para seguimiento de tendencias de sensores',
             'Monitoreo de conectividad de dispositivos ESP32',
         ]),
    ]

    for img_name, title, subtitle, bullets in admin_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="Panel de Admin \u2014 ")

    # ==========================================
    # Benefits Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Beneficios", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    merits = [
        ("Cero P\u00e9rdidas por Ausencias",
         "La confirmaci\u00f3n con prepago elimina cancelaciones de \u00faltimo momento.\nCambie entre prepago / postpago con un clic.",
         "\U0001F6E1"),
        ("Reservas F\u00e1ciles por LINE",
         "Los clientes reservan directamente desde el chat de LINE.\nNo requiere instalaci\u00f3n de nueva aplicaci\u00f3n.",
         "\U0001F4AC"),
        ("Aumente sus Ingresos",
         "Los pedidos QR aumentan la tasa de pedidos adicionales.\nEl soporte multiling\u00fce captura visitantes internacionales.",
         "\U0001F4C8"),
        ("Ahorre Horas de Trabajo",
         "Automatice reservas, pedidos, turnos y n\u00f3mina.\nLibere a su personal para enfocarse en el servicio al cliente.",
         "\u23F1"),
        ("Reduzca Costos",
         "Libretas de reservas, notas de pedido, tarjetas de tiempo...\nElimine costos de papel y esfuerzo manual.",
         "\U0001F4B0"),
        ("Decisiones Basadas en Datos",
         "Visualice ventas, reservas y rendimiento del personal.\nTome decisiones informadas respaldadas por datos reales.",
         "\U0001F4CA"),
    ]

    for i, (title, desc, icon) in enumerate(merits):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), WHITE, corner_radius=True)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                 icon, font_size=30, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.35),
                 title, font_size=16, bold=True, color=BROWN)

        add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # Getting Started Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "C\u00f3mo Empezar", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    steps = [
        ("PASO 1", "Cont\u00e1ctenos", "Ll\u00e1menos o use\nnuestro formulario de contacto"),
        ("PASO 2", "Consulta", "Conocemos su flujo de trabajo\ny puntos de dolor actuales"),
        ("PASO 3", "Demo y Cotizaci\u00f3n", "Vea el producto en acci\u00f3n.\nLe proponemos el mejor plan"),
        ("PASO 4", "Configuraci\u00f3n", "Registre info de tienda y personal.\nLe asistimos en la configuraci\u00f3n"),
        ("PASO 5", "\u00a1En Vivo!", "Despu\u00e9s de una breve gu\u00eda,\n\u00a1est\u00e1 listo para comenzar!"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.6)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BROWN
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            add_text(slide, x + Inches(1.9), y + Inches(0.2), Inches(0.6), Inches(0.6),
                     "\u2192", font_size=28, color=BROWN, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.2), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(1.5), LIGHT_GRAY, corner_radius=True)
    add_text(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1.2),
             "Comience a usarlo el mismo d\u00eda.\n"
             "No se requiere hardware especial ni instalaci\u00f3n de sistema. Acceda al instante desde su smartphone, tablet o PC.\n"
             "Ofrecemos soporte continuo, capacitaci\u00f3n y actualizaciones despu\u00e9s del despliegue.",
             font_size=13, color=DARK)

    # ==========================================
    # Pricing Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Planes y Precios", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "* Los precios son de referencia. Personalizamos planes seg\u00fan sus necesidades.",
             font_size=12, color=GRAY)

    plans = [
        ("Plan B\u00e1sico", "Desde $XX/mes", ACCENT_GREEN,
         "Para tiendas peque\u00f1as", [
             "Gesti\u00f3n de reservas (LINE y email)",
             "Gesti\u00f3n de personal (hasta 5)",
             "Gesti\u00f3n del sitio web",
             "Soporte por email y chat",
         ]),
        ("Plan Est\u00e1ndar", "Desde $XX/mes", ACCENT_BLUE,
         "\u00a1Recomendado!", [
             "Todas las funciones del Plan B\u00e1sico",
             "Pedidos QR en mesa",
             "Gesti\u00f3n de inventario",
             "Gesti\u00f3n de turnos",
             "Personal ilimitado",
             "Soporte telef\u00f3nico",
         ]),
        ("Plan Premium", "Desde $XX/mes", BROWN,
         "Multi-tienda / Todas las funciones", [
             "Todas las funciones del Plan Est\u00e1ndar",
             "N\u00f3mina y asistencia",
             "Panel de ventas y an\u00e1lisis",
             "Soporte multiling\u00fce (7 idiomas)",
             "Gesti\u00f3n multi-tienda",
             "Soporte prioritario",
             "Desarrollo personalizado",
         ]),
    ]

    for i, (name, price, color, badge, features) in enumerate(plans):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.2)

        card = add_shape(slide, x, y, w, h, LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, w, Inches(0.08), color)

        if badge == "\u00a1Recomendado!":
            badge_shape = add_shape(slide, x + Inches(0.8), y - Inches(0.15), Inches(2.2), Inches(0.35), color, corner_radius=True)
            tf = badge_shape.text_frame
            p = tf.paragraphs[0]
            p.text = "\u2605 \u00a1Recomendado!"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER

        add_text(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.5),
                 price, font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

        add_shape(slide, x + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(0.02), GRAY)

        add_text(slide, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(0.3),
                 badge if badge != "\u00a1Recomendado!" else "Para tiendas medianas",
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9),
                        w - Inches(0.6), h - Inches(2.1),
                        features, font_size=12, icon="\u2713")

    add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "* Tarifas de instalaci\u00f3n y personalizaci\u00f3n se cotizan por separado.  * Precios no incluyen impuestos.",
             font_size=10, color=GRAY)

    # ==========================================
    # FAQ Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Preguntas Frecuentes", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    faqs = [
        ("P. \u00bfNecesito alg\u00fan equipo especial?",
         "R. No se requiere hardware especial. \u00daselo en su smartphone, tablet o PC existente.\n"
         "   Para pedidos en mesa, simplemente imprima c\u00f3digos QR desde el panel de administraci\u00f3n."),
        ("P. \u00bfPuede usarlo personal sin conocimientos t\u00e9cnicos?",
         "R. S\u00ed. Si puede usar un smartphone, puede usar TimeBaiBai. La interfaz es intuitiva y simple.\n"
         "   Tambi\u00e9n proporcionamos capacitaci\u00f3n pr\u00e1ctica durante la incorporaci\u00f3n."),
        ("P. \u00bfPuedo agregar o cambiar funciones despu\u00e9s?",
         "R. Por supuesto. Puede actualizar o cambiar su plan en cualquier momento.\n"
         "   Recomendamos comenzar con las funciones b\u00e1sicas y escalar seg\u00fan sea necesario."),
        ("P. \u00bfQu\u00e9 tan seguros est\u00e1n mis datos?",
         "R. Los datos de clientes est\u00e1n cifrados en reposo. Toda comunicaci\u00f3n usa cifrado SSL.\n"
         "   La informaci\u00f3n de tarjetas se procesa por servicios compatibles con PCI DSS \u2014 no almacenamos datos de tarjetas en nuestros servidores."),
        ("P. \u00bfPuedo usarlo para m\u00faltiples ubicaciones?",
         "R. S\u00ed. El plan Premium soporta gesti\u00f3n multi-tienda.\n"
         "   Gestione personal e inventario por ubicaci\u00f3n mientras ve reportes consolidados."),
    ]

    for i, (q, a) in enumerate(faqs):
        y = Inches(1.2) + i * Inches(1.2)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(1.05), WHITE, corner_radius=True)

        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(11.73), Inches(0.3),
                 q, font_size=13, bold=True, color=BROWN)
        add_text(slide, Inches(0.8), y + Inches(0.4), Inches(11.73), Inches(0.6),
                 a, font_size=11, color=DARK)

    # ==========================================
    # Contact Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "Cont\u00e1ctenos",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             "Solicite una demo, cotizaci\u00f3n o consulte cualquier duda.",
             font_size=18, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.33), Inches(0.03), WHITE)

    contact_card = add_shape(slide, Inches(3), Inches(3.8), Inches(7.33), Inches(2.5), WHITE, corner_radius=True)

    contacts = [
        "TEL:     000-0000-0000",
        "Email:   info@example.com",
        "Web:    https://timebaibai.com",
        "",
        "Nombre de la Empresa",
        "Direcci\u00f3n",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, Inches(4), Inches(4.0) + i * Inches(0.35),
                 Inches(5), Inches(0.35),
                 line, font_size=15 if i < 3 else 13,
                 bold=(i < 3), color=DARK if i < 3 else GRAY,
                 align=PP_ALIGN.LEFT)

    add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "\u00A9 2026 TimeBaiBai. All Rights Reserved.",
             font_size=10, color=BEIGE, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# Output
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    pptx_path = os.path.join(output_dir, 'TimeBaiBai_Service_Overview_ES.pptx')
    prs.save(pptx_path)
    print(f"\u2705 PPT saved: {pptx_path}")

    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            print(f"\u2705 PDF saved: {pdf_path}")
        else:
            raise FileNotFoundError("LibreOffice not found")
    except Exception as e:
        print(f"\u26A0\uFE0F  PDF conversion skipped ({e})")
        print(f"   To create PDF from PPTX:")
        print(f"   - Open in PowerPoint / Keynote and export as PDF")
        print(f"   - Or: brew install libreoffice && soffice --headless --convert-to pdf '{pptx_path}'")
