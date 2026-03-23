#!/usr/bin/env python3
"""
Sales Deck Generator — English Version
Outputs PPT + PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ==============================
# Color Palette
# ==============================
BROWN = RGBColor(0x8C, 0x87, 0x6C)
BEIGE = RGBColor(0xF1, 0xF0, 0xEC)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="\u2713"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing


def add_feature_card(slide, left, top, width, height, title, items, accent_color=BROWN):
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=True)
    line = add_shape(slide, left + Inches(0.15), top + Inches(0.08),
                     width - Inches(0.3), Inches(0.04), accent_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.2),
             width - Inches(0.4), Inches(0.4),
             title, font_size=15, bold=True, color=accent_color)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6),
                    width - Inches(0.4), height - Inches(0.8),
                    items, font_size=11, icon="\u25CF")


def add_picture_fit(slide, img_path, left, top, max_width, max_height):
    with Image.open(img_path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = max_width / max_height
    if ratio > box_ratio:
        w = max_width
        h = max_width / ratio
    else:
        h = max_height
        w = max_height * ratio
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    return slide.shapes.add_picture(img_path, int(x), int(y), width=int(w), height=int(h))


def add_number_highlight(slide, left, top, number, label, color=BROWN):
    add_text(slide, left, top, Inches(2), Inches(0.6),
             number, font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2), Inches(0.3),
             label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ==============================
# Slide Creation
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # Slide 1: Cover
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "All-in-One Store DX Platform",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "Reservations \u2022 Orders \u2022 Payments \u2022 Staff Management \u2022 Payroll\nAll in one system.",
             font_size=22, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5), Inches(4.0), Inches(3.33), Inches(0.03), WHITE)

    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "TimeBaiBai  Service Overview",
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "2026  /  Company Name Inc.",
             font_size=14, color=BEIGE, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 2: Pain Points
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Do any of these sound familiar?", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    problems = [
        ("Scattered Reservations",
         "Phone, LINE, email... different booking methods\ncause double-booking anxiety.",
         "\U0001F4DE"),
        ("No-Shows & Last-Minute Cancellations",
         "Customers book but don't show up.\nLost revenue and wasted prep.",
         "\U0001F630"),
        ("Overwhelmed Order Taking",
         "Can't take orders during rush hour.\nOrder mistakes keep happening.",
         "\U0001F4DD"),
        ("Inventory Blind Spots",
         "Items run out without warning.\nNo idea when to reorder.",
         "\U0001F4E6"),
    ]

    for i, (title, desc, icon) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.6) + row * Inches(2.6)

        card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), LIGHT_GRAY, corner_radius=True)

        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8),
                 icon, font_size=36, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK)

        add_text(slide, x + Inches(1.2), y + Inches(0.75), Inches(4), Inches(1.2),
                 desc, font_size=14, color=GRAY)

    # ==========================================
    # Slide 3: Solution — What is TimeBaiBai
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "TimeBaiBai solves it all", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "One-stop platform for everything your store needs. Get started with just a smartphone.",
             font_size=16, color=DARK)

    highlights = [
        ("6+", "Core Modules"),
        ("7", "Languages Supported"),
        ("24h", "Auto Booking"),
        ("$0", "Paper Cost Savings"),
    ]
    for i, (num, label) in enumerate(highlights):
        x = Inches(1.2) + i * Inches(3)
        add_shape(slide, x - Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.4), WHITE, corner_radius=True)
        add_number_highlight(slide, x, Inches(2.5), num, label)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
             "Reservations, table ordering, inventory, shift management, payroll, website management...\n"
             "Manage everything from a single dashboard \u2014 no more juggling separate tools.\n\n"
             "Setup is simple. Start using it the same day your account is created.\n"
             "No complicated installation or hardware required. Just access it from your phone or PC.",
             font_size=15, color=DARK)

    # ==========================================
    # Slide 4: Features (1/4) Reservations & Orders
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Key Features  \u2460  Reservations / Table Ordering",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Online Reservation System", [
            "Book via LINE \u2014 complete the reservation right from the chat",
            "Email reservations also supported for customer convenience",
            "Prepaid booking confirmation \u2014 auto-confirm on payment (prevents no-shows)",
            "Switch between prepaid / postpaid with one click in admin",
            "Staff calendars show real-time availability",
            "Two booking flows: by date or by staff member",
            "QR code check-in (zero reception effort)",
            "Auto LINE notification to staff on booking confirmation",
            "7-language support (Japanese, English, Chinese, Korean, Spanish, Portuguese)",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "QR Table Ordering System", [
            "Place QR codes on each table to get started",
            "Customers scan with their phone \u2192 browse menu \u2192 order",
            "Photo menus help customers visualize dishes",
            "Category tabs for quick item discovery",
            "Cart function for easy additional orders",
            "Real-time order status tracking",
            "Cash, credit card, and e-money payment support",
            "Bulk QR code generation and printable download from admin",
        ], ACCENT_GREEN)

    # ==========================================
    # Slide 5: Features (2/4) Inventory & Staff
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Key Features  \u2461  Inventory / Staff & Shift Management",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Inventory Management", [
            "Real-time stock level tracking for all products",
            "Auto-deduct inventory linked to incoming orders",
            "Automatic alerts when stock runs low",
            "QR code receiving (no barcode scanner needed)",
            "Full history of all stock-in / stock-out transactions (audit-ready)",
            "Organized by product category for easy management",
            "Unified inventory for EC shop and in-store stock",
            "Stocktake function for bulk adjustments",
        ], ACCENT_ORANGE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Staff & Shift Management", [
            "Staff submit shift preferences from their phones",
            "Three-level requests: Available / Preferred / Unavailable",
            "Managers finalize shifts from the admin dashboard",
            "Auto-scheduling for efficient staff allocation",
            "Confirmed shifts auto-notified to staff via LINE",
            "Staff profiles with photos and descriptions",
            "Staff types: Cast (service) / Store Staff",
            "Role-based permissions: Owner / Manager / Staff / Developer",
        ], ACCENT_BLUE)

    # ==========================================
    # Slide 6: Features (3/4) Payroll & Website
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Key Features  \u2462  Payroll / Website Management",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Payroll & Attendance Management", [
            "Auto-generate attendance data from shift records",
            "Auto-classify regular, overtime, night, and holiday hours",
            "Payroll calculation for hourly and salaried employees",
            "Auto-calculate social insurance premiums",
            "Auto-apply income tax and resident tax deductions",
            "One-click payslip generation",
            "Bank transfer CSV export (Zengin format)",
            "Easy setup for allowances (commuting, housing, family)",
        ], ACCENT_GREEN)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Website & Marketing Management", [
            "Edit your store's website from the admin panel",
            "Hero banners (sliders) to showcase your store",
            "Easy update for news and media features",
            "Banner ad placement and management",
            "Social media integration (X / Instagram embeds)",
            "Custom HTML blocks for flexible layouts",
            "Privacy policy and legal pages editable from admin",
            "Built-in EC shop for online sales",
        ], ACCENT_ORANGE)

    # ==========================================
    # Slide 7: Features (4/4) Analytics & Security
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Key Features  \u2463  Data Analytics / Security",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "Data Analytics & AI", [
            "Sales dashboard with daily/weekly/monthly/yearly trend charts",
            "Menu Engineering (Star/Dog classification) for profit improvement",
            "ABC analysis, RFM analysis, cohort analysis, basket analysis",
            "NPS (customer satisfaction) auto-aggregation for retention",
            "AI predicts optimal staffing by time slot",
            "Demand & revenue forecasting for procurement and staffing",
            "Business insights auto-detect anomalies and opportunities",
            "KPI scorecards for at-a-glance business health",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "Security & Audit", [
            "Full access logging (login, operations, anomaly detection)",
            "IP-based rate limiting blocks unauthorized access automatically",
            "Customer PII encrypted with AES for secure storage",
            "Automated security audit (12 checks daily)",
            "Role-based access control (Owner / Manager / Staff / Developer)",
            "Complete audit trail: operation history and shift change logs",
            "SSL encryption and CSRF protection built in",
            "PCI DSS-compliant external payment service integration",
        ], ACCENT_ORANGE)

    # ==========================================
    # Screenshot Slides
    # ==========================================
    screenshots_dir = os.path.dirname(os.path.abspath(__file__))

    def _add_desktop_slide(img_name, title, subtitle, bullets, header_prefix=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        full_title = f"{header_prefix}{title}" if header_prefix else title
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 full_title, font_size=28, bold=True, color=BROWN)
        add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)
        img_path = os.path.join(screenshots_dir, img_name)
        add_shape(slide, Inches(0.4), Inches(1.2), Inches(8.2), Inches(5.8),
                  LIGHT_GRAY, corner_radius=True)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.6))
            except Exception:
                add_text(slide, Inches(2), Inches(3.5), Inches(4), Inches(1),
                         f'[{title}]', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(0.5),
                 subtitle, font_size=18, bold=True, color=DARK)
        add_bullet_list(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(5),
                        bullets, font_size=12, icon="\u2713")

    # ── Customer-Facing Screens ──
    frontend_screens = [
        ('screenshots/en/front_top_desktop.png', 'Top Page',
         'The first screen customers see',
         [
             'Hero banners showcase your store\'s appeal',
             'Clear booking flow with multiple options',
             'Staff ranking to promote popular cast members',
             'Auto-display news and media features',
             '7-language support for international customers',
         ]),
        ('screenshots/en/front_staff_list_desktop.png', 'Staff Profiles & Store Info',
         'Present your staff attractively',
         [
             'Staff photos and profile descriptions',
             'Display specialties and qualifications',
             'Direct link to booking calendar',
             'Store info (hours, map, access directions)',
             'Easy content updates from admin panel',
         ]),
        ('screenshots/en/front_news_desktop.png', 'News & Announcements',
         'Keep customers informed',
         [
             'Post new menu items, campaigns, and promotions',
             'One-click publishing from admin panel',
             'Auto-sorted by date for fresh content',
             'Social media integration for broader reach',
             'Drive repeat visits with timely updates',
         ]),
        ('screenshots/en/front_shop_desktop.png', 'Online Shop',
         'Expand your sales channels with EC',
         [
             'Sell your products online',
             'Product search and category filtering',
             'Shopping cart for multi-item purchases',
             'Credit card and e-money payments',
             'Unified inventory prevents stock discrepancies',
         ]),
        ('screenshots/en/table_order_vp.png', 'QR Table Ordering',
         'Easy ordering from smartphone',
         [
             'Just scan the QR code on the table',
             'Photo menus convey the dish visually',
             'Category tabs for quick item discovery',
             'Cart function for easy additional orders',
             'Real-time order status tracking',
         ]),
        ('screenshots/en/booking_calendar_desktop.png', 'Booking Calendar',
         '24/7 online reservations',
         [
             'View staff availability on calendar at a glance',
             'Customers pick their preferred date and time slot',
             'Prepaid confirmation prevents no-shows',
             'LINE and email booking support',
             'Auto-notification to staff on confirmation',
         ]),
    ]

    for img_name, title, subtitle, bullets in frontend_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="Customer View \u2014 ")

    # ── Mobile-Responsive ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Fully Mobile-Responsive", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "Both the customer site and admin dashboard work seamlessly on smartphones.",
             font_size=16, color=DARK)

    mobile_items = [
        ('screenshots/en/front_top_mobile_vp.png', 'Customer Top Page'),
        ('screenshots/en/admin_dashboard_mobile_vp.png', 'Admin Dashboard'),
    ]

    for i, (img_name, label) in enumerate(mobile_items):
        x = Inches(1.5) + i * Inches(5.5)
        y = Inches(1.8)

        add_shape(slide, x, y, Inches(4.2), Inches(5.0), WHITE, corner_radius=True)

        img_path = os.path.join(screenshots_dir, img_name)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                x + Inches(0.6), y + Inches(0.15),
                                Inches(3.0), Inches(4.4))
            except Exception:
                add_text(slide, x + Inches(0.5), y + Inches(2), Inches(3.2), Inches(1),
                         f'[{label}]', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(4.5), Inches(4.2), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ── Admin Screens ──
    admin_screens = [
        ('screenshots/en/dashboard_sales_vp.png', 'Sales Dashboard',
         'Business overview at a glance',
         [
             'Daily / weekly / monthly / yearly revenue trends',
             'Reservation KPIs (count, cancellation rate) in real time',
             'Three-channel breakdown: Total, EC Shop, In-store Menu',
             'KPI scorecards for key metrics overview',
             'Business Insights auto-suggest improvements',
         ]),
        ('screenshots/en/dashboard_menu_eng_vp.png', 'Sales Analytics (Menu Analysis)',
         'AI-powered profit optimization',
         [
             'Menu Engineering: Star / Plowhorse / Puzzle / Dog 4-quadrant classification',
             'ABC Analysis: Pareto principle to identify top revenue drivers',
             'Revenue Forecast: Prophet AI model predicts 14 days ahead',
             'Hourly sales heatmap for optimal staffing decisions',
             'AOV (Average Order Value) trends for pricing strategy review',
         ]),
        ('screenshots/en/dashboard_rfm_vp.png', 'Sales Analytics (Customer Analysis)',
         'Data-driven customer understanding',
         [
             'RFM Analysis: auto-segment customers by recency, frequency, monetary',
             'Cohort Analysis: visualize monthly new-customer retention rates',
             'Basket Analysis: discover cross-sell patterns and recommendations',
             'AI Summary: analysis results summarized in natural language',
             'Recommended Actions: AI auto-suggests specific improvement steps',
         ]),
        ('screenshots/en/shift_calendar_vp.png', 'Shift Calendar',
         'Streamlined shift management',
         [
             'Staff submit shift preferences from their smartphones',
             'Calendar view of all staff schedules at a glance',
             'Auto-scheduling with one-click assignment',
             'Auto-detect understaffed days with alerts',
             'Confirmed shifts auto-notified via LINE',
         ]),
        ('screenshots/en/pos_vp.png', 'POS Register',
         'Smooth checkout experience',
         [
             'Category tabs for quick product selection',
             'Cash, credit card, PayPay, and transit IC support',
             'Auto receipt generation and printing',
             'Kitchen display integration for order notifications',
             'Daily and monthly sales reports auto-generated',
         ]),
        ('screenshots/en/inventory_vp.png', 'Inventory Management',
         'Prevent stockouts',
         [
             'Real-time stock levels for all products',
             'Auto-deduct on order receipt',
             'Low stock alerts and notifications',
             'QR code receiving (no barcode scanner needed)',
             'Full stock-in/out history for audit compliance',
         ]),
        ('screenshots/en/customer_feedback_vp.png', 'Customer Feedback (NPS)',
         'Quantify the voice of your customers',
         [
             'NPS (Net Promoter Score) auto-aggregation and trending',
             'Promoter / Passive / Detractor color-coded display',
             'Comment-based feedback for actionable improvements',
             'Linked to order data for service quality analysis',
             'Monthly and staff-level satisfaction trends',
         ]),
        ('screenshots/en/attendance_board_vp.png', 'Attendance Board',
         'Real-time staff attendance status',
         [
             'View on-duty, on-break, and not-clocked-in status at a glance',
             'Auto-refresh every 30 seconds for live tracking',
             'Three clock-in methods: QR code, PIN, or smartphone',
             'Auto-classify overtime, night, and holiday hours',
             'Seamless data flow from attendance to payroll',
         ]),
        ('screenshots/en/iot_sensors_vp.png', 'IoT Sensor Monitor',
         'Real-time store environment monitoring',
         [
             'Temperature, humidity, pressure, gas levels in real time',
             'PIR motion sensor for automatic visitor counting',
             'Alert notifications on anomalies (gas leaks, high temps, etc.)',
             'Time-series charts to track sensor data trends',
             'ESP32 device connectivity monitoring',
         ]),
    ]

    for img_name, title, subtitle, bullets in admin_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="Admin Panel \u2014 ")

    # ==========================================
    # Benefits Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Benefits", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    merits = [
        ("Zero No-Show Losses",
         "Prepaid booking confirmation eliminates last-minute cancellations.\nToggle between prepaid / postpaid with one click.",
         "\U0001F6E1"),
        ("Effortless LINE Booking",
         "Customers book right from LINE chat.\nNo new app installation required.",
         "\U0001F4AC"),
        ("Boost Revenue",
         "QR ordering increases add-on order rate.\nMulti-language support captures international visitors.",
         "\U0001F4C8"),
        ("Save Hours of Work",
         "Automate booking, ordering, shifts, and payroll.\nFree your staff to focus on customer service.",
         "\u23F1"),
        ("Cut Costs",
         "Reservation books, order slips, time cards...\nEliminate paper costs and manual effort.",
         "\U0001F4B0"),
        ("Data-Driven Decisions",
         "Visualize sales, bookings, and staff performance.\nMake informed decisions backed by real data.",
         "\U0001F4CA"),
    ]

    for i, (title, desc, icon) in enumerate(merits):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), WHITE, corner_radius=True)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                 icon, font_size=30, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.35),
                 title, font_size=16, bold=True, color=BROWN)

        add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # Getting Started Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Getting Started", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    steps = [
        ("STEP 1", "Contact Us", "Reach out via phone\nor our contact form"),
        ("STEP 2", "Consultation", "We learn about your\ncurrent workflow and pain points"),
        ("STEP 3", "Demo & Quote", "See the product in action.\nWe propose the best plan for you"),
        ("STEP 4", "Account Setup", "Register store info and staff.\nWe assist with initial setup"),
        ("STEP 5", "Go Live", "After a quick walkthrough,\nyou're ready to go!"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.6)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BROWN
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            add_text(slide, x + Inches(1.9), y + Inches(0.2), Inches(0.6), Inches(0.6),
                     "\u2192", font_size=28, color=BROWN, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.2), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(1.5), LIGHT_GRAY, corner_radius=True)
    add_text(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1.2),
             "Start using it as soon as the same day.\n"
             "No special hardware or system installation required. Access instantly from your smartphone, tablet, or PC.\n"
             "We provide ongoing support, training, and feature updates after deployment.",
             font_size=13, color=DARK)

    # ==========================================
    # Pricing Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Pricing Plans", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "* Prices below are for reference. We customize plans to fit your needs.",
             font_size=12, color=GRAY)

    plans = [
        ("Light Plan", "From $XX/mo", ACCENT_GREEN,
         "For small stores", [
             "Reservation management (LINE & email)",
             "Staff management (up to 5)",
             "Website management",
             "Email & chat support",
         ]),
        ("Standard Plan", "From $XX/mo", ACCENT_BLUE,
         "Recommended!", [
             "All Light Plan features",
             "QR table ordering",
             "Inventory management",
             "Shift management",
             "Unlimited staff",
             "Phone support",
         ]),
        ("Premium Plan", "From $XX/mo", BROWN,
         "Multi-store / Full features", [
             "All Standard features",
             "Payroll & attendance",
             "Sales dashboard & analytics",
             "Multi-language support (7 languages)",
             "Multi-store management",
             "Priority support",
             "Custom development",
         ]),
    ]

    for i, (name, price, color, badge, features) in enumerate(plans):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.2)

        card = add_shape(slide, x, y, w, h, LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, w, Inches(0.08), color)

        if badge == "Recommended!":
            badge_shape = add_shape(slide, x + Inches(0.8), y - Inches(0.15), Inches(2.2), Inches(0.35), color, corner_radius=True)
            tf = badge_shape.text_frame
            p = tf.paragraphs[0]
            p.text = "\u2605 Recommended!"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER

        add_text(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.5),
                 price, font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

        add_shape(slide, x + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(0.02), GRAY)

        add_text(slide, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(0.3),
                 badge if badge != "Recommended!" else "For mid-size stores",
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9),
                        w - Inches(0.6), h - Inches(2.1),
                        features, font_size=12, icon="\u2713")

    add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "* Setup fees and customization charges quoted separately.  * Prices exclude tax.",
             font_size=10, color=GRAY)

    # ==========================================
    # FAQ Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "Frequently Asked Questions", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    faqs = [
        ("Q. Do I need any special equipment?",
         "A. No special hardware required. Use it on your existing smartphone, tablet, or PC.\n"
         "   For table ordering, just print QR codes from the admin panel."),
        ("Q. Can non-technical staff use it?",
         "A. Yes. If you can use a smartphone, you can use TimeBaiBai. The interface is intuitive and simple.\n"
         "   We also provide hands-on training during onboarding."),
        ("Q. Can I add or change features later?",
         "A. Absolutely. You can upgrade or change your plan at any time.\n"
         "   We recommend starting with core features and scaling up as needed."),
        ("Q. How secure is my data?",
         "A. Customer data is encrypted at rest. All communication uses SSL encryption.\n"
         "   Credit card information is processed by PCI DSS-compliant payment services \u2014 no card data is stored on our servers."),
        ("Q. Can I use it for multiple locations?",
         "A. Yes. The Premium plan supports multi-store management.\n"
         "   Manage staff and inventory per location while viewing consolidated reports across all stores."),
    ]

    for i, (q, a) in enumerate(faqs):
        y = Inches(1.2) + i * Inches(1.2)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(1.05), WHITE, corner_radius=True)

        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(11.73), Inches(0.3),
                 q, font_size=13, bold=True, color=BROWN)
        add_text(slide, Inches(0.8), y + Inches(0.4), Inches(11.73), Inches(0.6),
                 a, font_size=11, color=DARK)

    # ==========================================
    # Contact Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "Get in Touch",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             "Request a demo, get a quote, or ask us anything.",
             font_size=18, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.33), Inches(0.03), WHITE)

    contact_card = add_shape(slide, Inches(3), Inches(3.8), Inches(7.33), Inches(2.5), WHITE, corner_radius=True)

    contacts = [
        "TEL:     000-0000-0000",
        "Email:   info@example.com",
        "Web:    https://timebaibai.com",
        "",
        "Company Name Inc.",
        "Address line here",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, Inches(4), Inches(4.0) + i * Inches(0.35),
                 Inches(5), Inches(0.35),
                 line, font_size=15 if i < 3 else 13,
                 bold=(i < 3), color=DARK if i < 3 else GRAY,
                 align=PP_ALIGN.LEFT)

    add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "\u00A9 2026 TimeBaiBai. All Rights Reserved.",
             font_size=10, color=BEIGE, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# Output
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    pptx_path = os.path.join(output_dir, 'TimeBaiBai_Service_Overview_EN.pptx')
    prs.save(pptx_path)
    print(f"\u2705 PPT saved: {pptx_path}")

    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            print(f"\u2705 PDF saved: {pdf_path}")
        else:
            raise FileNotFoundError("LibreOffice not found")
    except Exception as e:
        print(f"\u26A0\uFE0F  PDF conversion skipped ({e})")
        print(f"   To create PDF from PPTX:")
        print(f"   - Open in PowerPoint / Keynote and export as PDF")
        print(f"   - Or: brew install libreoffice && soffice --headless --convert-to pdf '{pptx_path}'")
