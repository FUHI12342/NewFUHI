#!/usr/bin/env python3
"""
Sales Deck Generator — Korean Version (한국어)
Outputs PPT + PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ==============================
# Color Palette
# ==============================
BROWN = RGBColor(0x8C, 0x87, 0x6C)
BEIGE = RGBColor(0xF1, 0xF0, 0xEC)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0x8B, 0xCA)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)

FONT_NAME = 'Malgun Gothic'


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, spacing=Pt(6), icon="\u2713"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = spacing


def add_feature_card(slide, left, top, width, height, title, items, accent_color=BROWN):
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=True)
    line = add_shape(slide, left + Inches(0.15), top + Inches(0.08),
                     width - Inches(0.3), Inches(0.04), accent_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.2),
             width - Inches(0.4), Inches(0.4),
             title, font_size=15, bold=True, color=accent_color)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6),
                    width - Inches(0.4), height - Inches(0.8),
                    items, font_size=11, icon="\u25CF")


def add_picture_fit(slide, img_path, left, top, max_width, max_height):
    with Image.open(img_path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = max_width / max_height
    if ratio > box_ratio:
        w = max_width
        h = max_width / ratio
    else:
        h = max_height
        w = max_height * ratio
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    return slide.shapes.add_picture(img_path, int(x), int(y), width=int(w), height=int(h))


def add_number_highlight(slide, left, top, number, label, color=BROWN):
    add_text(slide, left, top, Inches(2), Inches(0.6),
             number, font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2), Inches(0.3),
             label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ==============================
# Slide Creation
# ==============================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ==========================================
    # Slide 1: Cover
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "올인원 매장 DX 플랫폼",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "예약 \u2022 주문 \u2022 결제 \u2022 직원 관리 \u2022 급여\n하나의 시스템으로 모두 해결.",
             font_size=22, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5), Inches(4.0), Inches(3.33), Inches(0.03), WHITE)

    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "TimeBaiBai  서비스 소개",
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "2026  /  회사명",
             font_size=14, color=BEIGE, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 2: Pain Points
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "이런 고민, 혹시 겪고 계신가요?", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    problems = [
        ("분산된 예약 채널",
         "전화, LINE, 이메일... 다양한 예약 방식으로\n이중 예약 걱정이 끊이지 않습니다.",
         "\U0001F4DE"),
        ("노쇼 및 갑작스러운 취소",
         "고객이 예약하고 나타나지 않습니다.\n매출 손실과 헛된 준비 시간.",
         "\U0001F630"),
        ("주문 접수의 혼란",
         "피크 시간대에 주문을 감당할 수 없습니다.\n주문 실수가 반복됩니다.",
         "\U0001F4DD"),
        ("재고 사각지대",
         "예고 없이 품절이 발생합니다.\n언제 발주해야 할지 알 수 없습니다.",
         "\U0001F4E6"),
    ]

    for i, (title, desc, icon) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.6) + row * Inches(2.6)

        card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), LIGHT_GRAY, corner_radius=True)

        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8),
                 icon, font_size=36, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK)

        add_text(slide, x + Inches(1.2), y + Inches(0.75), Inches(4), Inches(1.2),
                 desc, font_size=14, color=GRAY)

    # ==========================================
    # Slide 3: Solution — What is TimeBaiBai
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "TimeBaiBai가 모든 것을 해결합니다", font_size=32, bold=True, color=BROWN)

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "매장에 필요한 모든 것을 하나의 플랫폼에서. 스마트폰 하나로 시작하세요.",
             font_size=16, color=DARK)

    highlights = [
        ("6+", "핵심 모듈"),
        ("7", "지원 언어"),
        ("24h", "자동 예약"),
        ("$0", "종이 비용 절감"),
    ]
    for i, (num, label) in enumerate(highlights):
        x = Inches(1.2) + i * Inches(3)
        add_shape(slide, x - Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.4), WHITE, corner_radius=True)
        add_number_highlight(slide, x, Inches(2.5), num, label)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
             "예약, 테이블 주문, 재고, 근무 관리, 급여 계산, 웹사이트 관리...\n"
             "하나의 대시보드에서 모든 업무를 관리하세요 \u2014 여러 도구를 번갈아 사용할 필요가 없습니다.\n\n"
             "설정이 간단합니다. 계정 생성 당일 바로 사용할 수 있습니다.\n"
             "복잡한 설치나 하드웨어가 필요 없습니다. 스마트폰이나 PC에서 바로 접속하세요.",
             font_size=15, color=DARK)

    # ==========================================
    # Slide 4: Features (1/4) Reservations & Orders
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "주요 기능  \u2460  예약 / 테이블 주문",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "온라인 예약 시스템", [
            "LINE으로 예약 \u2014 채팅에서 바로 예약 완료",
            "이메일 예약도 지원",
            "선불 예약 확인 \u2014 결제 시 자동 확정 (노쇼 방지)",
            "관리자 화면에서 선불/후불 원클릭 전환",
            "직원 캘린더에 실시간 가용 현황 표시",
            "두 가지 예약 흐름: 날짜 기준 또는 직원 기준",
            "QR코드 체크인 (접수 인력 불필요)",
            "예약 확정 시 직원에게 LINE 자동 알림",
            "7개 언어 지원 (일본어, 영어, 중국어, 한국어, 스페인어, 포르투갈어)",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "QR 테이블 주문 시스템", [
            "각 테이블에 QR코드를 배치하면 준비 완료",
            "고객이 스마트폰으로 스캔 \u2192 메뉴 탐색 \u2192 주문",
            "사진 메뉴로 요리를 시각적으로 전달",
            "카테고리 탭으로 빠른 메뉴 검색",
            "장바구니 기능으로 추가 주문도 간편",
            "실시간 주문 상태 추적",
            "현금, 신용카드, 전자화폐 결제 지원",
            "관리자 화면에서 QR코드 일괄 생성 및 인쇄용 다운로드",
        ], ACCENT_GREEN)

    # ==========================================
    # Slide 5: Features (2/4) Inventory & Staff
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "주요 기능  \u2461  재고 / 직원 근무 관리",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "재고 관리", [
            "모든 상품의 실시간 재고 현황",
            "수신 주문과 연동한 재고 자동 차감",
            "재고 부족 시 자동 알림",
            "QR코드 입고 (바코드 스캐너 불필요)",
            "모든 입출고 이력 관리 (감사 대비)",
            "상품 카테고리별 정리",
            "EC 쇼핑몰과 매장 재고 통합 관리",
            "일괄 조정을 위한 재고 조사 기능",
        ], ACCENT_ORANGE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "직원 근무 관리", [
            "직원이 스마트폰에서 근무 희망 제출",
            "3단계 요청: 가능 / 희망 / 불가",
            "관리자가 대시보드에서 근무 확정",
            "효율적 인력 배치를 위한 자동 스케줄링",
            "확정된 근무표 LINE으로 직원에게 자동 알림",
            "사진과 소개가 포함된 직원 프로필",
            "직원 유형: 캐스트(서비스) / 매장 스태프",
            "역할별 권한: 오너 / 매니저 / 스태프 / 개발자",
        ], ACCENT_BLUE)

    # ==========================================
    # Slide 6: Features (3/4) Payroll & Website
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "주요 기능  \u2462  급여 / 웹사이트 관리",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "급여 및 출퇴근 관리", [
            "근무 기록에서 출퇴근 데이터 자동 생성",
            "일반, 야간, 휴일 근무 자동 분류",
            "시급제 및 월급제 직원 급여 계산",
            "사회보험료 자동 산출",
            "소득세 및 주민세 자동 공제",
            "원클릭 급여명세서 생성",
            "은행 이체 CSV 내보내기",
            "각종 수당 간편 설정 (교통, 주거, 가족)",
        ], ACCENT_GREEN)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "웹사이트 및 마케팅 관리", [
            "관리자 화면에서 매장 웹사이트 편집",
            "히어로 배너(슬라이더)로 매장 어필",
            "뉴스 및 미디어 소개 간편 업데이트",
            "배너 광고 배치 및 관리",
            "소셜 미디어 연동 (X / Instagram 임베드)",
            "유연한 레이아웃을 위한 커스텀 HTML 블록",
            "개인정보보호 방침 등 법적 페이지 관리자 편집 가능",
            "온라인 판매를 위한 내장 EC 쇼핑몰",
        ], ACCENT_ORANGE)

    # ==========================================
    # Slide 7: Features (4/4) Analytics & Security
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "주요 기능  \u2463  데이터 분석 / 보안",
             font_size=28, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_feature_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
        "데이터 분석 및 AI", [
            "일별/주별/월별/연별 매출 추이 대시보드",
            "메뉴 엔지니어링 (Star/Dog 분류)으로 수익 개선",
            "ABC 분석, RFM 분석, 코호트 분석, 장바구니 분석",
            "NPS (고객 만족도) 자동 집계",
            "AI가 시간대별 최적 인력 배치 예측",
            "수요 및 매출 예측으로 발주·인력 최적화",
            "비즈니스 인사이트가 이상 징후와 기회를 자동 감지",
            "KPI 스코어카드로 사업 건전성 한눈에 파악",
        ], ACCENT_BLUE)

    add_feature_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
        "보안 및 감사", [
            "전체 접속 로그 (로그인, 작업, 이상 탐지)",
            "IP 기반 속도 제한으로 무단 접근 자동 차단",
            "고객 개인정보 AES 암호화 저장",
            "자동 보안 감사 (매일 12개 항목 검사)",
            "역할 기반 접근 제어 (오너/매니저/스태프/개발자)",
            "완전한 감사 추적: 작업 이력 및 근무 변경 로그",
            "SSL 암호화 및 CSRF 보호 기본 제공",
            "PCI DSS 준수 외부 결제 서비스 연동",
        ], ACCENT_ORANGE)

    # ==========================================
    # Screenshot Slides
    # ==========================================
    screenshots_dir = os.path.dirname(os.path.abspath(__file__))

    def _add_desktop_slide(img_name, title, subtitle, bullets, header_prefix=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        full_title = f"{header_prefix}{title}" if header_prefix else title
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 full_title, font_size=28, bold=True, color=BROWN)
        add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)
        img_path = os.path.join(screenshots_dir, img_name)
        add_shape(slide, Inches(0.4), Inches(1.2), Inches(8.2), Inches(5.8),
                  LIGHT_GRAY, corner_radius=True)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.6))
            except Exception:
                add_text(slide, Inches(2), Inches(3.5), Inches(4), Inches(1),
                         f'[{title}]', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(0.5),
                 subtitle, font_size=18, bold=True, color=DARK)
        add_bullet_list(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(5),
                        bullets, font_size=12, icon="\u2713")

    # ── Customer-Facing Screens ──
    frontend_screens = [
        ('screenshots/ko/front_top_desktop.png', '메인 페이지',
         '고객이 처음 보는 화면',
         [
             '히어로 배너로 매장의 매력을 어필',
             '다양한 옵션의 명확한 예약 흐름',
             '인기 캐스트를 홍보하는 직원 랭킹',
             '뉴스 및 미디어 소개 자동 표시',
             '외국인 고객을 위한 7개 언어 지원',
         ]),
        ('screenshots/ko/front_staff_list_desktop.png', '직원 프로필 및 매장 정보',
         '직원을 매력적으로 소개',
         [
             '직원 사진 및 프로필 소개',
             '전문 분야 및 자격 표시',
             '예약 캘린더로 바로 연결',
             '매장 정보 (영업시간, 지도, 오시는 길)',
             '관리자 화면에서 간편하게 콘텐츠 업데이트',
         ]),
        ('screenshots/ko/front_news_desktop.png', '뉴스 및 공지',
         '고객에게 최신 정보 전달',
         [
             '신메뉴, 캠페인, 프로모션 게시',
             '관리자 화면에서 원클릭 게시',
             '최신 콘텐츠 우선 표시를 위한 날짜순 자동 정렬',
             '폭넓은 도달을 위한 소셜 미디어 연동',
             '시의적절한 업데이트로 재방문 유도',
         ]),
        ('screenshots/ko/front_shop_desktop.png', '온라인 쇼핑몰',
         'EC로 판매 채널 확장',
         [
             '온라인으로 상품 판매',
             '상품 검색 및 카테고리 필터링',
             '복수 상품 구매를 위한 장바구니',
             '신용카드 및 전자화폐 결제',
             '통합 재고로 재고 불일치 방지',
         ]),
        ('screenshots/ko/table_order_vp.png', 'QR 테이블 주문',
         '스마트폰으로 간편 주문',
         [
             '테이블의 QR코드를 스캔하기만 하면 됩니다',
             '사진 메뉴로 요리를 시각적으로 전달',
             '카테고리 탭으로 빠른 메뉴 검색',
             '장바구니 기능으로 추가 주문도 간편',
             '실시간 주문 상태 추적',
         ]),
        ('screenshots/ko/booking_calendar_desktop.png', '예약 캘린더',
         '24시간 온라인 예약',
         [
             '캘린더에서 직원 가용 현황을 한눈에 확인',
             '고객이 원하는 날짜와 시간대를 선택',
             '선불 확인으로 노쇼 방지',
             'LINE 및 이메일 예약 지원',
             '확정 시 직원에게 자동 알림',
         ]),
    ]

    for img_name, title, subtitle, bullets in frontend_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="고객 화면 \u2014 ")

    # ── Mobile-Responsive ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "완벽한 모바일 대응", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "고객 사이트와 관리자 대시보드 모두 스마트폰에서 완벽하게 작동합니다.",
             font_size=16, color=DARK)

    mobile_items = [
        ('screenshots/ko/front_top_mobile_vp.png', '고객 메인 페이지'),
        ('screenshots/ko/admin_dashboard_mobile_vp.png', '관리자 대시보드'),
    ]

    for i, (img_name, label) in enumerate(mobile_items):
        x = Inches(1.5) + i * Inches(5.5)
        y = Inches(1.8)

        add_shape(slide, x, y, Inches(4.2), Inches(5.0), WHITE, corner_radius=True)

        img_path = os.path.join(screenshots_dir, img_name)
        if os.path.exists(img_path):
            try:
                add_picture_fit(slide, img_path,
                                x + Inches(0.6), y + Inches(0.15),
                                Inches(3.0), Inches(4.4))
            except Exception:
                add_text(slide, x + Inches(0.5), y + Inches(2), Inches(3.2), Inches(1),
                         f'[{label}]', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(4.5), Inches(4.2), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ── Admin Screens ──
    admin_screens = [
        ('screenshots/ko/dashboard_sales_vp.png', '매출 대시보드',
         '사업 현황을 한눈에',
         [
             '일별 / 주별 / 월별 / 연별 매출 추이',
             '예약 KPI (건수, 취소율) 실시간 확인',
             '3채널 분석: 전체, EC 쇼핑몰, 매장 메뉴',
             'KPI 스코어카드로 주요 지표 한눈에 파악',
             '비즈니스 인사이트가 개선점을 자동 제안',
         ]),
        ('screenshots/ko/dashboard_menu_eng_vp.png', '매출 분석 (메뉴 분석)',
         'AI 기반 수익 최적화',
         [
             '메뉴 엔지니어링: Star / Plowhorse / Puzzle / Dog 4분류',
             'ABC 분석: 파레토 법칙으로 주요 매출 동력 파악',
             '매출 예측: Prophet AI 모델로 14일 앞 예측',
             '시간대별 매출 히트맵으로 최적 인력 배치',
             'AOV (객단가) 추이로 가격 전략 검토',
         ]),
        ('screenshots/ko/dashboard_rfm_vp.png', '매출 분석 (고객 분석)',
         '데이터 기반 고객 이해',
         [
             'RFM 분석: 최근성, 빈도, 금액으로 고객 자동 세분화',
             '코호트 분석: 월별 신규 고객 유지율 시각화',
             '장바구니 분석: 교차 판매 패턴과 추천 발견',
             'AI 요약: 분석 결과를 자연어로 요약',
             '추천 조치: AI가 구체적인 개선 단계를 자동 제안',
         ]),
        ('screenshots/ko/shift_calendar_vp.png', '근무 캘린더',
         '효율적인 근무 관리',
         [
             '직원이 스마트폰에서 근무 희망 제출',
             '캘린더 뷰에서 전체 직원 일정을 한눈에 확인',
             '원클릭 배정의 자동 스케줄링',
             '인력 부족일 자동 감지 및 알림',
             '확정된 근무표 LINE으로 자동 알림',
         ]),
        ('screenshots/ko/pos_vp.png', 'POS 레지스터',
         '원활한 결제 경험',
         [
             '카테고리 탭으로 빠른 상품 선택',
             '현금, 신용카드, PayPay, 교통IC 지원',
             '영수증 자동 생성 및 인쇄',
             '주방 디스플레이 연동으로 주문 알림',
             '일별 및 월별 매출 보고서 자동 생성',
         ]),
        ('screenshots/ko/inventory_vp.png', '재고 관리',
         '품절 방지',
         [
             '모든 상품의 실시간 재고 현황',
             '주문 접수 시 자동 차감',
             '재고 부족 알림 및 통지',
             'QR코드 입고 (바코드 스캐너 불필요)',
             '감사 대비를 위한 전체 입출고 이력',
         ]),
        ('screenshots/ko/customer_feedback_vp.png', '고객 피드백 (NPS)',
         '고객의 목소리를 수치화',
         [
             'NPS (순추천지수) 자동 집계 및 추이',
             '추천자 / 중립자 / 비추천자 색상 구분 표시',
             '실행 가능한 개선을 위한 코멘트 기반 피드백',
             '주문 데이터와 연동한 서비스 품질 분석',
             '월별 및 직원별 만족도 추이',
         ]),
        ('screenshots/ko/attendance_board_vp.png', '출퇴근 현황판',
         '실시간 직원 출퇴근 상태',
         [
             '근무중, 휴식중, 미출근 상태를 한눈에 확인',
             '실시간 추적을 위한 30초 자동 새로고침',
             '3가지 출근 방법: QR코드, PIN, 스마트폰',
             '초과근무, 야간, 휴일 근무 자동 분류',
             '출퇴근에서 급여까지 원활한 데이터 연동',
         ]),
        ('screenshots/ko/iot_sensors_vp.png', 'IoT 센서 모니터',
         '매장 환경 실시간 모니터링',
         [
             '온도, 습도, 기압, 가스 수치 실시간 확인',
             'PIR 모션 센서로 자동 방문객 집계',
             '이상 발생 시 알림 (가스 누출, 고온 등)',
             '시계열 차트로 센서 데이터 추이 추적',
             'ESP32 디바이스 연결 상태 모니터링',
         ]),
    ]

    for img_name, title, subtitle, bullets in admin_screens:
        _add_desktop_slide(img_name, title, subtitle, bullets,
                           header_prefix="관리자 화면 \u2014 ")

    # ==========================================
    # Benefits Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "도입 효과", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    merits = [
        ("노쇼 손실 제로",
         "선불 예약 확인으로 갑작스러운 취소를 방지합니다.\n선불/후불 원클릭 전환.",
         "\U0001F6E1"),
        ("LINE으로 간편 예약",
         "고객이 LINE 채팅에서 바로 예약합니다.\n새로운 앱 설치가 필요 없습니다.",
         "\U0001F4AC"),
        ("매출 향상",
         "QR 주문으로 추가 주문율 증가.\n다국어 지원으로 외국인 고객 확보.",
         "\U0001F4C8"),
        ("업무 시간 절감",
         "예약, 주문, 근무, 급여를 자동화합니다.\n직원이 고객 서비스에 집중할 수 있습니다.",
         "\u23F1"),
        ("비용 절감",
         "예약장, 주문 전표, 출근 카드...\n종이 비용과 수작업을 없앱니다.",
         "\U0001F4B0"),
        ("데이터 기반 의사결정",
         "매출, 예약, 직원 성과를 시각화합니다.\n실제 데이터에 기반한 현명한 의사결정.",
         "\U0001F4CA"),
    ]

    for i, (title, desc, icon) in enumerate(merits):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), WHITE, corner_radius=True)

        add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                 icon, font_size=30, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.35),
                 title, font_size=16, bold=True, color=BROWN)

        add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                 desc, font_size=12, color=DARK)

    # ==========================================
    # Getting Started Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "도입 절차", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), BROWN)

    steps = [
        ("STEP 1", "문의", "전화 또는\n문의 양식으로 연락"),
        ("STEP 2", "상담", "현재 업무 흐름과\n고민 사항을 파악합니다"),
        ("STEP 3", "데모 & 견적", "제품을 직접 체험하세요.\n최적의 플랜을 제안합니다"),
        ("STEP 4", "계정 설정", "매장 정보와 직원을 등록.\n초기 설정을 도와드립니다"),
        ("STEP 5", "운영 시작", "간단한 안내 후\n바로 시작하세요!"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.6)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BROWN
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            add_text(slide, x + Inches(1.9), y + Inches(0.2), Inches(0.6), Inches(0.6),
                     "\u2192", font_size=28, color=BROWN, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.2), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(1.5), LIGHT_GRAY, corner_radius=True)
    add_text(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1.2),
             "당일 바로 사용 가능합니다.\n"
             "별도의 하드웨어나 시스템 설치가 필요 없습니다. 스마트폰, 태블릿, PC에서 바로 접속하세요.\n"
             "도입 후에도 지속적인 지원, 교육 및 기능 업데이트를 제공합니다.",
             font_size=13, color=DARK)

    # ==========================================
    # Pricing Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "요금 안내", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "* 아래 가격은 참고용입니다. 고객님의 필요에 맞춰 맞춤 플랜을 제공합니다.",
             font_size=12, color=GRAY)

    plans = [
        ("라이트 플랜", "$XX~/월", ACCENT_GREEN,
         "소규모 매장용", [
             "예약 관리 (LINE & 이메일)",
             "직원 관리 (최대 5명)",
             "웹사이트 관리",
             "이메일 & 채팅 지원",
         ]),
        ("스탠다드 플랜", "$XX~/월", ACCENT_BLUE,
         "추천!", [
             "라이트 플랜의 모든 기능",
             "QR 테이블 주문",
             "재고 관리",
             "근무 관리",
             "직원 수 무제한",
             "전화 지원",
         ]),
        ("프리미엄 플랜", "$XX~/월", BROWN,
         "다점포 / 전체 기능", [
             "스탠다드 플랜의 모든 기능",
             "급여 및 출퇴근 관리",
             "매출 대시보드 및 분석",
             "다국어 지원 (7개 언어)",
             "다점포 관리",
             "우선 지원",
             "맞춤 개발",
         ]),
    ]

    for i, (name, price, color, badge, features) in enumerate(plans):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.2)

        card = add_shape(slide, x, y, w, h, LIGHT_GRAY, corner_radius=True)
        add_shape(slide, x, y, w, Inches(0.08), color)

        if badge == "추천!":
            badge_shape = add_shape(slide, x + Inches(0.8), y - Inches(0.15), Inches(2.2), Inches(0.35), color, corner_radius=True)
            tf = badge_shape.text_frame
            p = tf.paragraphs[0]
            p.text = "\u2605 추천!"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER

        add_text(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.5),
                 price, font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

        add_shape(slide, x + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(0.02), GRAY)

        add_text(slide, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(0.3),
                 badge if badge != "추천!" else "중간 규모 매장용",
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9),
                        w - Inches(0.6), h - Inches(2.1),
                        features, font_size=12, icon="\u2713")

    add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "* 초기 설정비 및 커스터마이즈 비용은 별도 견적.  * 가격은 세금 별도.",
             font_size=10, color=GRAY)

    # ==========================================
    # FAQ Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BEIGE)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
             "자주 묻는 질문", font_size=32, bold=True, color=BROWN)
    add_shape(slide, Inches(0.8), Inches(0.9), Inches(2), Inches(0.04), BROWN)

    faqs = [
        ("Q. 특별한 장비가 필요한가요?",
         "A. 별도의 하드웨어가 필요 없습니다. 기존 스마트폰, 태블릿, PC에서 사용하세요.\n"
         "   테이블 주문의 경우 관리자 화면에서 QR코드를 인쇄하면 됩니다."),
        ("Q. IT에 익숙하지 않은 직원도 사용할 수 있나요?",
         "A. 네. 스마트폰을 사용할 수 있다면 TimeBaiBai도 사용할 수 있습니다.\n"
         "   도입 시 실무 교육도 제공합니다."),
        ("Q. 나중에 기능을 추가하거나 변경할 수 있나요?",
         "A. 물론입니다. 언제든지 플랜을 업그레이드하거나 변경할 수 있습니다.\n"
         "   핵심 기능부터 시작하여 필요에 따라 확장하는 것을 추천합니다."),
        ("Q. 데이터는 안전한가요?",
         "A. 고객 데이터는 암호화하여 저장됩니다. 모든 통신은 SSL 암호화를 사용합니다.\n"
         "   신용카드 정보는 PCI DSS 준수 결제 서비스에서 처리하며, 당사 서버에 카드 정보를 저장하지 않습니다."),
        ("Q. 여러 매장에서 사용할 수 있나요?",
         "A. 네. 프리미엄 플랜은 다점포 관리를 지원합니다.\n"
         "   매장별로 직원과 재고를 관리하면서 전체 매장의 통합 보고서를 확인할 수 있습니다."),
    ]

    for i, (q, a) in enumerate(faqs):
        y = Inches(1.2) + i * Inches(1.2)
        card = add_shape(slide, Inches(0.5), y, Inches(12.33), Inches(1.05), WHITE, corner_radius=True)

        add_text(slide, Inches(0.8), y + Inches(0.08), Inches(11.73), Inches(0.3),
                 q, font_size=13, bold=True, color=BROWN)
        add_text(slide, Inches(0.8), y + Inches(0.4), Inches(11.73), Inches(0.6),
                 a, font_size=11, color=DARK)

    # ==========================================
    # Contact Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BROWN)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
             "문의하기",
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             "데모, 견적 또는 궁금한 점이 있으시면 문의해 주세요.",
             font_size=18, color=BEIGE, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.33), Inches(0.03), WHITE)

    contact_card = add_shape(slide, Inches(3), Inches(3.8), Inches(7.33), Inches(2.5), WHITE, corner_radius=True)

    contacts = [
        "TEL:     000-0000-0000",
        "Email:   info@example.com",
        "Web:    https://timebaibai.com",
        "",
        "회사명",
        "주소",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, Inches(4), Inches(4.0) + i * Inches(0.35),
                 Inches(5), Inches(0.35),
                 line, font_size=15 if i < 3 else 13,
                 bold=(i < 3), color=DARK if i < 3 else GRAY,
                 align=PP_ALIGN.LEFT)

    add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "\u00A9 2026 TimeBaiBai. All Rights Reserved.",
             font_size=10, color=BEIGE, align=PP_ALIGN.CENTER)

    return prs


# ==============================
# Output
# ==============================
if __name__ == '__main__':
    output_dir = os.path.dirname(os.path.abspath(__file__))

    prs = create_presentation()

    pptx_path = os.path.join(output_dir, 'TimeBaiBai_Service_Overview_KO.pptx')
    prs.save(pptx_path)
    print(f"\u2705 PPT saved: {pptx_path}")

    pdf_path = pptx_path.replace('.pptx', '.pdf')
    try:
        import subprocess
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            print(f"\u2705 PDF saved: {pdf_path}")
        else:
            raise FileNotFoundError("LibreOffice not found")
    except Exception as e:
        print(f"\u26A0\uFE0F  PDF conversion skipped ({e})")
        print(f"   To create PDF from PPTX:")
        print(f"   - Open in PowerPoint / Keynote and export as PDF")
        print(f"   - Or: brew install libreoffice && soffice --headless --convert-to pdf '{pptx_path}'")
